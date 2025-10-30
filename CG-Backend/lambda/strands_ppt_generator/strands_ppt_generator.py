"""
PowerPoint Presentation Generator - Strands Agent
==================================================
Generates engaging PowerPoint presentations from theory book content using Strands Agents.
Uses existing images from the book to avoid regenerating visuals.

Expected event parameters:
    - course_bucket: S3 bucket name
    - project_folder: Project folder path
    - book_version_key: S3 key to specific book version JSON (optional, defaults to latest)
    - model_provider: 'bedrock' or 'openai'
    - slides_per_lesson: Number of slides per lesson (default: 5-7)
    - presentation_style: 'professional' | 'educational' | 'modern' (default: 'professional')
"""

import os
import json
import boto3
import re
import requests
from datetime import datetime
from typing import Dict, List, Any, Optional
from botocore.config import Config

# Configure boto3 with extended timeouts
boto_config = Config(
    read_timeout=600,
    connect_timeout=60,
    retries={'max_attempts': 3, 'mode': 'adaptive'}
)

# AWS Clients
s3_client = boto3.client('s3')
bedrock_client = boto3.client('bedrock-runtime', region_name='us-east-1', config=boto_config)
secrets_client = boto3.client('secretsmanager', region_name='us-east-1')

# Model Configuration
DEFAULT_BEDROCK_MODEL = "us.anthropic.claude-sonnet-4-5-20250929-v1:0"
DEFAULT_OPENAI_MODEL = "gpt-5"


def get_secret(secret_name: str) -> dict:
    """Retrieve secret from AWS Secrets Manager."""
    try:
        response = secrets_client.get_secret_value(SecretId=secret_name)
        return json.loads(response['SecretString'])
    except Exception as e:
        print(f"⚠️ Error retrieving secret {secret_name}: {e}")
        raise


def configure_bedrock_model(model_id: str = DEFAULT_BEDROCK_MODEL):
    """Configure Bedrock model for Strands Agent."""
    try:
        from strands.models import BedrockModel
        
        model = BedrockModel(
            model_id=model_id,
            client=bedrock_client
        )
        print(f"✅ Configured Bedrock model: {model_id}")
        return model
    except Exception as e:
        print(f"❌ Failed to configure Bedrock model: {e}")
        raise


def configure_openai_model(model_id: str = DEFAULT_OPENAI_MODEL) -> Any:
    """Configure OpenAI model for Strands Agent."""
    try:
        from strands.models.openai import OpenAIModel
        
        # Get API key from Secrets Manager
        secret_data = get_secret('aurora/openai-api-key')
        api_key = secret_data.get('api_key')
        
        if not api_key:
            raise ValueError("OpenAI API key not found in secrets")
        
        model = OpenAIModel(
            client_args={"api_key": api_key},
            model_id=model_id,
            streaming=False
        )
        print(f"✅ Configured OpenAI model: {model_id}")
        return model
    except Exception as e:
        print(f"❌ Failed to configure OpenAI model: {e}")
        raise


def load_book_from_s3(bucket: str, book_key: str) -> Dict:
    """Load book JSON from S3."""
    try:
        print(f"📖 Loading book from s3://{bucket}/{book_key}")
        print(f"   Bucket: '{bucket}'")
        print(f"   Key: '{book_key}'")
        
        # Check if the object exists first
        try:
            s3_client.head_object(Bucket=bucket, Key=book_key)
            print(f"✅ Object exists, proceeding with download")
        except Exception as head_err:
            print(f"❌ Object does not exist or no access: {head_err}")
            # Try listing objects in the folder to help debug
            try:
                folder = '/'.join(book_key.split('/')[:-1])
                print(f"🔍 Listing objects in folder: {folder}/")
                list_response = s3_client.list_objects_v2(Bucket=bucket, Prefix=folder + '/', MaxKeys=10)
                if 'Contents' in list_response:
                    print(f"   Found {len(list_response['Contents'])} objects:")
                    for obj in list_response['Contents']:
                        print(f"   - {obj['Key']}")
                else:
                    print(f"   No objects found in folder")
            except Exception as list_err:
                print(f"❌ Could not list folder contents: {list_err}")
            raise head_err
        
        response = s3_client.get_object(Bucket=bucket, Key=book_key)
        book_data = json.loads(response['Body'].read().decode('utf-8'))
        print(f"✅ Loaded book with {len(book_data.get('lessons', []))} lessons")
        return book_data
    except Exception as e:
        print(f"❌ Failed to load book: {e}")
        print(f"   Error type: {type(e).__name__}")
        raise


def extract_images_from_content(content: str) -> List[Dict[str, str]]:
    """Extract image references from lesson content."""
    images = []

    # Match markdown image syntax: ![alt](url)
    img_pattern = r'!\[([^\]]*)\]\(([^\)]+)\)'
    matches = re.findall(img_pattern, content)

    print(f"🔍 Found {len(matches)} image matches in content")

    for alt_text, img_url in matches:
        print(f"📷 Processing image: {alt_text} -> {img_url}")
        images.append({
            'alt_text': alt_text,
            'url': img_url,
            'type': 'diagram' if 'diagram' in alt_text.lower() else 'image'
        })

    # Also look for VISUAL tags which are common in our books
    visual_pattern = r'\[VISUAL([^\]]*)\]'
    visual_matches = re.findall(visual_pattern, content)

    print(f"🔍 Found {len(visual_matches)} VISUAL tags in content")

    for visual_desc in visual_matches:
        visual_desc = visual_desc.strip()
        if visual_desc:
            # Extract the actual description from formats like ": description"
            if visual_desc.startswith(':'):
                visual_desc = visual_desc[1:].strip()
            
            print(f"📷 Processing VISUAL: {visual_desc}")
            images.append({
                'alt_text': visual_desc if visual_desc else 'Visual element',
                'url': None,  # Will be filled by searching S3 or using placeholder
                'type': 'diagram' if 'diagram' in visual_desc.lower() else 'image'
            })

    return images


def generate_presentation_structure(
    book_data: Dict,
    model,
    slides_per_lesson: int = 6,
    presentation_style: str = 'professional',
    course_bucket: str = None,
    project_folder: str = None
) -> Dict:
    """
    Use Strands Agent to generate PowerPoint presentation structure from book content.
    Creates engaging, professional slides optimized for course delivery.
    """
    from strands import Agent

    # Extract course metadata
    metadata = book_data.get('metadata', {})
    course_title = metadata.get('title', 'Course Presentation')
    lessons = book_data.get('lessons', [])

    print(f"\n🎨 Generating AI-powered presentation structure for: {course_title}")
    print(f"📊 Lessons to convert: {len(lessons)}")
    print(f"🎯 Slides per lesson: {slides_per_lesson}")
    print(f"✨ Style: {presentation_style}")

    # Check if book has modules instead of lessons
    if len(lessons) == 0 and 'modules' in book_data:
        print("⚠️ Book uses 'modules' structure, converting to lessons...")
        modules = book_data.get('modules', [])
        for module in modules:
            if 'lessons' in module and isinstance(module['lessons'], list):
                lessons.extend(module['lessons'])
        print(f"✅ Extracted {len(lessons)} lessons from {len(modules)} modules")

    if len(lessons) == 0:
        raise ValueError("No lessons found in book data. Cannot generate presentation.")

    # Build comprehensive course context
    course_context = f"""
COURSE: {course_title}
TOTAL LESSONS: {len(lessons)}
PRESENTATION STYLE: {presentation_style}
SLIDES PER LESSON: {slides_per_lesson}

LESSON TITLES:
"""
    for i, lesson in enumerate(lessons, 1):
        course_context += f"{i}. {lesson.get('title', f'Lesson {i}')}\n"

    # Create AI-powered PPT Designer Agent
    ppt_designer = Agent(
        model=model,
        system_prompt=f"""You are a senior instructional designer and PowerPoint expert specializing in creating engaging, professional presentations for corporate training and educational courses.

Style: {presentation_style}
- Professional: Clean, corporate design with blue color scheme, structured layouts, business-focused
- Educational: Student-friendly, colorful, emphasizes examples and practice exercises, interactive elements
- Modern: Minimalist, high-impact visuals, brief text, dynamic layouts, contemporary design

COURSE CONTEXT:
{course_context}

EXPERTISE REQUIREMENTS:
1. Each slide MUST have:
    - A compelling, descriptive title (not generic like "Introduction" or "Overview")
    - 3-7 bullet points maximum (concise, action-oriented, impactful) - NEVER create empty bullet lists
    - Clear indication if existing image should be used
    - CRITICAL: Every content slide MUST have bullets array with at least 3 meaningful bullet points

2. Advanced Slide Types to Use:
    - Title Slide: Course introduction with compelling hook and value proposition
    - Learning Objectives Slide: Specific, measurable outcomes with business impact
    - Content Slide: Main concepts with practical applications and real-world examples
    - Image Slide: Full-screen diagram/visual with detailed caption and learning context
    - Process Slide: Step-by-step workflows with clear progression
    - Comparison Slide: Side-by-side concepts, before/after scenarios, pros/cons
    - Case Study Slide: Real-world applications with outcomes
    - Interactive Slide: Discussion questions, exercises, reflection points
    - Summary Slide: Key takeaways, action items, and next steps

3. Strategic Image Usage:
    - ALWAYS reference existing images from the lesson when available
    - Specify which image to use: "USE_IMAGE: [exact alt_text or description]"
    - Don't request new images - use what exists in the book
    - Place images strategically to enhance learning and retention
    - Use images to break up text-heavy sections

4. Professional Course Delivery Elements:
    - Include specific learning objectives for each lesson (measurable outcomes)
    - Add estimated timing for each section
    - Include discussion questions and interaction points
    - Add practical exercises and hands-on activities
    - Create reflection points and knowledge checks
    - Include real-world applications and case studies
    - Add summary slides with key takeaways and action items
    - Include next steps and homework assignments

5. Content Adaptation for Maximum Engagement:
    - Convert complex paragraphs into clear, memorable bullet points
    - Extract key terminology for emphasis and discussion
    - Identify practical examples that resonate with the audience
    - Create "what this means for you" connections
    - Build concepts progressively with clear learning progression
    - Use storytelling techniques to maintain interest
    - Include surprises and "aha" moments

6. Contextual Understanding and Flow:
    - Maintain cohesion between slides within a lesson
    - Create smooth transitions between lessons
    - Build upon concepts progressively
    - Reference previous lessons when relevant
    - Create anticipation for upcoming content
    - End each lesson with clear next steps

OUTPUT FORMAT: Return a JSON structure with this exact format:
{{
    "presentation_title": "{course_title}",
    "total_slides": 45,
    "slides": [
        {{
            "slide_number": 1,
            "slide_type": "title",
            "title": "{course_title}",
            "subtitle": "Professional Training Program",
            "notes": "Welcome participants and set expectations"
        }},
        {{
            "slide_number": 2,
            "slide_type": "content",
            "title": "What You'll Achieve Today",
            "bullets": [
                "Master practical skills for immediate application",
                "Understand real-world scenarios and challenges",
                "Develop strategies for long-term success"
            ],
            "image_reference": "USE_IMAGE: Course overview diagram",
            "notes": "Connect objectives to participant goals"
        }},
        {{
            "slide_number": 3,
            "slide_type": "image",
            "title": "System Architecture Deep Dive",
            "image_reference": "USE_IMAGE: Complete system architecture",
            "caption": "Understanding the components and their critical interactions",
            "notes": "Walk through each component and explain relationships"
        }}
    ]
}}
""",
        tools=[]
    )

    # Process lessons with AI agent for professional slide creation
    all_slides = []
    slide_counter = 1

    # Add title slide
    all_slides.append({
        "slide_number": slide_counter,
        "slide_type": "title",
        "title": course_title,
        "subtitle": f"Professional Training • {len(lessons)} Lessons",
        "notes": "Welcome and course overview"
    })
    slide_counter += 1

    # Process each lesson with AI for professional slide design
    for lesson_idx, lesson in enumerate(lessons, 1):
        lesson_title = lesson.get('title', f'Lesson {lesson_idx}')
        lesson_content = lesson.get('content', '')

        print(f"\n📝 Processing Lesson {lesson_idx}/{len(lessons)}: {lesson_title}")

        # Extract images available in this lesson
        available_images = extract_images_from_content(lesson_content)
        print(f"🖼️  Found {len(available_images)} images in this lesson")
        
        # Track used images to avoid repetition
        used_image_indices = set()

        # Use AI agent to create professional slides for this lesson
        lesson_prompt = f"""
Create professional PowerPoint slides for this lesson:

LESSON {lesson_idx}: {lesson_title}

CONTENT:
{lesson_content}

Available images: {len(available_images)} images found
Image descriptions: {[img['alt_text'] for img in available_images]}

Create {slides_per_lesson} engaging slides that are optimized for course delivery.
Include learning objectives, key concepts, practical examples, and discussion points.

Focus on:
- Creating contextually connected slides with smooth flow
- Using specific, compelling titles for each slide
- Including practical, real-world examples
- Adding discussion questions and interaction points
- Leveraging available images strategically
- Building concepts progressively within the lesson
"""

        # Get AI-generated slide structure for this lesson
        print(f"🤖 Calling AI agent to generate slides...")
        
        # Use Strands Agent to generate slides - proper invocation
        ai_response = ppt_designer(lesson_prompt)
        
        print(f"📝 AI Response type: {type(ai_response)}")
        
        # Handle AgentResult object from Strands
        if hasattr(ai_response, 'output'):
            ai_response = ai_response.output
            print(f"📝 Extracted output from AgentResult")
        elif hasattr(ai_response, 'text'):
            ai_response = ai_response.text
            print(f"📝 Extracted text from AgentResult")
        
        # Ensure we have a string response
        if not isinstance(ai_response, str):
            ai_response = str(ai_response)
        
        print(f"📝 AI Response length: {len(ai_response)} chars")
        print(f"📝 AI Response (first 500 chars): {ai_response[:500]}")

        # Strip markdown code blocks if present
        ai_response = ai_response.strip()
        if ai_response.startswith('```json'):
            ai_response = ai_response[7:]  # Remove ```json
        elif ai_response.startswith('```'):
            ai_response = ai_response[3:]  # Remove ```
        
        if ai_response.endswith('```'):
            ai_response = ai_response[:-3]  # Remove trailing ```
        
        ai_response = ai_response.strip()
        
        print(f"📝 Cleaned response length: {len(ai_response)} chars")
        print(f"📝 Cleaned response (first 200 chars): {ai_response[:200]}")
        print(f"📝 Cleaned response (last 200 chars): {ai_response[-200:]}")

        # Parse JSON response
        try:
            lesson_slides_data = json.loads(ai_response)
        except json.JSONDecodeError as json_err:
            print(f"❌ Failed to parse AI response as JSON: {json_err}")
            print(f"📄 Error position - Line: {json_err.lineno}, Column: {json_err.colno}, Char: {json_err.pos}")
            
            # Show context around the error
            if json_err.pos:
                start = max(0, json_err.pos - 200)
                end = min(len(ai_response), json_err.pos + 200)
                print(f"📄 Context around error:")
                print(f"   ...{ai_response[start:json_err.pos]}<<<ERROR>>>{ai_response[json_err.pos:end]}...")
            
            # Try to fix common JSON issues
            print(f"🔧 Attempting to fix common JSON issues...")
            
            # Replace smart quotes with regular quotes
            fixed_response = ai_response.replace('"', '"').replace('"', '"')
            fixed_response = fixed_response.replace(''', "'").replace(''', "'")
            
            # Try parsing again
            try:
                lesson_slides_data = json.loads(fixed_response)
                print(f"✅ Fixed JSON by replacing smart quotes!")
            except json.JSONDecodeError as json_err2:
                print(f"❌ Still failed after fixing quotes: {json_err2}")
                print(f"📄 Saving full response to /tmp/failed_json_lesson_{lesson_idx}.txt for debugging")
                
                # Save to file for debugging
                with open(f"/tmp/failed_json_lesson_{lesson_idx}.txt", "w") as f:
                    f.write(ai_response)
                
                raise ValueError(f"AI agent did not return valid JSON for lesson {lesson_idx}. JSON Error: {json_err}. Response saved to /tmp/failed_json_lesson_{lesson_idx}.txt")

        print(f"✅ AI generated {len(lesson_slides_data.get('slides', []))} slides for lesson {lesson_idx}")

        # Validate that we got slides
        if not lesson_slides_data.get('slides') or len(lesson_slides_data.get('slides', [])) == 0:
            raise ValueError(f"AI agent returned empty slides array for lesson {lesson_idx}")

        # Process and enhance the AI-generated slides
        for slide_data in lesson_slides_data.get('slides', []):
            slide_data['slide_number'] = slide_counter
            slide_data['lesson_number'] = lesson_idx
            slide_data['lesson_title'] = lesson_title

            # Ensure image references use available images
            if 'image_reference' in slide_data and available_images:
                # Map AI image references to actual available images
                img_ref = slide_data['image_reference']
                if 'USE_IMAGE:' in img_ref:
                    requested_img = img_ref.split('USE_IMAGE:')[1].strip()
                    print(f"🔗 Mapping AI image request: '{requested_img}' to available images")

                    # Find best matching available image that hasn't been used
                    best_match = None
                    best_score = 0

                    for img_idx, img in enumerate(available_images):
                        # Skip already used images
                        if img_idx in used_image_indices:
                            continue
                            
                        # Calculate similarity score
                        score = 0
                        if requested_img.lower() in img['alt_text'].lower():
                            score += 10
                        if requested_img.lower() == img['alt_text'].lower():
                            score += 20
                        if any(word in img['alt_text'].lower() for word in requested_img.lower().split()):
                            score += 5

                        if score > best_score:
                            best_score = score
                            best_match = (img_idx, img)
                    
                    # If no unused images with good score, try to find any unused image
                    if not best_match and available_images:
                        for img_idx, img in enumerate(available_images):
                            if img_idx not in used_image_indices:
                                best_match = (img_idx, img)
                                best_score = 1
                                break
                    
                    # If all images used, allow reuse but mark it
                    if not best_match and available_images:
                        best_match = (0, available_images[0])
                        best_score = 0
                        print(f"⚠️ All images used, reusing first image")

                    if best_match:
                        img_idx, img = best_match
                        used_image_indices.add(img_idx)
                        print(f"✅ Mapped to: {img['alt_text']} (score: {best_score}, index: {img_idx})")
                        
                        # If the image has no URL (from VISUAL tag), search S3 for it
                        if not img.get('url') or img['url'] is None:
                            print(f"🔍 Image has no URL, searching S3 for lesson images...")
                            if course_bucket and project_folder:
                                s3_image_url = None
                                try:
                                    # Try multiple possible path patterns
                                    # Based on actual S3 structure: {project}/images/XX-XX-XXXX.png
                                    search_paths = [
                                        f"{project_folder}/images/",  # Most common pattern
                                        f"{project_folder}/lessons/lesson_{lesson_idx}/images/",
                                        f"{project_folder}/lesson_{lesson_idx}/images/",
                                        f"{project_folder}/images/lesson_{lesson_idx}/",
                                    ]
                                    
                                    for search_path in search_paths:
                                        print(f"📁 Searching S3: s3://{course_bucket}/{search_path}")
                                        
                                        try:
                                            list_response = s3_client.list_objects_v2(
                                                Bucket=course_bucket,
                                                Prefix=search_path,
                                                MaxKeys=100
                                            )
                                            
                                            if 'Contents' in list_response and len(list_response['Contents']) > 0:
                                                print(f"✅ Found {len(list_response['Contents'])} objects in {search_path}")
                                                
                                                # Find images that match this lesson number
                                                # Image naming pattern: XX-YY-ZZZZ.png where XX is module/lesson number
                                                lesson_num_str = f"{lesson_idx:02d}"  # Format as 2-digit: 01, 02, etc.
                                                matching_images = []
                                                
                                                for s3_obj in list_response['Contents']:
                                                    key = s3_obj['Key']
                                                    if key.lower().endswith(('.png', '.jpg', '.jpeg')):
                                                        filename = key.split('/')[-1]
                                                        # Check if filename starts with lesson number (e.g., "01-")
                                                        if filename.startswith(f"{lesson_num_str}-"):
                                                            matching_images.append(key)
                                                            print(f"📸 Found matching image: {filename}")
                                                
                                                if matching_images:
                                                    # Use the first matching image for this lesson
                                                    # You could also randomly select or use img_idx to cycle through them
                                                    selected_key = matching_images[img_idx % len(matching_images)] if img_idx < len(matching_images) else matching_images[0]
                                                    s3_image_url = f"s3://{course_bucket}/{selected_key}"
                                                    print(f"✅ Selected S3 image: {s3_image_url}")
                                                    slide_data['image_url'] = s3_image_url
                                                    img['url'] = s3_image_url  # Update the image dict too
                                                    break
                                                else:
                                                    print(f"⚠️ No images matching lesson {lesson_num_str} pattern found")
                                            else:
                                                print(f"⚠️ No objects found at {search_path}")
                                        except Exception as path_err:
                                            print(f"⚠️ Search failed for {search_path}: {path_err}")
                                            continue
                                    
                                    if not s3_image_url:
                                        print(f"⚠️ No images found in any search path for lesson {lesson_idx}")
                                except Exception as s3_err:
                                    print(f"❌ S3 search failed: {s3_err}")
                        else:
                            slide_data['image_url'] = img['url']
                        
                        slide_data['image_reference'] = f"USE_IMAGE: {img['alt_text']}"

            # Only search S3 for images if slide explicitly needs one but has no URL yet
            # AND we haven't exhausted our available images
            if slide_data.get('slide_type') in ['image'] and (not slide_data.get('image_url') or not slide_data['image_url']):
                if available_images and len(used_image_indices) < len(available_images):
                    # Use next unused image from available_images
                    for img_idx, img in enumerate(available_images):
                        if img_idx not in used_image_indices:
                            used_image_indices.add(img_idx)
                            slide_data['image_url'] = img['url']
                            slide_data['image_reference'] = f"USE_IMAGE: {img['alt_text']}"
                            print(f"🔗 Assigned unused image: {img['alt_text']} (index: {img_idx})")
                            break

            all_slides.append(slide_data)
            slide_counter += 1

        print(f"📈 Total slides so far: {len(all_slides)}")

    # Add final summary slide
    all_slides.append({
        "slide_number": slide_counter,
        "slide_type": "summary",
        "title": "Course Summary & Next Steps",
        "bullets": [
            f"Completed {len(lessons)} comprehensive lessons",
            "Developed practical skills for immediate application",
            "Ready to apply knowledge in professional scenarios",
            "Continue learning with advanced topics and certifications"
        ],
        "notes": "Thank participants and discuss next steps and resources"
    })

    presentation_structure = {
        "presentation_title": course_title,
        "total_slides": len(all_slides),
        "total_lessons": len(lessons),
        "style": presentation_style,
        "generated_at": datetime.now().isoformat(),
        "slides": all_slides
    }

    print(f"\n✅ AI-powered presentation structure complete: {len(all_slides)} slides")
    return presentation_structure


def generate_pptx_file(presentation_structure: Dict, book_data: Dict) -> bytes:
    """
    Generate actual PPTX file from the presentation structure.
    Uses python-pptx library.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        import io
        import requests
        from urllib.parse import urlparse
        
        print("🎨 Creating PowerPoint file...")
        
        prs = Presentation()
        # Set 16:9 widescreen format (standard for modern displays)
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)
        # cache slide dimensions for layout helpers
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # small helper to horizontally center a shape and set a relative width
        def center_shape_horizontally(shape, width_ratio: float = 0.8):
            """Set shape.width to width_ratio * slide_width and center it horizontally."""
            try:
                desired_width = int(slide_width * width_ratio)
                shape.width = desired_width
                shape.left = int((slide_width - desired_width) // 2)
            except Exception:
                # If shape doesn't support width/left or values fail, ignore silently
                pass

        # In-memory image cache to avoid double downloads: {url: bytes}
        image_cache: Dict[str, bytes] = {}

        # Helpers for adaptive layout decisions
        def estimate_text_word_count(slide_data: Dict) -> int:
            """Count words in title, bullets and caption to estimate text length."""
            words = 0
            title = slide_data.get('title', '') or ''
            caption = slide_data.get('caption', '') or ''
            bullets = slide_data.get('bullets', []) or []
            words += len(title.split())
            words += len(caption.split())
            for b in bullets:
                words += len(str(b).split())
            return words

        # Delegate robust image tasks to image_manager
        try:
            from . import image_manager
        except Exception:
            # fallback to local helper if relative import fails (module may be run as script)
            try:
                import image_manager
            except Exception:
                image_manager = None

        def get_image_dimensions(image_url: str) -> Optional[Dict[str, int]]:
            """Try to retrieve image and return its pixel dimensions {width, height}.
            Returns None on failure.
            """
            # If we have the helper module, use it (it handles caching and Pillow checks)
            if image_manager is not None:
                try:
                    # Try to use cache first
                    if image_url in image_cache:
                        dims = image_manager.get_image_size_from_bytes(image_cache[image_url])
                        if dims:
                            return {'width': dims[0], 'height': dims[1]}

                    raw = image_manager.fetch_image_bytes(image_url, s3_client=s3_client)
                    if not raw:
                        return None
                    # store in cache
                    image_cache[image_url] = raw
                    dims = image_manager.get_image_size_from_bytes(raw)
                    if dims:
                        return {'width': dims[0], 'height': dims[1]}
                    return None
                except Exception:
                    return None

            # Fallback: no image_manager available -> return None
            return None

        def decide_layout_for_slide(slide_data: Dict) -> str:
            """Decide layout: 'image-only', 'split-right', 'split-top', or 'text-only'."""
            WORD_SHORT = 35
            WORD_MEDIUM = 100

            text_words = estimate_text_word_count(slide_data)
            image_url = slide_data.get('image_url')

            if not image_url:
                return 'text-only'

            dims = get_image_dimensions(image_url)
            # If we couldn't get dimensions, be conservative and use split-right
            if not dims:
                if text_words <= WORD_SHORT:
                    return 'image-only'
                return 'split-right'

            w = dims.get('width', 0)
            h = dims.get('height', 1)
            if h == 0:
                h = 1
            aspect = w / h

            # Decision rules
            if text_words <= WORD_SHORT:
                # Favor large visual slide when text is short
                return 'image-only'
            if text_words <= WORD_MEDIUM:
                # Medium text: if image is wide, place on top; if tall or square, place on right
                if aspect >= 1.3:
                    return 'split-top'
                return 'split-right'
            # Long text -> text only (split into more slides upstream)
            return 'text-only'
        
        # Define slide layouts
        title_slide_layout = prs.slide_layouts[0]  # Title Slide
        content_slide_layout = prs.slide_layouts[1]  # Title and Content
        blank_slide_layout = prs.slide_layouts[6]  # Blank
        
        style = presentation_structure.get('style', 'professional')
        
        # Color schemes by style
        color_schemes = {
            'professional': {
                'primary': RGBColor(0, 51, 102),  # Dark Blue
                'secondary': RGBColor(68, 114, 196),  # Medium Blue
                'accent': RGBColor(255, 192, 0)  # Gold
            },
            'educational': {
                'primary': RGBColor(46, 139, 87),  # Sea Green
                'secondary': RGBColor(70, 130, 180),  # Steel Blue
                'accent': RGBColor(255, 140, 0)  # Dark Orange
            },
            'modern': {
                'primary': RGBColor(33, 33, 33),  # Dark Gray
                'secondary': RGBColor(96, 125, 139),  # Blue Gray
                'accent': RGBColor(0, 150, 136)  # Teal
            }
        }
        
        colors = color_schemes.get(style, color_schemes['professional'])
        
        # Process each slide
        for slide_data in presentation_structure.get('slides', []):
            slide_type = slide_data.get('slide_type', 'content')
            slide_title = slide_data.get('title', '')
            
            # Skip blank slides (no content, no image, no bullets)
            if slide_type == 'content':
                has_bullets = slide_data.get('bullets') and len(slide_data.get('bullets', [])) > 0
                has_image = slide_data.get('image_url') or slide_data.get('image_reference')
                has_title = slide_title and slide_title.strip()
                
                if not has_bullets and not has_image and not has_title:
                    print(f"⚠️ Skipping blank slide (no content, no image, no title)")
                    continue
                
                # Skip slides with title but no bullets and no image
                if has_title and not has_bullets and not has_image:
                    print(f"⚠️ Skipping content slide with only title: '{slide_title}'")
                    continue

            if slide_type == 'title':
                slide = prs.slides.add_slide(title_slide_layout)
                
                # First set the title and subtitle text
                title_shape = slide.shapes.title
                subtitle = slide.placeholders[1]
                title_shape.text = slide_title
                subtitle.text = slide_data.get('subtitle', '')
                
                # Enhanced title formatting
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(54)
                title_para.font.bold = True
                title_para.font.color.rgb = colors['primary']  # Use primary color instead of white
                title_para.alignment = PP_ALIGN.CENTER

                # Center the title and subtitle shapes horizontally on the slide
                center_shape_horizontally(title_shape, width_ratio=0.82)
                try:
                    if subtitle is not None:
                        center_shape_horizontally(subtitle, width_ratio=0.6)
                except Exception:
                    pass
                
                # Subtitle formatting
                if subtitle.has_text_frame:
                    subtitle_para = subtitle.text_frame.paragraphs[0]
                    subtitle_para.font.size = Pt(24)
                    subtitle_para.font.color.rgb = colors['secondary']
                    subtitle_para.alignment = PP_ALIGN.CENTER
            elif slide_type == 'content':
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Add accent bar on left side
                accent_bar = slide.shapes.add_shape(
                    1,  # Rectangle
                    Inches(0), Inches(0), Inches(0.3), Inches(7.5)
                )
                fill = accent_bar.fill
                fill.solid()
                fill.fore_color.rgb = colors['accent']
                
                # Title with enhanced styling
                title_shape = slide.shapes.title
                title_shape.text = slide_title
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(28)  # SMALLER, professional title
                title_para.font.bold = True
                title_para.font.color.rgb = colors['primary']
                # Ensure the title is visually centered (placeholder defaults can be off)
                title_para.alignment = PP_ALIGN.CENTER
                center_shape_horizontally(title_shape, width_ratio=0.85)
                
                # Add subtle background to title
                title_shape.fill.solid()
                title_shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
                
                # Adaptive layout selection based on text length and image dimensions
                image_url = slide_data.get('image_url')
                layout = decide_layout_for_slide(slide_data) if image_url else 'text-only'

                # Default assignments
                content_left = Inches(2.0)
                content_width = Inches(9.5)
                image_left = None
                image_top = None
                image_width = None
                image_height = None

                if layout == 'text-only':
                    has_image = False
                    content_left = Inches(2.0)
                    content_width = Inches(9.5)
                elif layout == 'image-only':
                    has_image = True
                    # Large full-width image with short caption below
                    content_left = Inches(2.0)
                    content_width = Inches(9.5)
                    image_left = Inches(1.0)
                    image_top = Inches(1.6)
                    image_width = Inches(11.3)
                    image_height = Inches(4.8)
                elif layout == 'split-top':
                    has_image = True
                    # Top image spanning the page, text below
                    image_left = Inches(1.0)
                    image_top = Inches(1.2)
                    image_width = Inches(11.3)
                    image_height = Inches(3.2)
                    content_left = Inches(1.0)
                    content_width = Inches(9.0)
                else:  # split-right and default conservative split
                    has_image = True
                    content_left = Inches(0.6)
                    content_width = Inches(5.8)
                    image_left = Inches(7.0)
                    image_top = Inches(2.0)
                    image_width = Inches(5.5)
                    image_height = Inches(4.5)
                
                # Content body with conservative spacing
                body = slide.placeholders[1]
                body.left = content_left
                body.top = Inches(1.8)
                body.width = content_width
                body.height = Inches(5.2)
                
                text_frame = body.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                
                # STRICT bullet limits for professional appearance
                bullets = slide_data.get('bullets', [])
                max_bullets = 5 if has_image else 6  # STRICT LIMITS
                bullets_to_show = bullets[:max_bullets]
                
                for i, bullet in enumerate(bullets_to_show):
                    p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    # Smaller, professional font sizes
                    p.font.size = Pt(14) if has_image else Pt(16)  # SMALLER FONTS
                    p.font.name = 'Calibri'
                    # Tight spacing for clean look
                    p.space_before = Pt(6) if has_image else Pt(8)
                    p.space_after = Pt(4) if has_image else Pt(6)
                    
                    # Alternate bullet colors for visual interest
                    if i % 2 == 0:
                        p.font.color.rgb = RGBColor(33, 37, 41)
                    else:
                        p.font.color.rgb = RGBColor(52, 58, 64)
                    
                    # Use custom bullet character
                    p.text = f"▸ {bullet}"
                
                # Add warning if bullets were truncated
                if len(bullets) > max_bullets:
                    print(f"⚠️ Truncated {len(bullets)} bullets to {max_bullets} for better layout")
                
                # Add image if referenced (on the right side)
                image_url = slide_data.get('image_url')
                if image_url and has_image:
                    try:
                        img_data = None
                        # Try cache first
                        if image_url in image_cache:
                            try:
                                img_data = io.BytesIO(image_cache[image_url])
                            except Exception:
                                img_data = None

                        # If not cached, download and cache
                        if not img_data:
                            if 's3.amazonaws.com' in image_url or image_url.startswith('s3://'):
                                try:
                                    # Parse S3 URL
                                    if image_url.startswith('s3://'):
                                        parts = image_url[5:].split('/', 1)
                                        bucket = parts[0]
                                        key = parts[1] if len(parts) > 1 else ''
                                    else:
                                        url_parts = image_url.replace('https://', '').replace('http://', '').split('/')
                                        bucket = url_parts[0].split('.')[0]
                                        key = '/'.join(url_parts[1:])

                                    print(f"🔄 Retrieving from S3: s3://{bucket}/{key}")
                                    s3_response = s3_client.get_object(Bucket=bucket, Key=key)
                                    raw = s3_response['Body'].read()
                                    img_data = io.BytesIO(raw)
                                    image_cache[image_url] = raw
                                    print(f"✅ Retrieved image from S3: {bucket}/{key}")
                                except Exception as s3_error:
                                    print(f"❌ S3 access failed: {s3_error}")
                                    # Fall back to HTTP
                                    if image_url.startswith('http'):
                                        try:
                                            img_resp = requests.get(image_url, timeout=10)
                                            if img_resp.status_code == 200:
                                                raw = img_resp.content
                                                img_data = io.BytesIO(raw)
                                                image_cache[image_url] = raw
                                                print(f"✅ Downloaded via HTTP: {image_url}")
                                        except Exception:
                                            pass
                            elif image_url.startswith('http'):
                                # Regular HTTP URL
                                try:
                                    img_resp = requests.get(image_url, timeout=10)
                                    if img_resp.status_code == 200:
                                        raw = img_resp.content
                                        img_data = io.BytesIO(raw)
                                        image_cache[image_url] = raw
                                        print(f"✅ Downloaded image: {image_url}")
                                except requests.exceptions.RequestException as e:
                                    print(f"⚠️ Network error downloading image: {e}")

                        if img_data:
                            # Add image with optimized size and position (prevents overwhelming)
                            # Use the variables set earlier based on layout
                            if image_left is not None:
                                try:
                                    # Ensure we have raw bytes
                                    raw_bytes = None
                                    if isinstance(img_data, io.BytesIO):
                                        raw_bytes = img_data.getvalue()
                                    elif isinstance(img_data, (bytes, bytearray)):
                                        raw_bytes = bytes(img_data)

                                    pic_stream = None
                                    DPI = 150  # target DPI for resizing (tunable)
                                    if raw_bytes and image_manager is not None:
                                        try:
                                            max_w = int(image_width.inches * DPI)
                                            max_h = int(image_height.inches * DPI)
                                            resized = image_manager.resize_image_bytes(raw_bytes, max_w, max_h, fmt='PNG')
                                            if resized:
                                                pic_stream = io.BytesIO(resized)
                                        except Exception:
                                            pic_stream = None

                                    # If we couldn't resize, use the original stream
                                    if not pic_stream:
                                        try:
                                            pic_stream = img_data if isinstance(img_data, io.BytesIO) else io.BytesIO(raw_bytes)
                                        except Exception:
                                            pic_stream = img_data

                                    pic = slide.shapes.add_picture(pic_stream, image_left, image_top, width=image_width)

                                    # If the image is too tall after width constraint, adjust it
                                    if pic.height > image_height:
                                        # Recalculate to fit height instead
                                        aspect_ratio = pic.width / pic.height
                                        pic.height = image_height
                                        pic.width = int(image_height * aspect_ratio)
                                        # Re-center horizontally if narrower
                                        if pic.width < image_width:
                                            pic.left = image_left + (image_width - pic.width) // 2

                                    print(f"✅ Inserted image: {pic.width.inches:.2f}\"x{pic.height.inches:.2f}\" at ({pic.left.inches:.2f}\", {pic.top.inches:.2f}\")")
                                except Exception as e:
                                    print(f"⚠️ Error inserting/resizing image: {e}")
                        else:
                            print(f"⚠️ Could not access image, adding visual placeholder")
                            # Try to create a professional-looking placeholder image via image_manager
                            if image_left is not None and image_manager is not None:
                                try:
                                    placeholder_text = slide_data.get('image_reference') or slide_title or 'Visual'
                                    # Create placeholder at target pixel size using DPI
                                    DPI = 150
                                    px_w = int(image_width.inches * DPI)
                                    px_h = int(image_height.inches * DPI)
                                    ph = image_manager.create_placeholder_image(text=placeholder_text, size=(px_w, px_h))
                                    if ph:
                                        pic_stream = io.BytesIO(ph)
                                        pic = slide.shapes.add_picture(pic_stream, image_left, image_top, width=image_width)
                                        try:
                                            pic.name = 'aurora_placeholder'
                                        except Exception:
                                            pass
                                    else:
                                        raise Exception('placeholder generation returned None')
                                except Exception as pe:
                                    print(f"⚠️ Placeholder image creation failed: {pe}. Attempting inline PNG creation or shape fallback.")
                                    # Try to create a minimal PNG inline using Pillow if available
                                    inline_ph = None
                                    try:
                                        from PIL import Image, ImageDraw, ImageFont
                                        pw = int(image_width.inches * DPI)
                                        ph_h = int(image_height.inches * DPI)
                                        im = Image.new('RGB', (max(1, pw), max(1, ph_h)), color=(245, 247, 250))
                                        draw = ImageDraw.Draw(im)
                                        try:
                                            font = ImageFont.truetype('DejaVuSans-Bold.ttf', 24)
                                        except Exception:
                                            try:
                                                font = ImageFont.load_default()
                                            except Exception:
                                                font = None
                                        text = placeholder_text if placeholder_text else 'Visual'
                                        try:
                                            bbox = draw.textbbox((0, 0), text, font=font)
                                            tw = bbox[2] - bbox[0]
                                            th = bbox[3] - bbox[1]
                                        except Exception:
                                            try:
                                                tw, th = font.getsize(text) if font else (int(pw * 0.6), 24)
                                            except Exception:
                                                tw, th = (int(pw * 0.6), 24)
                                        tx = max(0, (pw - tw) // 2)
                                        ty = max(0, (ph_h - th) // 2)
                                        try:
                                            draw.text((tx, ty), text, fill=(108, 117, 125), font=font)
                                        except Exception:
                                            pass
                                        out = io.BytesIO()
                                        im.save(out, format='PNG')
                                        inline_ph = out.getvalue()
                                    except Exception:
                                        inline_ph = None

                                    if inline_ph:
                                        try:
                                            pic_stream = io.BytesIO(inline_ph)
                                            pic = slide.shapes.add_picture(pic_stream, image_left, image_top, width=image_width)
                                            try:
                                                pic.name = 'aurora_placeholder'
                                            except Exception:
                                                pass
                                        except Exception:
                                            inline_ph = None

                                    # Final fallback: shape placeholder
                                    if not inline_ph:
                                        placeholder = slide.shapes.add_shape(
                                            1,  # Rectangle
                                            image_left, image_top, image_width, image_height
                                        )
                                        fill = placeholder.fill
                                        fill.solid()
                                        fill.fore_color.rgb = RGBColor(240, 240, 240)
                                        line = placeholder.line
                                        line.color.rgb = colors['secondary']
                                        line.width = Pt(2)
                                        text_frame = placeholder.text_frame
                                        text_frame.text = "📊 Visual [PLACEHOLDER]"
                                        try:
                                            placeholder.name = 'aurora_placeholder'
                                        except Exception:
                                            pass
                                        p = text_frame.paragraphs[0]
                                        p.alignment = PP_ALIGN.CENTER
                                        p.font.size = Pt(14)
                                        p.font.color.rgb = colors['secondary']
                            elif image_left is not None:
                                # No image_manager available; draw a simple rectangle placeholder
                                placeholder = slide.shapes.add_shape(
                                    1,  # Rectangle
                                    image_left, image_top, image_width, image_height
                                )
                                fill = placeholder.fill
                                fill.solid()
                                fill.fore_color.rgb = RGBColor(240, 240, 240)
                                line = placeholder.line
                                line.color.rgb = colors['secondary']
                                line.width = Pt(2)
                                text_frame = placeholder.text_frame
                                text_frame.text = "📊 Visual [PLACEHOLDER]"
                                try:
                                    placeholder.name = 'aurora_placeholder'
                                except Exception:
                                    pass
                                p = text_frame.paragraphs[0]
                                p.alignment = PP_ALIGN.CENTER
                                p.font.size = Pt(14)
                                p.font.color.rgb = colors['secondary']
                    except Exception as e:
                        print(f"⚠️ Error inserting image: {e}")
            elif slide_type == 'image':
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add gradient background
                background = slide.shapes.add_shape(
                    1,  # Rectangle
                    Inches(0), Inches(0), Inches(13.333), Inches(7.5)
                )
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(245, 247, 250)
                
                # Title with modern styling
                # Create a centered textbox for the title to ensure perfect centering
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(1))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = slide_data.get('title', '')
                p.font.size = Pt(40)
                p.font.bold = True
                p.font.color.rgb = colors['primary']
                p.alignment = PP_ALIGN.CENTER
                # Center the textbox horizontally using our helper
                center_shape_horizontally(txBox, width_ratio=0.92)
                
                # Add image with border and shadow
                image_url = slide_data.get('image_url')
                if image_url:
                    try:
                        img_data = None
                        # Try cache first
                        if image_url in image_cache:
                            try:
                                img_data = io.BytesIO(image_cache[image_url])
                            except Exception:
                                img_data = None

                        # If not cached, download and cache
                        if not img_data:
                            if 's3.amazonaws.com' in image_url or image_url.startswith('s3://'):
                                try:
                                    if image_url.startswith('s3://'):
                                        parts = image_url[5:].split('/', 1)
                                        bucket = parts[0]
                                        key = parts[1] if len(parts) > 1 else ''
                                    else:
                                        url_parts = image_url.replace('https://', '').replace('http://', '').split('/')
                                        bucket = url_parts[0].split('.')[0]
                                        key = '/'.join(url_parts[1:])

                                    print(f"🔄 Retrieving image from S3: s3://{bucket}/{key}")
                                    s3_response = s3_client.get_object(Bucket=bucket, Key=key)
                                    raw = s3_response['Body'].read()
                                    img_data = io.BytesIO(raw)
                                    image_cache[image_url] = raw
                                    print(f"✅ Retrieved image from S3")
                                except Exception as s3_error:
                                    print(f"❌ S3 access failed: {s3_error}")

                        if img_data:
                            # Center the image
                            left = Inches(1.5)
                            top = Inches(1.5)
                            width = Inches(7)
                            try:
                                raw_bytes = None
                                if isinstance(img_data, io.BytesIO):
                                    raw_bytes = img_data.getvalue()
                                elif isinstance(img_data, (bytes, bytearray)):
                                    raw_bytes = bytes(img_data)

                                pic_stream = None
                                DPI = 150
                                if raw_bytes and image_manager is not None:
                                    try:
                                        max_w = int(width.inches * DPI)
                                        max_h = int((Inches(4).inches) * DPI)  # approximate height limit
                                        resized = image_manager.resize_image_bytes(raw_bytes, max_w, max_h, fmt='PNG')
                                        if resized:
                                            pic_stream = io.BytesIO(resized)
                                    except Exception:
                                        pic_stream = None

                                if not pic_stream:
                                    try:
                                        pic_stream = img_data if isinstance(img_data, io.BytesIO) else io.BytesIO(raw_bytes)
                                    except Exception:
                                        pic_stream = img_data

                                pic = slide.shapes.add_picture(pic_stream, left, top, width=width)
                                print(f"✅ Inserted image into slide")
                            except Exception as e:
                                print(f"⚠️ Error inserting image into image slide: {e}")
                        else:
                            # Professional placeholder - try to generate via image_manager where possible
                            if image_manager is not None:
                                try:
                                    left = Inches(2)
                                    top = Inches(2)
                                    width = Inches(6)
                                    height = Inches(4)
                                    DPI = 150
                                    px_w = int(width.inches * DPI)
                                    px_h = int(height.inches * DPI)
                                    ph = image_manager.create_placeholder_image(text=slide_data.get('caption', 'Diagram'), size=(px_w, px_h))
                                    if ph:
                                        pic_stream = io.BytesIO(ph)
                                        pic = slide.shapes.add_picture(pic_stream, left, top, width=width)
                                        try:
                                            pic.name = 'aurora_placeholder'
                                        except Exception:
                                            pass
                                    else:
                                        raise Exception('placeholder returned None')
                                except Exception as e:
                                    print(f"⚠️ Failed to create placeholder image for image slide: {e}. Attempting inline PNG creation or shape fallback.")
                                    inline_ph = None
                                    try:
                                        from PIL import Image, ImageDraw, ImageFont
                                        pw = int(width.inches * DPI)
                                        ph_h = int(height.inches * DPI)
                                        im = Image.new('RGB', (max(1, pw), max(1, ph_h)), color=(245, 247, 250))
                                        draw = ImageDraw.Draw(im)
                                        try:
                                            font = ImageFont.truetype('DejaVuSans-Bold.ttf', 36)
                                        except Exception:
                                            try:
                                                font = ImageFont.load_default()
                                            except Exception:
                                                font = None
                                        text = slide_data.get('caption', 'Diagram')
                                        try:
                                            bbox = draw.textbbox((0, 0), text, font=font)
                                            tw = bbox[2] - bbox[0]
                                            th = bbox[3] - bbox[1]
                                        except Exception:
                                            try:
                                                tw, th = font.getsize(text) if font else (int(pw * 0.6), 24)
                                            except Exception:
                                                tw, th = (int(pw * 0.6), 24)
                                        tx = max(0, (pw - tw) // 2)
                                        ty = max(0, (ph_h - th) // 2)
                                        try:
                                            draw.text((tx, ty), text, fill=(108, 117, 125), font=font)
                                        except Exception:
                                            pass
                                        out = io.BytesIO()
                                        im.save(out, format='PNG')
                                        inline_ph = out.getvalue()
                                    except Exception:
                                        inline_ph = None

                                    if inline_ph:
                                        try:
                                            pic_stream = io.BytesIO(inline_ph)
                                            pic = slide.shapes.add_picture(pic_stream, left, top, width=width)
                                            try:
                                                pic.name = 'aurora_placeholder'
                                            except Exception:
                                                pass
                                        except Exception:
                                            inline_ph = None

                                    if not inline_ph:
                                        # fallback to a simple shape
                                        left = Inches(2)
                                        top = Inches(2)
                                        width = Inches(6)
                                        height = Inches(4)
                                        placeholder = slide.shapes.add_shape(1, left, top, width, height)
                                        fill = placeholder.fill
                                        fill.solid()
                                        fill.fore_color.rgb = RGBColor(255, 255, 255)
                                        line = placeholder.line
                                        line.color.rgb = colors['secondary']
                                        line.width = Pt(3)
                                        text_frame = placeholder.text_frame
                                        text_frame.text = "📊 Diagram [PLACEHOLDER]"
                                        try:
                                            placeholder.name = 'aurora_placeholder'
                                        except Exception:
                                            pass
                                        p = text_frame.paragraphs[0]
                                        p.alignment = PP_ALIGN.CENTER
                                        p.font.size = Pt(32)
                                        p.font.color.rgb = colors['secondary']
                            else:
                                left = Inches(2)
                                top = Inches(2)
                                width = Inches(6)
                                height = Inches(4)
                                placeholder = slide.shapes.add_shape(1, left, top, width, height)
                                fill = placeholder.fill
                                fill.solid()
                                fill.fore_color.rgb = RGBColor(255, 255, 255)
                                line = placeholder.line
                                line.color.rgb = colors['secondary']
                                line.width = Pt(3)
                                text_frame = placeholder.text_frame
                                text_frame.text = "📊 Diagram [PLACEHOLDER]"
                                try:
                                    placeholder.name = 'aurora_placeholder'
                                except Exception:
                                    pass
                                p = text_frame.paragraphs[0]
                                p.alignment = PP_ALIGN.CENTER
                                p.font.size = Pt(32)
                                p.font.color.rgb = colors['secondary']
                    except Exception as e:
                        print(f"⚠️ Error inserting image: {e}")
                
                # Add caption with modern styling
                caption = slide_data.get('caption', '')
                if caption:
                    caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
                    tf2 = caption_box.text_frame
                    p2 = tf2.paragraphs[0]
                    p2.text = caption
                    p2.font.size = Pt(16)
                    p2.font.italic = True
                    p2.font.color.rgb = RGBColor(108, 117, 125)
                    p2.alignment = PP_ALIGN.CENTER
            elif slide_type == 'summary':
                slide = prs.slides.add_slide(content_slide_layout)
                
                # Add decorative background elements
                accent_shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    Inches(0), Inches(6.5), Inches(13.333), Inches(1)
                )
                fill = accent_shape.fill
                fill.solid()
                fill.fore_color.rgb = colors['accent']
                
                # Title with icon
                title_shape = slide.shapes.title
                title_shape.text = slide_title or 'Key Takeaways'
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.font.size = Pt(40)
                title_para.font.bold = True
                title_para.font.color.rgb = colors['primary']
                # Center the title placeholder horizontally
                title_para.alignment = PP_ALIGN.CENTER
                center_shape_horizontally(title_shape, width_ratio=0.85)
                
                # Summary points with checkmarks
                body = slide.placeholders[1]
                text_frame = body.text_frame
                text_frame.clear()
                
                for i, bullet in enumerate(slide_data.get('bullets', [])):
                    p = text_frame.add_paragraph()
                    p.text = f"✓ {bullet}"
                    p.level = 0
                    p.font.size = Pt(22)
                    p.font.bold = True
                    p.font.color.rgb = colors['primary']
                    p.space_before = Pt(18)
                    p.space_after = Pt(12)
        
        # Save to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        print(f"✅ PowerPoint file created: {len(presentation_structure.get('slides', []))} slides")
        return pptx_buffer.getvalue()
        
    except ImportError as e:
        print(f"❌ Missing python-pptx library: {e}")
        print("💡 Install with: pip install python-pptx")
        raise
    except Exception as e:
        print(f"❌ Error creating PPTX file: {e}")
        raise


def lambda_handler(event, context):
    """
    Main Lambda handler for PPT generation.
    """
    try:
        print("=" * 80)
        print("🎯 STRANDS PPT GENERATOR")
        print("=" * 80)
        print(f"📥 Event: {json.dumps(event, indent=2)}")
        
        # Parse event - handle API Gateway body format
        if 'body' in event:
            # API Gateway POST request - body is a JSON string
            body = json.loads(event['body']) if isinstance(event['body'], str) else event['body']
        else:
            # Direct Lambda invocation
            body = event
        
        print(f"📦 Parsed body: {json.dumps(body, indent=2)}")
        
        # Extract parameters
        course_bucket = body.get('course_bucket')
        project_folder = body.get('project_folder')
        book_version_key = body.get('book_version_key')
        book_type = body.get('book_type', 'theory').lower()  # 'theory' or 'lab'
        model_provider = body.get('model_provider', 'bedrock').lower()
        slides_per_lesson = int(body.get('slides_per_lesson', 6))
        presentation_style = body.get('presentation_style', 'professional')
        
        if not course_bucket or not project_folder:
            raise ValueError("course_bucket and project_folder are required")
        
        # Auto-discover book file if no version specified (filter by book_type)
        if not book_version_key:
            print(f"\n🔍 Auto-discovering {book_type} book file in {project_folder}/book/")
            book_prefix = f"{project_folder}/book/"
            response = s3_client.list_objects_v2(
                Bucket=course_bucket,
                Prefix=book_prefix,
                MaxKeys=50
            )
            
            if 'Contents' in response:
                # Look for files ending with _data.json, filter by book_type
                json_files = []
                for obj in response['Contents']:
                    key = obj['Key']
                    last_modified = obj.get('LastModified')
                    filename = key.split('/')[-1]
                    
                    if filename.endswith('_data.json'):
                        # Filter by book type
                        is_lab_guide = 'Lab_Guide' in filename or 'LabGuide' in filename
                        
                        if book_type == 'lab' and is_lab_guide:
                            json_files.append((key, last_modified))
                            print(f"   📋 Found lab guide: {filename}")
                        elif book_type == 'theory' and not is_lab_guide:
                            json_files.append((key, last_modified))
                            print(f"   📚 Found theory book: {filename}")
                
                if json_files:
                    # Sort by last modified (most recent first)
                    json_files.sort(key=lambda x: x[1], reverse=True)
                    book_version_key = json_files[0][0]
                    print(f"✅ Auto-discovered {book_type} book: {book_version_key}")
                else:
                    raise ValueError(f"No {book_type} book files found in {book_prefix}. Found {len(response.get('Contents', []))} total files.")
            else:
                raise ValueError(f"No files found in {book_prefix}")
        
        print(f"\n📚 Book version: {book_version_key}")
        print(f"🤖 Model provider: {model_provider}")
        print(f"📊 Slides per lesson: {slides_per_lesson}")
        print(f"🎨 Style: {presentation_style}")
        
        # Configure AI model
        if model_provider == 'openai':
            model = configure_openai_model()
        else:
            model = configure_bedrock_model()
        
        # Load book data
        book_data = load_book_from_s3(course_bucket, book_version_key)
        
        # Generate presentation structure
        presentation_structure = generate_presentation_structure(
            book_data,
            model,
            slides_per_lesson,
            presentation_style,
            course_bucket,
            project_folder
        )
        
        # Save presentation structure as JSON
        structure_key = f"{project_folder}/presentations/presentation_structure.json"
        s3_client.put_object(
            Bucket=course_bucket,
            Key=structure_key,
            Body=json.dumps(presentation_structure, indent=2),
            ContentType='application/json'
        )
        print(f"💾 Saved structure: s3://{course_bucket}/{structure_key}")
        
        # Generate PPTX file
        try:
            print("🎨 Starting PPTX file generation...")
            print(f"📊 Presentation structure has {len(presentation_structure.get('slides', []))} slides")
            
            pptx_bytes = generate_pptx_file(presentation_structure, book_data)
            
            print(f"✅ PPTX file generated: {len(pptx_bytes)} bytes")
            
            # Save PPTX to S3 in the presentations folder
            pptx_title = presentation_structure.get('presentation_title', 'presentation')
            # Sanitize only the filename, not the folder path
            sanitized_filename = pptx_title.replace(' ', '_').replace('/', '-').replace(':', '-').replace('\\', '-')
            pptx_key = f"{project_folder}/presentations/{sanitized_filename}.pptx"
            
            print(f"📁 Saving to: {pptx_key}")
            
            s3_client.put_object(
                Bucket=course_bucket,
                Key=pptx_key,
                Body=pptx_bytes,
                ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            print(f"💾 Saved PPTX: s3://{course_bucket}/{pptx_key}")
            
        except Exception as e:
            error_details = {
                'error_type': type(e).__name__,
                'error_message': str(e),
                'error_args': str(e.args) if hasattr(e, 'args') else None
            }
            print(f"⚠️ PPTX generation failed (structure still saved)")
            print(f"Error details: {json.dumps(error_details, indent=2)}")
            import traceback
            print("Full traceback:")
            traceback.print_exc()
            pptx_key = None
        
        # Return success response with CORS headers
        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Headers': 'Content-Type,X-Amz-Date,Authorization,X-Api-Key,X-Amz-Security-Token',
                'Access-Control-Allow-Methods': 'POST,OPTIONS',
                'Content-Type': 'application/json'
            },
            'body': json.dumps({
                'message': 'Presentation generated successfully',
                'presentation_title': presentation_structure['presentation_title'],
                'total_slides': presentation_structure['total_slides'],
                'structure_s3_key': structure_key,
                'pptx_s3_key': pptx_key,
                'generated_at': presentation_structure['generated_at']
            })
        }
        
    except Exception as e:
        error_msg = f"Error generating presentation: {str(e)}"
        print(f"❌ {error_msg}")
        import traceback
        traceback.print_exc()
        
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Headers': 'Content-Type,X-Amz-Date,Authorization,X-Api-Key,X-Amz-Security-Token',
                'Access-Control-Allow-Methods': 'POST,OPTIONS',
                'Content-Type': 'application/json'
            },
            'body': json.dumps({
                'error': error_msg,
                'traceback': traceback.format_exc()
            })
        }
