import io
import os
import sys

# allow import when running from repo root
ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)

from pptx import Presentation
from strands_ppt_generator import generate_pptx_file


def rects_overlap(a, b) -> bool:
    a_l, a_t, a_w, a_h = a
    b_l, b_t, b_w, b_h = b
    a_r = a_l + a_w
    a_b = a_t + a_h
    b_r = b_l + b_w
    b_b = b_t + b_h
    horiz = not (a_r <= b_l or b_r <= a_l)
    vert = not (a_b <= b_t or b_b <= a_t)
    return horiz and vert


def test_no_overlap_after_generation():
    # Create a content slide with long text and a large image that would normally overlap
    presentation_structure = {
        'presentation_title': 'Overlap Test',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'content',
                'title': 'Long Text with Image',
                'bullets': ['word'] * 200,
                'image_url': 's3://nonexistent/huge.png'
            }
        ],
        'style': 'professional'
    }

    pptx_bytes = generate_pptx_file(presentation_structure, {})
    prs = Presentation(io.BytesIO(pptx_bytes))
    # inspect slide shapes to ensure no shape overlaps the body placeholder
    slide = prs.slides[0]
    # find body placeholder rectangle
    body = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            body = shape
            break

    assert body is not None
    body_rect = (int(body.left), int(body.top), int(body.width), int(body.height))

    # find any picture or shape that has aurora_placeholder name or is a picture
    for shape in slide.shapes:
        if getattr(shape, 'shape_type', None) in (13, 1, 17):
            # PICTURE=13, AUTO_SHAPE=1, TEXT_BOX=17
            rect = (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
            # skip the body itself
            if rect == body_rect:
                continue
            assert not rects_overlap(body_rect, rect), f"Found overlap between body and shape at {rect}"
