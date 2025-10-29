import io
import os
import sys
from pptx import Presentation

# Ensure the generator package directory is importable when tests run from repo root
ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)

from strands_ppt_generator import generate_pptx_file


def test_placeholder_inserted_for_missing_image():
    """Integration test: when an image slide has no image, generator should insert a placeholder PNG
    and tag it with 'aurora_placeholder'."""
    # Minimal presentation structure with a single image slide that has no image_url
    presentation_structure = {
        'presentation_title': 'Test Presentation',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'image',
                'title': 'Test Visual Slide',
                    'caption': 'Test Diagram',
                    # Provide a non-existent S3 URL to force the download path and placeholder logic
                    'image_url': 's3://nonexistent-bucket/does-not-exist.png',
            }
        ],
        'style': 'professional'
    }

    # book_data can be an empty dict for this test
    pptx_bytes = generate_pptx_file(presentation_structure, {})
    assert pptx_bytes is not None and len(pptx_bytes) > 0

    prs = Presentation(io.BytesIO(pptx_bytes))
    found = False
    for slide in prs.slides:
        for shape in slide.shapes:
            name = getattr(shape, 'name', None)
            text = None
            try:
                text = shape.text if shape.has_text_frame else None
            except Exception:
                text = None
            if name == 'aurora_placeholder':
                found = True
                break
            if text and '[PLACEHOLDER]' in text:
                found = True
                break
        if found:
            break

    assert found, "Expected to find a placeholder (aurora_placeholder) in generated PPTX"
