"""
HTML to PowerPoint Converter
=============================
Parses HTML infographic and converts to PowerPoint using browser-rendered layout as source of truth.
This ensures HTML and PPT are EXACTLY synchronized.
"""

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import logging
import io
import re
from typing import Dict, Optional

# Use same logger as main module so logs appear in CloudWatch
logger = logging.getLogger("aurora.infographic_generator")


def convert_html_to_pptx_new(
    html_content: str, 
    structure: Dict,
    course_bucket: str = None, 
    project_folder: str = None, 
    template_key: str = None,
    s3_client = None
) -> bytes:
    """
    Convert HTML infographic to PowerPoint by parsing the HTML structure.
    HTML is the single source of truth - PPT mirrors HTML exactly.
    
    Architecture:
    1. Parse HTML with BeautifulSoup
    2. Extract each slide's content (title, bullets, images)
    3. Create PPT slide with exact same content
    4. No layout calculations - just extract and render
    
    Args:
        html_content: Generated HTML with slides
        structure: Original JSON structure (for metadata only)
        course_bucket: S3 bucket for images
        project_folder: Project folder path
        template_key: Optional template S3 key
        s3_client: boto3 S3 client
    
    Returns:
        PPT file as bytes
    """
    logger.info("="*80)
    logger.info("🎨 STARTING HTML-TO-PPT CONVERSION")
    logger.info("="*80)
    logger.info(f"HTML length: {len(html_content)} chars")
    logger.info(f"Template: {template_key if template_key else 'None (clean design)'}")
    logger.info(f"Course bucket: {course_bucket}")
    logger.info(f"Project folder: {project_folder}")
    
    # Parse HTML
    soup = BeautifulSoup(html_content, 'lxml')
    slides_html = soup.find_all('div', class_='slide')
    
    logger.info(f"📄 Found {len(slides_html)} slides in HTML")
    
    # Load or create presentation
    prs = _load_presentation(template_key, s3_client)
    
    # Get layouts
    layouts = _get_slide_layouts(prs)
    
    # Color scheme
    colors = _get_color_scheme(structure.get('style', 'professional'))
    
    logger.info(f"🚀 ABOUT TO START MAIN LOOP - Processing {len(slides_html)} slides")
    logger.info(f"   Layouts available: {list(layouts.keys())}")
    logger.info(f"   Colors: {colors}")
    
    # Track overflow slides for reporting
    overflow_slides = []
    
    # Process each HTML slide
    for slide_idx, slide_html in enumerate(slides_html):
        logger.info(f"\n{'='*80}")
        logger.info(f"🔄 Processing HTML slide {slide_idx + 1}/{len(slides_html)}")
        
        # Extract slide data from HTML
        slide_data = _extract_slide_data(slide_html)
        
        logger.info(f"   Title: {slide_data['title'][:60]}")
        logger.info(f"   Subtitle: {slide_data.get('subtitle', 'None')}")
        logger.info(f"   Blocks: {len(slide_data['content_blocks'])}")
        logger.info(f"   Images: {len(slide_data['images'])}")
        logger.info(f"   Layout type: {slide_data.get('layout_type', 'normal')}")
        
        # SPECIAL HANDLING: For branded title slides (course/module/lesson)
        if slide_data.get('layout_type') in ['course-title', 'module-title', 'lesson-title']:
            _create_branded_title_slide(prs, layouts['blank'], slide_data, colors, course_bucket, project_folder, s3_client)
            continue  # Skip normal slide creation
        
        # Check if slide has actual text content (not just empty blocks)
        has_text_content = any(
            block.get('type') in ['heading', 'bullets', 'callout'] and 
            (block.get('text') or block.get('items'))
            for block in slide_data['content_blocks']
        )
        
        # Check if slide has callouts - if so, we'll remove logo to prevent overlap
        has_callout_on_slide = any(block.get('type') == 'callout' for block in slide_data['content_blocks'])
        if has_callout_on_slide:
            logger.info(f"🚫 Slide will skip logo - callout detected")
            print(f"DEBUG: Found callout on slide - {slide_data.get('title', 'Unknown')}")
        
        has_images = bool(slide_data['images'])
        
        # Create PPT slide using content placeholder
        # NOTE: NO SPLITTING HERE - HTML already decided the slide structure
        # If HTML has content in one slide, PPT must keep it in one slide
        slide = prs.slides.add_slide(layouts['content'])
        
        # CRITICAL FIX: Clear all placeholder text in the slide
        # Placeholders may have default text like "Untitled" that we need to remove
        placeholder_count = 0
        for shape in slide.placeholders:
            try:
                if hasattr(shape, 'text_frame'):
                    # Clear the placeholder text
                    shape.text_frame.clear()
                    shape.text = ""
                    placeholder_count += 1
                    logger.info(f"  🧹 Cleared placeholder {shape.placeholder_format.idx}: '{shape.name}'")
            except Exception as e:
                logger.debug(f"  ⚠️ Could not clear placeholder: {e}")
        
        # Also remove any visible shapes with placeholder text
        shapes_to_remove = []
        for shape in slide.shapes:
            try:
                # Skip placeholders (we already cleared them above)
                if shape.is_placeholder:
                    continue
                
                # Check if it has text that says "Untitled" or common placeholder text
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip().lower()
                    if text in ['untitled', 'title', 'click to add title', 'click to add subtitle', 'add title', 'add subtitle']:
                        shapes_to_remove.append(shape)
                        logger.info(f"  🗑️ Marking shape with placeholder text for removal: '{text}'")
                        
            except Exception as e:
                logger.debug(f"  ⚠️ Could not check shape: {e}")
        
        # Remove marked shapes
        for shape in shapes_to_remove:
            try:
                sp = shape.element
                sp.getparent().remove(sp)
                logger.info(f"  ✅ Removed shape successfully")
            except Exception as e:
                logger.warning(f"  ⚠️ Could not remove shape: {e}")
        
        logger.info(f"  🧹 Cleaned {placeholder_count} placeholders + {len(shapes_to_remove)} shapes from slide")
        
        # Add logo to regular slides (but NOT lateral image - that's only for title slides)
        # Skip logo on slides with callouts to prevent overlap
        if not has_callout_on_slide:
            _add_logo_to_regular_slide(slide, course_bucket, project_folder, s3_client)
        else:
            logger.info(f"⏭️  Skipping logo - callout present on slide")
        
        # Set title and capture its height for proper subtitle positioning
        title_height = _set_slide_title(slide, slide_data['title'], colors)
        
        # Set subtitle (if exists) and get its total height for content positioning
        subtitle_height = 0  # Default if no subtitle
        if slide_data.get('subtitle'):
            subtitle_height = _set_slide_subtitle(slide, slide_data['subtitle'], colors, title_height=title_height)
        
        # Add ALL content blocks (no splitting - HTML already decided structure)
        # Pass subtitle_height so content block can position itself correctly
        _add_content_blocks(slide, slide_data['content_blocks'], colors, has_images=has_images and has_text_content, subtitle_height=subtitle_height)
        
        # Add images if present
        if slide_data['images'] and course_bucket and project_folder:
            _add_images(slide, slide_data['images'], course_bucket, project_folder, structure, s3_client, colors)
    
    # Save to bytes
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    logger.info(f"\n✅ Created PowerPoint with {len(slides_html)} slides from HTML")
    return pptx_buffer.getvalue()


def _load_presentation(template_key: str, s3_client) -> Presentation:
    """Load template or create blank presentation."""
    if template_key and s3_client:
        try:
            logger.info(f"📥 Loading template from S3: {template_key}")
            template_response = s3_client.get_object(
                Bucket='crewai-course-artifacts', 
                Key=template_key
            )
            template_stream = io.BytesIO(template_response['Body'].read())
            prs = Presentation(template_stream)
            logger.info(f"✅ Template loaded! {len(prs.slide_layouts)} layouts available")
            return prs
        except Exception as e:
            logger.warning(f"⚠️ Failed to load template: {e}")
    
    # Default presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    return prs


def _get_slide_layouts(prs: Presentation) -> Dict:
    """Get slide layouts from template or defaults."""
    # Log all available layouts
    logger.info(f"📋 Available layouts in presentation:")
    for idx, layout in enumerate(prs.slide_layouts):
        placeholder_count = len(layout.placeholders)
        logger.info(f"   Layout {idx}: '{layout.name}' ({placeholder_count} placeholders)")
    
    # Try to get truly blank layout (layout 6 is typically blank in default templates)
    blank_layout = None
    
    # First priority: Search for layout named "Blank"
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            blank_layout = layout
            logger.info(f"✅ Found blank layout by name: '{layout.name}'")
            break
    
    # Second priority: Search for layout with 0 placeholders
    if not blank_layout:
        for layout in prs.slide_layouts:
            if len(layout.placeholders) == 0:
                blank_layout = layout
                logger.info(f"✅ Found blank layout with 0 placeholders: '{layout.name}'")
                break
    
    # Fallback: Use layout with fewest placeholders
    if not blank_layout:
        blank_layout = min(prs.slide_layouts, key=lambda l: len(l.placeholders))
        logger.info(f"⚠️ Using layout '{blank_layout.name}' with {len(blank_layout.placeholders)} placeholders (fewest available)")
    
    return {
        'title': blank_layout,
        'content': blank_layout,  # Use blank for all slides
        'blank': blank_layout
    }


def _get_color_scheme(style: str) -> Dict:
    """Get color scheme based on style - UPDATED to match HTML cyan theme."""
    color_map = {
        'professional': {
            'primary': RGBColor(0, 60, 120),    # Dark blue #003c78 for headings
            'secondary': RGBColor(0, 153, 204),  # Cyan #0099cc for bullets (matching HTML)
            'accent': RGBColor(255, 204, 0),     # Yellow #ffcc00 for callouts
            'text': RGBColor(51, 51, 51)         # Dark gray #333333 for body text
        },
        'modern': {
            'primary': RGBColor(33, 33, 33),
            'secondary': RGBColor(0, 153, 204),  # Cyan for bullets
            'accent': RGBColor(0, 150, 136),
            'text': RGBColor(51, 51, 51)
        },
        'minimal': {
            'primary': RGBColor(0, 0, 0),
            'secondary': RGBColor(0, 153, 204),  # Cyan for bullets
            'accent': RGBColor(117, 117, 117),
            'text': RGBColor(51, 51, 51)
        }
    }
    return color_map.get(style, color_map['professional'])


def _extract_slide_data(slide_html) -> Dict:
    """
    Extract structured data from HTML slide.
    
    Returns:
        {
            'title': str,
            'subtitle': str or None,
            'content_blocks': [{'type': 'bullets', 'items': [...]}, ...],
            'images': [{'alt': str, 'caption': str}, ...],
            'layout_type': str  # 'normal', 'course-title', 'module-title', 'lesson-title'
        }
    """
    # Detect special title slides by CSS class
    layout_type = 'normal'
    slide_classes = slide_html.get('class', [])
    if 'course-title' in slide_classes:
        layout_type = 'course-title'
    elif 'module-title' in slide_classes:
        layout_type = 'module-title'
    elif 'lesson-title' in slide_classes:
        layout_type = 'lesson-title'
    
    # Extract title (different selector for special title slides)
    if layout_type in ['course-title', 'module-title', 'lesson-title']:
        title_elem = slide_html.find('h1', class_='main-title')
    else:
        title_elem = slide_html.find('h1', class_='slide-title')
    title = title_elem.get_text(strip=True) if title_elem else "Untitled"
    logger.info(f"  📝 Extracted title: {title[:50]}")
    
    # Extract subtitle (different selector for special title slides)
    if layout_type in ['course-title', 'module-title', 'lesson-title']:
        subtitle_elem = slide_html.find('p', class_='main-subtitle')
    else:
        subtitle_elem = slide_html.find('p', class_='slide-subtitle')
    subtitle = subtitle_elem.get_text(strip=True) if subtitle_elem else None
    if subtitle:
        logger.info(f"  📝 Extracted subtitle: {subtitle[:50]}")
    
    # Extract content blocks IN DOM ORDER
    # CRITICAL: Must preserve the order from HTML (heading → bullets → heading → bullets)
    content_blocks = []
    
    logger.info(f"  📦 Extracting content blocks in DOM order from .content-block divs")
    
    # Find all content-block divs and process them in order
    content_block_divs = slide_html.find_all('div', class_='content-block')
    logger.info(f"  🔍 Found {len(content_block_divs)} <div class='content-block'> containers")
    
    for block_idx, block_div in enumerate(content_block_divs):
        logger.info(f"    📦 Content block {block_idx + 1}:")
        
        # Check for heading within this block
        h2 = block_div.find('h2', class_='block-heading')
        if h2:
            heading_text = h2.get_text(strip=True)
            logger.info(f"      ✓ Heading: {heading_text[:50]}")
            content_blocks.append({
                'type': 'heading',
                'text': heading_text
            })
        
        # Check for bullets within this block
        ul = block_div.find('ul', class_='bullets')
        if ul:
            items = []
            for li in ul.find_all('li'):
                item_text = li.get_text(strip=True)
                # Preserve level information from HTML class
                if 'level-2' in li.get('class', []):
                    items.append({'text': item_text, 'level': 2})
                else:
                    items.append({'text': item_text, 'level': 1})
            
            if items:
                logger.info(f"      ✓ Bullets: {len(items)} items")
                for item in items[:2]:  # Log first 2
                    level_str = '  ' if item.get('level') == 2 else ''
                    logger.info(f"        {level_str}- {item['text'][:60]}...")
                content_blocks.append({
                    'type': 'bullets',
                    'items': items
                })
        
        # Check for callout within this block
        callout = block_div.find('div', class_='callout')
        if callout:
            callout_text = callout.get_text(strip=True)
            if callout_text:  # Only add if not empty
                logger.info(f"      ✓ Callout: {callout_text[:50]}")
                content_blocks.append({
                    'type': 'callout',
                    'text': callout_text
                })
            else:
                logger.info(f"      ⊘ Skipping empty callout")
    
    logger.info(f"  ✅ Total content blocks extracted (in order): {len(content_blocks)}")
    
    # Extract images
    images = []
    for img_container in slide_html.find_all('div', class_='image-container'):
        img_placeholder = img_container.find('div', class_='image-placeholder')
        caption_elem = img_container.find('p', class_='image-caption')
        
        if img_placeholder:
            # Extract image reference from placeholder text
            image_ref = img_placeholder.get_text(strip=True)
            images.append({
                'reference': image_ref,
                'caption': caption_elem.get_text(strip=True) if caption_elem else ''
            })
    
    return {
        'title': title,
        'subtitle': subtitle,
        'content_blocks': content_blocks,
        'images': images,
        'layout_type': layout_type
    }


def _set_slide_title(slide, title: str, colors: Dict):
    """Set slide title using text box with proper fixed height.
    
    Returns:
        float: Height of the title box in inches
    """
    try:
        # Calculate height based on actual PowerPoint rendering:
        # At 32pt bold, ~45 chars fit per line at 11.7" width
        # Line height: 32pt * 1.2 = 38.4pt ≈ 0.533"
        chars_per_line = 45
        num_lines = max(1, (len(title) + chars_per_line - 1) // chars_per_line)
        # Height: 0.6" for first line + 0.5" per additional line
        title_height = 0.6 if num_lines == 1 else 0.6 + (num_lines - 1) * 0.5
        
        # Create text box with calculated fixed height
        title_box = slide.shapes.add_textbox(
            Inches(0.8),              # left margin
            Inches(0.5),              # top margin
            Inches(11.7),             # width
            Inches(title_height)      # calculated height
        )
        
        # Add white background to cover any "Untitled" text from slide master
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        title_box.line.fill.background()
        
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        # Set title text
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = colors['primary']
        
        logger.info(f"✅ Set title: {title[:50]} ({num_lines} lines, height: {title_height:.2f}\")")
        return title_height
        
    except Exception as e:
        logger.error(f"❌ Failed to set title: {e}")
        return 0.6  # Return default height on error


def _set_slide_subtitle(slide, subtitle: str, colors: Dict, title_height: float = 0.6):
    """Set slide subtitle with dynamic positioning and proper fixed height.
    
    Args:
        title_height: Height of title box in inches
    
    Returns:
        float: Total height from slide top (0.5") to bottom of subtitle
    """
    try:
        # Calculate subtitle height based on actual PowerPoint rendering:
        # At 20pt regular, ~75 chars fit per line at 11.7" width
        # Line height: 20pt * 1.2 = 24pt ≈ 0.333"
        chars_per_line = 75
        num_lines = max(1, (len(subtitle) + chars_per_line - 1) // chars_per_line)
        # Height: 0.4" for first line + 0.35" per additional line
        subtitle_height = 0.4 if num_lines == 1 else 0.4 + (num_lines - 1) * 0.35
        
        # Position subtitle below title with 0.25" gap
        subtitle_top = 0.5 + title_height + 0.25
        
        # Create subtitle text box with calculated fixed height
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.8),              # left margin (same as title)
            Inches(subtitle_top),     # position below title
            Inches(11.7),             # width (same as title)
            Inches(subtitle_height)   # calculated height
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        # Set text with formatting
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = RGBColor(70, 130, 180)  # Steel blue
        
        logger.info(f"✅ Set subtitle: {subtitle[:50]} ({num_lines} lines, top: {subtitle_top:.2f}\", height: {subtitle_height:.2f}\")")
        
        # Return total height from slide top to bottom of subtitle
        return subtitle_top + subtitle_height
        
    except Exception as e:
        logger.error(f"❌ Failed to set subtitle: {e}")
        return 1.5  # Return default on error


def _create_branded_title_slide(prs, blank_layout, slide_data: Dict, colors: Dict, course_bucket: str, project_folder: str, s3_client):
    """Create a full-screen branded title slide with background images matching Netec design."""
    try:
        from PIL import Image as PILImage
        
        layout_type = slide_data.get('layout_type', 'normal')
        title = slide_data.get('title', '')
        subtitle = slide_data.get('subtitle', '')
        
        logger.info(f"🎨 Creating branded {layout_type} slide: {title[:40]}")
        
        # Create slide with blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Clear any placeholders
        for shape in slide.placeholders:
            try:
                if hasattr(shape, 'text_frame'):
                    shape.text_frame.clear()
                    shape.text = ""
            except:
                pass
        
        # COURSE TITLE SLIDE - Full background image (Netec_Portada_1.png)
        if layout_type == 'course-title':
            try:
                # Download and add full background image
                bg_key = 'PPT_Templates/Netec_Portada_1.png'
                logger.info(f"📥 Downloading background: s3://{course_bucket}/{bg_key}")
                bg_response = s3_client.get_object(Bucket=course_bucket, Key=bg_key)
                bg_bytes = bg_response['Body'].read()
                bg_stream = io.BytesIO(bg_bytes)
                
                # Add background image (full slide)
                slide.shapes.add_picture(
                    bg_stream,
                    Inches(0),      # left
                    Inches(0),      # top
                    width=Inches(13.333),   # full width
                    height=Inches(7.5)      # full height
                )
                logger.info(f"✅ Background image added")
                
                # Add title text (centered, white, large)
                title_box = slide.shapes.add_textbox(
                    Inches(1.0),    # left margin
                    Inches(2.8),    # centered vertically
                    Inches(11.333), # width
                    Inches(1.8)     # height
                )
                
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                title_para = title_frame.paragraphs[0]
                title_para.text = title
                title_para.font.size = Pt(60)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(255, 255, 255)  # White
                title_para.alignment = PP_ALIGN.CENTER
                
            except Exception as e:
                logger.error(f"❌ Failed to add course background: {e}")
                # Fallback to solid color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 188, 235)
        
        # MODULE TITLE SLIDE - White background with lateral image
        elif layout_type == 'module-title':
            try:
                # Set white background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
                
                # Add logo at top-left
                logo_key = 'logo/LogoNetec.png'
                logo_response = s3_client.get_object(Bucket=course_bucket, Key=logo_key)
                logo_bytes = logo_response['Body'].read()
                logo_stream = io.BytesIO(logo_bytes)
                
                slide.shapes.add_picture(
                    logo_stream,
                    Inches(0.5),   # left
                    Inches(0.5),   # top
                    width=Inches(2.0),
                    height=Inches(0.67)
                )
                
                # Add lateral decorative image on the right (MODULE - aligned to edge)
                lateral_key = 'PPT_Templates/Netec_Lateral_1.png'
                logger.info(f"📥 Downloading lateral image: s3://{course_bucket}/{lateral_key}")
                lateral_response = s3_client.get_object(Bucket=course_bucket, Key=lateral_key)
                lateral_bytes = lateral_response['Body'].read()
                lateral_stream = io.BytesIO(lateral_bytes)
                
                # Calculate width to cover from 8" to right edge (13.333")
                lateral_width = Inches(5.333)  # 13.333 - 8.0 = 5.333"
                
                # Add lateral image aligned to right edge
                slide.shapes.add_picture(
                    lateral_stream,
                    Inches(8.0),          # Start at 8" from left
                    Inches(0),            # Top edge
                    width=lateral_width,  # Exactly 5.333" to reach right edge
                    height=Inches(7.5)    # Full height
                )
                logger.info(f"✅ Module lateral image added (5.333\" wide, aligned to right edge)")
                
                # Add module title (left side, dark blue, large)
                title_box = slide.shapes.add_textbox(
                    Inches(0.5),    # left margin
                    Inches(2.5),    # centered vertically
                    Inches(6.5),    # width (left side only)
                    Inches(2.5)     # height
                )
                
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                title_para = title_frame.paragraphs[0]
                title_para.text = title
                title_para.font.size = Pt(48)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 60, 120)  # Dark blue
                title_para.alignment = PP_ALIGN.LEFT
                
            except Exception as e:
                logger.error(f"❌ Failed to add module images: {e}")
                import traceback
                logger.error(traceback.format_exc())
        
        # LESSON TITLE SLIDE - White background with lateral image + module subtitle
        elif layout_type == 'lesson-title':
            try:
                # Set white background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
                
                # Add logo at top-left
                logo_key = 'logo/LogoNetec.png'
                logo_response = s3_client.get_object(Bucket=course_bucket, Key=logo_key)
                logo_bytes = logo_response['Body'].read()
                logo_stream = io.BytesIO(logo_bytes)
                
                slide.shapes.add_picture(
                    logo_stream,
                    Inches(0.5),   # left
                    Inches(0.5),   # top
                    width=Inches(2.0),
                    height=Inches(0.67)
                )
                
                # Add lateral decorative image on the right (LESSON - aligned to edge)
                lateral_key = 'PPT_Templates/Netec_Lateral_1.png'
                logger.info(f"📥 Downloading lateral image: s3://{course_bucket}/{lateral_key}")
                lateral_response = s3_client.get_object(Bucket=course_bucket, Key=lateral_key)
                lateral_bytes = lateral_response['Body'].read()
                lateral_stream = io.BytesIO(lateral_bytes)
                
                # Calculate width to cover from 8" to right edge (13.333")
                lateral_width = Inches(5.333)  # 13.333 - 8.0 = 5.333"
                
                # Add lateral image aligned to right edge
                slide.shapes.add_picture(
                    lateral_stream,
                    Inches(8.0),          # Start at 8" from left
                    Inches(0),            # Top edge
                    width=lateral_width,  # Exactly 5.333" to reach right edge
                    height=Inches(7.5)    # Full height
                )
                logger.info(f"✅ Lesson lateral image added (5.333\" wide, aligned to right edge)")
                
                # Add lesson title (left side, dark blue, large)
                title_box = slide.shapes.add_textbox(
                    Inches(0.5),    # left margin
                    Inches(2.2),    # higher up for subtitle below
                    Inches(6.5),    # width (left side only)
                    Inches(1.8)     # height
                )
                
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                title_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
                
                title_para = title_frame.paragraphs[0]
                title_para.text = title
                title_para.font.size = Pt(42)
                title_para.font.bold = True
                title_para.font.color.rgb = RGBColor(0, 60, 120)  # Dark blue
                title_para.alignment = PP_ALIGN.LEFT
                
                # Add module subtitle below (cyan color)
                if subtitle:
                    subtitle_box = slide.shapes.add_textbox(
                        Inches(0.5),    # left margin
                        Inches(4.2),    # below title
                        Inches(6.5),    # width
                        Inches(0.8)     # height
                    )
                    
                    subtitle_frame = subtitle_box.text_frame
                    subtitle_frame.word_wrap = True
                    
                    subtitle_para = subtitle_frame.paragraphs[0]
                    subtitle_para.text = subtitle
                    subtitle_para.font.size = Pt(28)
                    subtitle_para.font.bold = False
                    subtitle_para.font.color.rgb = RGBColor(0, 188, 235)  # Cyan
                    subtitle_para.alignment = PP_ALIGN.LEFT
                
            except Exception as e:
                logger.error(f"❌ Failed to add lesson images: {e}")
                import traceback
                logger.error(traceback.format_exc())
        
        logger.info(f"✅ Created branded {layout_type} slide")
        
    except Exception as e:
        logger.error(f"❌ Failed to create branded title slide: {e}")
        import traceback
        logger.error(traceback.format_exc())


def _add_logo_to_regular_slide(slide, course_bucket: str, project_folder: str, s3_client):
    """Add Netec branding to regular content slides: small corner blocks + logo in bottom-right.
    
    Args:
        slide: PPT slide object
        course_bucket: S3 bucket name
        project_folder: Project folder path
        s3_client: boto3 S3 client
    """
    try:
        # 1. Add small dark blue corner blocks (small rectangles, not thin lines)
        # Left corner block
        left_frame = slide.shapes.add_shape(
            1,  # Rectangle shape type
            left=Inches(0),
            top=Inches(0.5),  # Align with title text position
            width=Inches(0.52),  # ~50px at 96 DPI - small block
            height=Inches(0.625)   # ~60px at 96 DPI - short rectangular block
        )
        left_frame.fill.solid()
        left_frame.fill.fore_color.rgb = RGBColor(0, 60, 120)  # Dark blue #003c78
        left_frame.line.fill.background()  # No border
        logger.info(f"✅ Added left corner block")
        
        # Right corner block
        right_frame = slide.shapes.add_shape(
            1,  # Rectangle shape type
            left=Inches(12.81),  # 13.333 - 0.52 = 12.81
            top=Inches(0.5),  # Align with title text position
            width=Inches(0.52),  # Small block
            height=Inches(0.625)  # ~60px at 96 DPI - short rectangular block
        )
        right_frame.fill.solid()
        right_frame.fill.fore_color.rgb = RGBColor(0, 60, 120)  # Dark blue #003c78
        right_frame.line.fill.background()  # No border
        logger.info(f"✅ Added right corner block")
        
        # 2. Add logo in BOTTOM-right corner (won't be hidden by title)
        logo_key = 'logo/LogoNetec.png'
        
        try:
            logo_response = s3_client.get_object(Bucket=course_bucket, Key=logo_key)
            logo_data = logo_response['Body'].read()
            logo_stream = io.BytesIO(logo_data)
            
            # Position: Bottom-right corner with margin
            slide.shapes.add_picture(
                logo_stream,
                left=Inches(11.1),  # Right side (13.333 - 2.0 - 0.2 margin)
                top=Inches(6.6),    # Bottom (7.5 - 0.67 - 0.2 margin)
                width=Inches(2.0),
                height=Inches(0.67)
            )
            logger.info(f"✅ Added logo to regular slide (bottom-right corner)")
            
        except Exception as e:
            logger.warning(f"⚠️ Could not add logo to regular slide: {e}")
            
    except Exception as e:
        logger.error(f"❌ Failed to add branding to regular slide: {e}")
        import traceback
        logger.error(traceback.format_exc())


def _add_content_blocks(slide, content_blocks: list, colors: Dict, has_images: bool = False, subtitle_height: float = 0):
    """Add content blocks to slide with improved overflow detection and positioning.
    
    Args:
        subtitle_height: Total height from top (0.5") to end of subtitle box (for accurate positioning)
    """
    if not content_blocks:
        logger.warning("⚠️ No content blocks to add")
        return
    
    try:
        # If slide has images, make content narrower (left side only)
        if has_images:
            content_width = 6.0   # Narrower to fit beside image (left half of slide)
            logger.info("📐 Using narrow layout (6.0\") for content - image on right")
        else:
            content_width = 11.7  # Full width (13.333 - 0.8 - 0.833 margins)
            logger.info("📐 Using full width layout (11.7\") for content")
        
        # Calculate content position more accurately
        # subtitle_height already includes: title_top + title_height + gap_to_subtitle + subtitle_height
        # So content should start right after subtitle + small gap
        if subtitle_height > 0:
            content_top = subtitle_height + 0.15  # Small gap (0.15") after subtitle
        else:
            # No subtitle: content starts after title area (0.5" top + ~0.7" default title height + 0.1" gap)
            content_top = 1.3
        
        # Ensure content doesn't start too late (max 2.8" from top)
        content_top = min(content_top, 2.8)
        
        # Calculate remaining space for content with safety margin at bottom (0.3")
        large_height = max(4.0, 7.5 - content_top - 0.3)  # Remaining space on slide
        
        # Log available space
        logger.info(f"📐 Content box: top={content_top:.2f}\", height={large_height:.2f}\", subtitle_height_used={subtitle_height:.2f}\"")
        
        content_box = slide.shapes.add_textbox(
            Inches(0.8),              # left margin (match HTML)
            Inches(content_top),      # top (dynamically positioned after subtitle)
            Inches(content_width),    # width (narrow if images, full if not)
            Inches(large_height)      # large height - text will use only what it needs
        )
        
        text_frame = content_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        
        logger.info(f"📝 Processing {len(content_blocks)} content blocks...")
        
        # Separate callouts from other content - we'll add them at the end
        regular_blocks = []
        callout_blocks = []
        for block in content_blocks:
            if block.get('type') == 'callout':
                callout_blocks.append(block)
            else:
                regular_blocks.append(block)
        
        # Check if there are any bullets in regular blocks (for callout centering logic)
        has_bullets = any(block.get('type') == 'bullets' for block in regular_blocks)
        logger.info(f"  Regular blocks: {len(regular_blocks)}, Callouts: {len(callout_blocks)}, Has bullets: {has_bullets}")
        
        # Add regular content blocks first (these flow naturally in text frame)
        first_paragraph = True
        total_bullets = 0  # Track total bullets for position estimation
        
        for block_idx, block in enumerate(regular_blocks):
            logger.info(f"  Block {block_idx + 1}: type={block.get('type')}, items={len(block.get('items', []))}")
            
            if block.get('type') == 'heading':
                # Add heading
                if first_paragraph:
                    p = text_frame.paragraphs[0]
                    first_paragraph = False
                else:
                    p = text_frame.add_paragraph()
                    
                p.text = block.get('text', '')
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = colors['primary']
                p.space_before = Pt(20)  # Add spacing before heading (matches HTML margin-top: 25px)
                p.space_after = Pt(8)
                logger.info(f"    ✓ Added heading: {block.get('text', '')[:40]}")
                
            elif block.get('type') == 'bullets':
                # Add bullet list with symbols matching HTML
                items = block.get('items', [])
                heading_text = block.get('heading', '')
                
                # Add heading if present (e.g., "¿Para quién es este curso?")
                if heading_text:
                    if first_paragraph:
                        p = text_frame.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = heading_text
                    p.font.size = Pt(24)
                    p.font.bold = True
                    p.font.color.rgb = colors['primary']
                    p.space_before = Pt(20)  # Add spacing before heading
                    p.space_after = Pt(8)
                    logger.info(f"    ✓ Added bullet section heading: {heading_text[:40]}")
                
                logger.info(f"    Adding {len(items)} bullets...")
                total_bullets += len(items)  # Track for callout positioning
                
                for item_idx, item in enumerate(items):
                    if first_paragraph:
                        p = text_frame.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = text_frame.add_paragraph()
                    
                    # Get item text and level from HTML structure
                    if isinstance(item, dict):
                        item_text = item.get('text', '')
                        item_level = item.get('level', 1)
                    else:
                        # Fallback for old format (plain strings)
                        item_text = str(item)
                        item_level = 1
                    
                    if item_level == 2:
                        # Second level: Cyan square bullet, indented
                        bullet_run = p.add_run()
                        bullet_run.text = "▪ "
                        bullet_run.font.size = Pt(20)
                        bullet_run.font.bold = True
                        bullet_run.font.color.rgb = colors.get('secondary', RGBColor(0, 188, 235))  # CYAN
                        
                        # Text in BLACK (matching HTML)
                        text_run = p.add_run()
                        text_run.text = item_text
                        text_run.font.size = Pt(18)
                        text_run.font.color.rgb = RGBColor(0, 0, 0)  # BLACK
                        
                        p.level = 1  # Indented level
                        p.line_spacing = 1.3
                        p.space_after = Pt(4)
                    else:
                        # First level: Yellow triangle bullet
                        bullet_run = p.add_run()
                        bullet_run.text = "▸ "
                        bullet_run.font.size = Pt(24)
                        bullet_run.font.bold = True
                        bullet_run.font.color.rgb = colors.get('accent', RGBColor(255, 204, 0))  # YELLOW #FFC000
                        
                        # Second run: BLACK text for item (matching HTML)
                        text_run = p.add_run()
                        text_run.text = item_text
                        text_run.font.size = Pt(20)
                        text_run.font.color.rgb = RGBColor(0, 0, 0)  # BLACK
                        
                        p.level = 0
                        p.line_spacing = 1.4
                        p.space_after = Pt(6)
                    
                    logger.info(f"      • {'  ' if item_level == 2 else ''}{item_text[:50]}...")
        
        # NOW add callouts AFTER all regular content (they won't overlap)
        if callout_blocks:
            logger.info(f"  Adding {len(callout_blocks)} callout(s) at end of content...")
            
            # If NO bullets on slide, vertically center callouts
            if not has_bullets and len(regular_blocks) == 0:
                # Pure callout slide - center vertically
                logger.info(f"  💡 No bullets detected - will center callouts vertically")
                callout_top = Inches(2.5)  # Roughly center vertically on slide
            else:
                # Calculate starting position for callouts based on content added
                # CRITICAL: Callouts must be BELOW all bullet content to avoid hiding text
                # Conservative estimate: title (1.7") + margin (1.5") + bullets (0.5" each) + safety buffer (1.0")
                callout_top = Inches(1.7 + 1.5 + (total_bullets * 0.5) + 1.0)  
                
                # IMPORTANT: Callouts should start at least 5.5 inches from top
                # This ensures they're in the bottom third of the slide (slide height ~7.5")
                callout_top = max(callout_top, Inches(5.5))
            
            for callout_idx, block in enumerate(callout_blocks):
                callout_text = block.get('text', '')
                
                callout_left = Inches(0.8)
                callout_width = Inches(11.7) if not has_images else Inches(6.0)
                
                # Calculate height based on text length - MORE ACCURATE
                text_length = len(callout_text)
                # At 18pt bold with emoji, approximately 70 chars per line
                chars_per_line = 70 if not has_images else 55
                estimated_lines = max(2, (text_length // chars_per_line) + 1)
                # More generous height: 0.25" per line + 0.4" padding (top/bottom margins)
                callout_height = Inches(0.25 * estimated_lines + 0.4)
                
                # For centered callouts (no bullets), position at center
                if not has_bullets and len(regular_blocks) == 0 and callout_idx == 0:
                    # Adjust to center considering callout height
                    # Slide height 7.5", try to position center of callout at middle of slide
                    callout_top = Inches(max(2.5, (7.5 - callout_height.inches) / 2))
                    logger.info(f"  Centering callout vertically at {callout_top}")
                
                # Add yellow background box
                callout_box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    callout_left, callout_top,
                    callout_width, callout_height
                )
                
                # Style the box: light yellow background, no border
                callout_box.fill.solid()
                callout_box.fill.fore_color.rgb = RGBColor(255, 250, 230)  # Light yellow #FFFAE6
                callout_box.line.fill.background()  # No border
                
                # Add yellow left border (thick line)
                left_border = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    callout_left, callout_top,
                    Inches(0.05), callout_height  # Thin vertical rectangle
                )
                left_border.fill.solid()
                left_border.fill.fore_color.rgb = RGBColor(255, 204, 0)  # Yellow #FFC000
                left_border.line.fill.background()  # No border
                
                # Add text with lightbulb emoji
                callout_text_frame = callout_box.text_frame
                callout_text_frame.margin_left = Inches(0.15)
                callout_text_frame.margin_right = Inches(0.1)
                callout_text_frame.margin_top = Inches(0.15)
                callout_text_frame.margin_bottom = Inches(0.1)
                callout_text_frame.word_wrap = True
                
                callout_p = callout_text_frame.paragraphs[0]
                callout_p.alignment = PP_ALIGN.LEFT
                
                callout_run = callout_p.add_run()
                callout_run.text = f"💡 {callout_text}"
                callout_run.font.size = Pt(18)
                callout_run.font.bold = True
                callout_run.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray text
                
                logger.info(f"    ✓ Added callout {callout_idx + 1} at top={callout_top}")
                
                # Update position for next callout
                callout_top += callout_height + Inches(0.2)
        
        # OVERFLOW DETECTION: Check if content exceeds slide boundaries
        # Get actual text box dimensions after all content is added
        actual_bottom = (content_box.top + content_box.height) / 914400  # Convert EMU to inches
        slide_height = 7.5  # Standard PPT slide height
        safety_margin = 0.3  # Bottom margin for footer/logo
        max_content_bottom = slide_height - safety_margin  # 7.2" is max safe position
        overflow_amount = actual_bottom - max_content_bottom
        
        if overflow_amount > 0.15:  # More than 0.15" overflow (after safety margin)
            logger.warning(f"⚠️  OVERFLOW DETECTED: Content extends {overflow_amount:.2f}\" beyond safe zone")
            logger.warning(f"   Text box ends at: {actual_bottom:.2f}\" (max safe: {max_content_bottom:.2f}\", slide height: {slide_height}\")")  
            logger.warning(f"   Content box top: {content_top:.2f}\", height: {large_height:.2f}\"")
            logger.warning(f"   Content blocks: {len(regular_blocks)} regular + {len(callout_blocks)} callouts")
            # Log block details for debugging
            total_chars = 0
            for idx, block in enumerate(content_blocks):
                if block.get('type') == 'heading':
                    chars = len(block.get('text', ''))
                    logger.warning(f"   Block {idx}: HEADING - {chars} chars")
                    total_chars += chars
                elif block.get('type') == 'bullets':
                    chars = sum(len(item.get('text', '') if isinstance(item, dict) else str(item)) for item in block.get('items', []))
                    logger.warning(f"   Block {idx}: BULLETS - {len(block.get('items', []))} items, {chars} chars total")
                    total_chars += chars
                elif block.get('type') == 'callout':
                    chars = len(block.get('text', ''))
                    logger.warning(f"   Block {idx}: CALLOUT - {chars} chars")
                    total_chars += chars
            logger.warning(f"   Total content: {total_chars} characters")
            logger.warning(f"💡 RECOMMENDATION: Split content across multiple slides or reduce bullet count")
        else:
            margin_available = max_content_bottom - actual_bottom
            logger.info(f"✅ Content fits safely: {actual_bottom:.2f}\" / {max_content_bottom:.2f}\" (safety margin: {margin_available:.2f}\")")  
        
        logger.info(f"✅ Successfully added {len(regular_blocks)} regular blocks + {len(callout_blocks)} callouts to slide")
        
    except Exception as e:
        logger.error(f"❌ Failed to add content blocks: {e}")
        import traceback
        logger.error(traceback.format_exc())


def _add_images(slide, images: list, course_bucket: str, project_folder: str, structure: Dict, s3_client, colors: Dict):
    """Download and add images to slide.
    
    Images are ALWAYS maximized to nearly full slide width (12" x 6") centered.
    No captions or side-by-side layouts - priority is maximum image size.
    """
    try:
        from PIL import Image
        
        # Add first image only
        if images:
            image_data = images[0]
            image_ref = image_data['reference']
            
            # Download image from S3
            image_bytes = _download_image_from_s3(
                image_ref, 
                course_bucket, 
                project_folder, 
                structure.get('image_url_mapping', {}),
                s3_client
            )
            
            if image_bytes:
                # Get image dimensions
                img = Image.open(io.BytesIO(image_bytes))
                width_px, height_px = img.size
                width_in = width_px / 96.0
                height_in = height_px / 96.0
                
                # ABSOLUTE MAXIMUM IMAGE SIZE - Fill entire content area
                # Use maximum possible dimensions while leaving minimal margins
                max_width = 12.0  # Maximum width (13.333" total, leave 0.67" margins)
                max_height = 6.0  # Maximum height (fill content area)
                img_top = 1.5     # Minimal space below title
                logger.info("📐 ABSOLUTE MAXIMUM image: 12\" x 6\" centered")
                
                # Scale to fit
                scale = min(max_width / width_in, max_height / height_in, 1.0)  # Don't upscale
                final_width = width_in * scale
                final_height = height_in * scale
                
                # Center image horizontally on slide
                img_left = (13.333 - final_width) / 2
                
                # Add to slide
                pic_stream = io.BytesIO(image_bytes)
                picture = slide.shapes.add_picture(
                    pic_stream, 
                    Inches(img_left), 
                    Inches(img_top),
                    width=Inches(final_width), 
                    height=Inches(final_height)
                )
                
                logger.info(f"🖼️ Added MAXIMIZED image {final_width:.2f}x{final_height:.2f} at ({img_left:.2f}, {img_top})")
                
                # NO CAPTION - removed to maximize space for larger images
                
    except Exception as e:
        logger.warning(f"⚠️ Failed to add images: {e}")


def _download_image_from_s3(image_ref: str, course_bucket: str, project_folder: str, 
                            image_url_mapping: Dict, s3_client) -> Optional[bytes]:
    """Download image from S3 based on reference."""
    try:
        # Clean reference
        clean_ref = image_ref.replace('USE_IMAGE: ', '').strip()
        
        # Check if it's a full S3 URL
        if clean_ref.startswith('http'):
            url_match = re.search(r'https://([^/]+)\.s3\.amazonaws\.com/(.+)', clean_ref)
            if url_match:
                bucket = url_match.group(1)
                key = url_match.group(2)
                logger.info(f"📥 Downloading from URL: s3://{bucket}/{key}")
                response = s3_client.get_object(Bucket=bucket, Key=key)
                return response['Body'].read()
        
        # Check URL mapping
        if clean_ref in image_url_mapping:
            mapped_url = image_url_mapping[clean_ref]
            logger.info(f"🔍 Found URL mapping: {clean_ref} → {mapped_url[:80]}")
            return _download_image_from_s3(mapped_url, course_bucket, project_folder, 
                                          image_url_mapping, s3_client)
        
        # Extract numeric ID and try legacy path
        match = re.search(r'(\d{2}-\d{2}-\d{4})', clean_ref)
        if match:
            image_id = match.group(1)
            
            if image_id in image_url_mapping:
                mapped_url = image_url_mapping[image_id]
                logger.info(f"🔍 Found ID mapping: {image_id} → {mapped_url[:80]}")
                return _download_image_from_s3(mapped_url, course_bucket, project_folder,
                                              image_url_mapping, s3_client)
            
            # Legacy path
            image_key = f"{project_folder}/images/{image_id}.png"
            logger.info(f"📥 Downloading legacy: s3://{course_bucket}/{image_key}")
            response = s3_client.get_object(Bucket=course_bucket, Key=image_key)
            return response['Body'].read()
        
        logger.warning(f"⚠️ Could not resolve image reference: {image_ref}")
        return None
        
    except Exception as e:
        logger.warning(f"⚠️ Failed to download image {image_ref}: {e}")
        return None
