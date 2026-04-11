"""
HTML to PPT Converter with Style Mapping
Converts HTML infographic slides to PowerPoint while preserving visual design.
Includes logo and image support.
"""

import io
import logging
import re
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger()
logger.setLevel(logging.INFO)

# =============================================================================
# DESIGN CONSTANTS (matching HTML CSS)
# =============================================================================

COLORS = {
    'header_dark': RGBColor(0x00, 0x33, 0x66),
    'header_light': RGBColor(0x46, 0x82, 0xB4),
    'bullet_marker': RGBColor(0xFF, 0xC0, 0x00),
    'code_bg': RGBColor(0x1E, 0x1E, 0x1E),
    'code_text': RGBColor(0xD4, 0xD4, 0xD4),
    'text_dark': RGBColor(0x00, 0x33, 0x66),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'callout_bg': RGBColor(0xFF, 0xC0, 0x00),
    'primary': RGBColor(0x00, 0x33, 0x66),
    'secondary': RGBColor(0x46, 0x82, 0xB4),
    'text_black': RGBColor(0x33, 0x33, 0x33),
}

FONTS = {
    'title': 'Neue Haas Grotesk Text Pro',
    'body': 'Neue Haas Grotesk Text Pro', 
    'code': 'Consolas',
}

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layout measurements
HEADER_HEIGHT = Inches(1.2)
CONTENT_TOP = Inches(1.4)
CONTENT_LEFT = Inches(0.5)
CONTENT_WIDTH = Inches(12.333)
CONTENT_HEIGHT = Inches(5.5)

# Logo settings
LOGO_WIDTH = Inches(1.5)
LOGO_HEIGHT = Inches(0.6)
LOGO_RIGHT_MARGIN = Inches(0.3)
LOGO_BOTTOM_MARGIN = Inches(0.3)

# =============================================================================
# SUPPLEMENTARY DATA HELPERS (for backward-compatible enrichment)
# =============================================================================

# Noise items that are AI section headings, not real objectives/topics
_NOISE_HEADINGS = {
    'visión general del concepto', 'detalles técnicos', 'aplicación práctica',
    'puntos clave', 'próximos pasos', 'recursos adicionales',
    'información general', 'temas principales del capítulo',
    'lecciones incluidas',
}


def _extract_objectives_from_md(content: str) -> list:
    """Extract learning objectives bullet items from lesson markdown."""
    if not content:
        return []
    pattern = re.compile(
        r'^#{2,3}\s*(?:Objetivos[^\n]*|Learning\s+Objectives[^\n]*)\n+'
        r'(?:[^\n]*\n)*?'
        r'((?:\s*-\s+.+\n?)+)',
        re.MULTILINE | re.IGNORECASE,
    )
    m = pattern.search(content)
    if not m:
        return []
    items = []
    for line in m.group(1).strip().split('\n'):
        line = line.strip()
        if line.startswith('- '):
            text = re.sub(r'\*\*([^*]+)\*\*', r'\1', line[2:].strip())
            # Skip glossary-style definitions (term + colon + definition)
            # Real objectives start with a verb, not a noun phrase with colon
            if re.match(r'^[A-Za-zÁ-ÿ]+[^:]{0,40}\)\s*:', text):
                continue
            if text.lower().strip() in _NOISE_HEADINGS:
                continue
            items.append(text)
    return items[:6]


def _extract_intro_from_md(content: str) -> str:
    """Extract the first paragraph of ## Introducción / ## Introduction."""
    if not content:
        return ''
    pattern = re.compile(
        r'^#{2,3}\s*(?:Introducción|Introduction)[^\n]*\n+(.+?)(?:\n\n|\n##|\Z)',
        re.MULTILINE | re.IGNORECASE | re.DOTALL,
    )
    m = pattern.search(content)
    if not m:
        return ''
    para = m.group(1).strip().split('\n\n')[0].replace('\n', ' ').strip()
    return para[:300] + '...' if len(para) > 300 else para


def _extract_resumen_items_from_md(content: str) -> list:
    """Extract summary items from Resumen del Capítulo markdown."""
    if not content:
        return []
    pattern = re.compile(
        r'###\s*(?:Temas\s+m[áa]s\s+importantes\s+cubiertos|Key\s+Topics\s+Covered)[^\n]*\n+'
        r'((?:\s*-\s+.+\n?)+)',
        re.MULTILINE | re.IGNORECASE,
    )
    m = pattern.search(content)
    if not m:
        return []
    items = []
    for line in m.group(1).strip().split('\n'):
        line = line.strip()
        if line.startswith('- '):
            text = line[2:].strip()
            # Filter out AI-generated section heading noise
            if text.lower() in _NOISE_HEADINGS:
                continue
            items.append(text)
    return items[:10]


def _parse_glossary_items_from_md(glossary_md: str) -> list:
    """Parse glossary markdown into list of (term, definition) tuples."""
    items = []
    if not glossary_md:
        return items
    seen_terms = set()
    # Noise words that indicate a section heading, not a real glossary term
    _skip_terms = {'próximos pasos', 'recursos adicionales', 'puntos clave',
                   'describir', 'explicar', 'comparar', 'aplicación práctica',
                   'detalles técnicos', 'visión general del concepto'}
    for line in glossary_md.strip().split('\n'):
        line = line.strip()
        if line.startswith('- **'):
            m = re.match(r'^-\s+\*\*([^*]+)\*\*:\s*(.+)$', line)
            if m:
                term = m.group(1).strip()
                defn = m.group(2).strip()
                # Skip question-style entries, pure verbs, or section headings
                term_lower = term.lower().replace('¿', '').replace('?', '').strip()
                if term_lower in _skip_terms:
                    continue
                if term.startswith('¿') or term.endswith('?'):
                    continue
                if term_lower in seen_terms:
                    continue
                seen_terms.add(term_lower)
                # Clean definition
                defn = re.sub(r'!\[[^\]]*\]\([^\)]+\)', '', defn)
                defn = re.sub(r'#{2,}\s*', '', defn)
                defn = re.sub(r'\*\*([^*]+)\*\*', r'\1', defn).strip()
                defn = re.sub(r'^[-–—]+\s*', '', defn).strip()
                # Strip leading term echo
                for prefix in [term, term.split(':')[0].strip()]:
                    if defn.lower().startswith(prefix.lower()):
                        defn = defn[len(prefix):].lstrip(' :,.-–—').strip()
                # If definition still looks like noise, drop it
                if defn.lower().replace('¿', '').replace('?', '') == term_lower:
                    defn = ''
                if 'concepto clave abordado' in defn.lower():
                    defn = ''
                if 'objetivos de aprendizaje' in defn.lower():
                    defn = ''
                if len(defn) > 120:
                    defn = defn[:117] + '...'
                items.append((term, defn))
    return items


def _build_supplementary_lookup(book_data: dict) -> dict:
    """Build a lookup dict keyed by module number with title, objectives, lessons.
    Used to enrich old HTML that lacks new classes."""
    supp = {}
    if not book_data:
        return supp

    outline_modules = book_data.get('outline_modules', [])
    book_modules = book_data.get('book_modules', [])

    for idx, om in enumerate(outline_modules):
        mod_num = idx + 1
        supp[mod_num] = {
            'title': om.get('title', ''),
            'objectives': [],
            'lessons': [],
            'resumen_items': [],
        }

    for idx, bm in enumerate(book_modules):
        mod_num = bm.get('module_number', idx + 1)
        if mod_num not in supp:
            supp[mod_num] = {'title': bm.get('module_title', ''), 'objectives': [], 'lessons': [], 'resumen_items': []}

        # Extract objectives from the first "Introducción" lesson
        # and summary items from the "Resumen del Capítulo" lesson
        for lesson in bm.get('lessons', []):
            lesson_title = lesson.get('title', '')
            content = lesson.get('content', '')
            supp[mod_num]['lessons'].append({
                'title': lesson_title,
                'content': content,
            })
            if lesson_title.lower().strip() in ('introducción', 'introduction') and not supp[mod_num]['objectives']:
                supp[mod_num]['objectives'] = _extract_objectives_from_md(content)
            if lesson_title.lower().strip() in ('resumen del capítulo', 'chapter summary', 'resumen', 'summary'):
                items = _extract_resumen_items_from_md(content)
                if items:
                    supp[mod_num]['resumen_items'] = items

        # Update title from book if outline didn't provide it
        if not supp[mod_num]['title']:
            supp[mod_num]['title'] = bm.get('module_title', f'Capítulo {mod_num}')

    # Extract glossary from metadata
    metadata = book_data.get('metadata', {}) if book_data else {}
    glossary_md = metadata.get('course_glossary', '')
    if not glossary_md:
        for ss in book_data.get('special_sections', []):
            if isinstance(ss, dict) and 'glosario' in (ss.get('title', '') or '').lower():
                glossary_md = ss.get('content', '')
                break
    supp['_glossary_items'] = _parse_glossary_items_from_md(glossary_md)

    return supp


def convert_html_to_pptx(html_content: str, s3_client=None, course_bucket: str = None, book_data: dict = None) -> bytes:
    """
    Convert HTML infographic slides to PowerPoint.
    book_data: optional dict with 'outline_modules' and 'book_modules' for
               enriching module-title / lesson-title slides from S3 data.
    """
    logger.info("Starting HTML to PPT conversion with style mapping...")
    
    soup = BeautifulSoup(html_content, 'html.parser')
    slides_html = soup.find_all('div', class_='slide')
    
    logger.info(f"Found {len(slides_html)} slides in HTML")
    
    logo_bytes = download_logo(s3_client, course_bucket)
    
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    blank_layout = prs.slide_layouts[6]
    
    # Store references for image downloads
    ctx = {'s3_client': s3_client, 'bucket': course_bucket}

    # Build supplementary lookup for backward-compatible enrichment
    supp = _build_supplementary_lookup(book_data) if book_data else {}
    current_module_num = 0  # tracks which module we're in (parsed from text)
    lesson_counter_in_module = 0  # counts lesson-title slides within current module
    seen_module_nums = set()  # track first-occurrence module-title slides
    # Bookend lesson titles to skip
    _bookend_titles = {'introducción', 'introduction', 'resumen del capítulo',
                       'chapter summary', 'resumen', 'summary'}
    _skipping_bookend = False  # when True, skip content slides until next lesson/module title
    _injected_summaries = set()  # track which modules already got a chapter summary

    # Pre-scan: find the index of the last slide (¡Gracias!) to inject glossary before it
    _last_slide_idx = len(slides_html) - 1
    _glossary_injected = False

    for idx, slide_html in enumerate(slides_html):
        logger.info(f"Processing slide {idx + 1}/{len(slides_html)}")
        
        # Check if this is a "structural" slide that ends skip mode
        is_structural = (slide_html.find(class_='module-title-slide') or
                         slide_html.find(class_='lesson-title-slide') or
                         slide_html.find(class_='course-title-slide') or
                         slide_html.find(class_='module-end-logo-slide') or
                         slide_html.find(class_='chapter-summary-slide') or
                         slide_html.find(class_='lab-intro-slide') or
                         slide_html.find(class_='glossary-slide') or
                         slide_html.find(class_='gracias-slide'))
        if is_structural:
            _skipping_bookend = False  # always reset on structural slides

        # If we're in skip mode (after a bookend lesson-title), skip content slides
        if _skipping_bookend and not is_structural:
            logger.info(f"   Skipping bookend content slide {idx + 1}")
            continue

        slide = None
        if slide_html.find(class_='intro-cover-slide'):
            slide = create_intro_cover_slide(prs, blank_layout, slide_html, logo_bytes, ctx)
        elif slide_html.find(class_='intro-legal-slide'):
            slide = create_intro_legal_slide(prs, blank_layout, slide_html, logo_bytes, ctx)
        elif slide_html.find(class_='intro-agenda-slide'):
            slide = create_intro_agenda_slide(prs, blank_layout, slide_html, logo_bytes)
        elif slide_html.find(class_='intro-content-slide'):
            slide = create_intro_content_slide(prs, blank_layout, slide_html, logo_bytes, ctx)
        elif slide_html.find(class_='chapter-summary-slide'):
            slide = create_chapter_summary_slide(prs, blank_layout, slide_html, logo_bytes, ctx)
        elif slide_html.find(class_='lab-intro-slide'):
            slide = create_lab_intro_slide(prs, blank_layout, slide_html, logo_bytes, ctx)
        elif slide_html.find(class_='glossary-slide'):
            slide = create_glossary_slide(prs, blank_layout, slide_html, logo_bytes)
        elif slide_html.find(class_='gracias-slide'):
            # Inject chapter summary for last module if not done yet
            if supp and current_module_num and current_module_num not in _injected_summaries:
                _ri = supp.get(current_module_num, {}).get('resumen_items', [])
                if _ri:
                    _injected_summaries.add(current_module_num)
                    _inject_chapter_summary(prs, blank_layout, logo_bytes, ctx,
                                            current_module_num, _ri,
                                            supp.get(current_module_num, {}).get('title', ''))
                    logger.info(f"   📋 Injected chapter summary for final module {current_module_num}")
            # Inject glossary before gracias if not already done
            if not _glossary_injected and supp:
                _inject_glossary(prs, blank_layout, logo_bytes, supp)
                _glossary_injected = True
            slide = create_gracias_slide(prs, blank_layout, slide_html, logo_bytes, ctx)
        elif slide_html.find(class_='course-title-slide'):
            # Backward compat: detect ¡Gracias! and route to gracias slide
            _ct_title = slide_html.find(class_='title')
            _ct_text = _ct_title.get_text(strip=True) if _ct_title else ''
            if 'gracias' in _ct_text.lower():
                # Inject chapter summary for last module if not done yet
                if supp and current_module_num and current_module_num not in _injected_summaries:
                    _ri = supp.get(current_module_num, {}).get('resumen_items', [])
                    if _ri:
                        _injected_summaries.add(current_module_num)
                        _inject_chapter_summary(prs, blank_layout, logo_bytes, ctx,
                                                current_module_num, _ri,
                                                supp.get(current_module_num, {}).get('title', ''))
                        logger.info(f"   📋 Injected chapter summary for final module {current_module_num}")
                # Inject glossary before gracias if not already done
                if not _glossary_injected and supp:
                    _inject_glossary(prs, blank_layout, logo_bytes, supp)
                    _glossary_injected = True
                slide = create_gracias_slide(prs, blank_layout, slide_html, logo_bytes, ctx)
            else:
                slide = create_course_title_slide(prs, blank_layout, slide_html, logo_bytes)
        elif slide_html.find(class_='module-title-slide'):
            # Parse module number from text instead of using sequential counter
            _mt = slide_html.find(class_='module-title-chapter') or slide_html.find(class_='title')
            _mt_text = _mt.get_text(strip=True) if _mt else ''
            _mn = re.search(r'(\d+)', _mt_text)
            if _mn:
                current_module_num = int(_mn.group(1))
            else:
                current_module_num += 1
            # Only create module-title on FIRST occurrence; skip duplicates
            if current_module_num not in seen_module_nums:
                lesson_counter_in_module = 0
                seen_module_nums.add(current_module_num)
                slide = create_module_title_slide(prs, blank_layout, slide_html, logo_bytes,
                                                  supp=supp, module_num=current_module_num)
            else:
                logger.info(f"   Skipping duplicate module-title for module {current_module_num}")
                _skipping_bookend = True  # also skip content after duplicate module-title
                continue
        elif slide_html.find(class_='lesson-title-slide'):
            # Skip bookend lessons (Introducción, Resumen del Capítulo)
            _lt_elem = slide_html.find(class_='title')
            _lt_text = _lt_elem.get_text(strip=True) if _lt_elem else ''
            # Strip any numbering prefix to check bare title
            _lt_bare = re.sub(r'^\d+\.\d+[:\s]+', '', _lt_text).strip().lower()
            if _lt_bare in _bookend_titles:
                logger.info(f"   Skipping bookend lesson: {_lt_text}")
                _skipping_bookend = True  # skip following content slides too
                continue
            lesson_counter_in_module += 1
            slide = create_lesson_title_slide(prs, blank_layout, slide_html, logo_bytes,
                                              supp=supp, module_num=current_module_num,
                                              lesson_num=lesson_counter_in_module)
        elif slide_html.find(class_='module-end-logo-slide'):
            # Inject chapter summary BEFORE the module-end logo if we have data
            _resumen_items = supp.get(current_module_num, {}).get('resumen_items', []) if supp else []
            if _resumen_items and current_module_num not in _injected_summaries:
                _injected_summaries.add(current_module_num)
                _inject_chapter_summary(prs, blank_layout, logo_bytes, ctx,
                                        current_module_num, _resumen_items,
                                        supp.get(current_module_num, {}).get('title', ''))
                logger.info(f"   📋 Injected chapter summary for module {current_module_num}")
            slide = create_module_end_logo_slide(prs, blank_layout, slide_html, logo_bytes)
        else:
            # Read the content-slide title for backward-compat checks
            _content_title_elem = slide_html.find(class_='slide-title') or slide_html.find(class_='slide-header') or slide_html.find('h2')
            _content_title_text = _content_title_elem.get_text(strip=True) if _content_title_elem else ''

            # Backward compat: skip per-lesson "Resumen" content slides
            if _content_title_text.lower().startswith('resumen:') or _content_title_text.lower().startswith('resumen -'):
                logger.info(f"   ⏭️ Skipping per-lesson summary slide: {_content_title_text}")
                continue
            if _content_title_text.lower() in ('resumen', 'puntos clave'):
                logger.info(f"   ⏭️ Skipping per-lesson summary slide: {_content_title_text}")
                continue

            # Backward compat: detect old lab overview slides (subtitle "Actividad Práctica")
            _subtitle_elem = slide_html.find(class_='slide-subtitle')
            _subtitle_text = _subtitle_elem.get_text(strip=True).lower() if _subtitle_elem else ''
            if 'actividad' in _subtitle_text or _content_title_text.lower().startswith('laboratorio'):
                logger.info(f"   🔬 Converting old lab overview to lab-intro: {_content_title_text}")
                slide = create_lab_intro_slide(prs, blank_layout, slide_html, logo_bytes, ctx)
            else:
                slide = create_content_slide(prs, blank_layout, slide_html, logo_bytes, ctx)
            
        # Add instructor notes if present
        if slide:
            notes_div = slide_html.find(class_='notes')
            if notes_div:
                notes_text = notes_div.get_text(strip=True)
                if notes_text:
                    try:
                        slide.notes_slide.notes_text_frame.text = notes_text
                        logger.info(f"   📝 Added instructor notes ({len(notes_text)} chars)")
                    except Exception as e:
                        logger.warning(f"   ⚠️ Could not add notes: {e}")
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    logger.info(f"✅ Created PowerPoint with {len(slides_html)} slides")
    return output.read()


def download_logo(s3_client, course_bucket: str) -> bytes:
    """Download logo from S3."""
    if not s3_client or not course_bucket:
        logger.warning("No S3 client or bucket provided, skipping logo")
        return None
    
    try:
        logo_key = "logo/LogoNetec.png"
        logger.info(f"Downloading logo from s3://{course_bucket}/{logo_key}")
        response = s3_client.get_object(Bucket=course_bucket, Key=logo_key)
        return response['Body'].read()
    except Exception as e:
        logger.warning(f"Could not download logo: {e}")
        return None


def add_logo(slide, logo_bytes):
    """Add logo to bottom-right corner of slide."""
    if not logo_bytes:
        return
    
    try:
        logo_stream = io.BytesIO(logo_bytes)
        left = SLIDE_WIDTH - LOGO_WIDTH - LOGO_RIGHT_MARGIN
        top = SLIDE_HEIGHT - LOGO_HEIGHT - LOGO_BOTTOM_MARGIN
        slide.shapes.add_picture(logo_stream, left, top, LOGO_WIDTH, LOGO_HEIGHT)
    except Exception as e:
        logger.warning(f"Could not add logo: {e}")


def add_logo_centered(slide, logo_bytes):
    """Add logo centered at bottom of slide (for title slides)."""
    if not logo_bytes:
        return
    
    try:
        logo_stream = io.BytesIO(logo_bytes)
        logo_w = Inches(2)  # Larger logo for title slides
        logo_h = Inches(0.8)
        left = (SLIDE_WIDTH - logo_w) / 2
        top = SLIDE_HEIGHT - logo_h - Inches(0.5)
        slide.shapes.add_picture(logo_stream, left, top, logo_w, logo_h)
    except Exception as e:
        logger.warning(f"Could not add centered logo: {e}")


def extract_s3_key_from_url(url: str) -> str:
    """Extract S3 key from presigned URL or direct S3 URL."""
    try:
        # URL format: https://bucket.s3.amazonaws.com/folder/images/file.png?...
        from urllib.parse import urlparse, unquote
        parsed = urlparse(url)
        # Get path after bucket name
        path = unquote(parsed.path)
        if path.startswith('/'):
            path = path[1:]
        return path
    except Exception as e:
        logger.warning(f"Could not extract S3 key: {e}")
        return None


def download_image_from_s3(s3_client, bucket: str, image_url: str) -> bytes:
    """Download image from S3 using the S3 client (avoids expired presigned URLs)."""
    if not s3_client or not bucket:
        return None
    
    s3_key = extract_s3_key_from_url(image_url)
    if not s3_key:
        return None
    
    try:
        logger.info(f"Downloading image from S3: {s3_key}")
        response = s3_client.get_object(Bucket=bucket, Key=s3_key)
        return response['Body'].read()
    except Exception as e:
        logger.warning(f"Could not download image from S3: {e}")
        return None


def add_content_image(slide, img_url: str, position: str, ctx: dict):
    """Add content image to slide."""
    img_bytes = download_image_from_s3(ctx.get('s3_client'), ctx.get('bucket'), img_url)
    if not img_bytes:
        logger.warning(f"Could not download image: {img_url[:60]}...")
        return
    
    try:
        img_stream = io.BytesIO(img_bytes)
        
        if position == 'left':
            left = CONTENT_LEFT
            top = ctx.get('content_top', CONTENT_TOP)
            width = Inches(5)
            height = Inches(4)
        else:
            left = Inches(7)
            top = ctx.get('content_top', CONTENT_TOP)
            width = Inches(5)
            height = Inches(4)
        
        slide.shapes.add_picture(img_stream, left, top, width, height)
        logger.info(f"✅ Added image at {position}")
    except Exception as e:
        logger.warning(f"Could not add image: {e}")


def download_image_bytes(img_url: str, ctx: dict) -> bytes:
    """Download an image from S3-backed URL or HTTP URL."""
    if not img_url:
        return None

    # Prefer S3 direct read when possible (works for presigned and S3 URLs)
    img_bytes = download_image_from_s3(ctx.get('s3_client'), ctx.get('bucket'), img_url)
    if img_bytes:
        return img_bytes

    try:
        response = requests.get(img_url, timeout=15)
        if response.status_code == 200:
            return response.content
    except Exception as e:
        logger.warning(f"Could not download image over HTTP: {e}")

    return None


def add_logo_bottom_left(slide, logo_bytes):
    """Add Netec logo at bottom-left (intro style)."""
    if not logo_bytes:
        return

    try:
        logo_stream = io.BytesIO(logo_bytes)
        slide.shapes.add_picture(
            logo_stream,
            Inches(0.22),
            Inches(6.58),
            width=Inches(1.95),
            height=Inches(0.73)
        )
    except Exception as e:
        logger.warning(f"Could not add bottom-left logo: {e}")


def _safe_add_picture(slide, image_bytes: bytes, left, top, width=None, height=None) -> bool:
    """Safely add image to slide, returning False instead of raising on unsupported formats."""
    if not image_bytes:
        return False
    try:
        stream = io.BytesIO(image_bytes)
        if width is not None and height is not None:
            slide.shapes.add_picture(stream, left, top, width=width, height=height)
        elif width is not None:
            slide.shapes.add_picture(stream, left, top, width=width)
        elif height is not None:
            slide.shapes.add_picture(stream, left, top, height=height)
        else:
            slide.shapes.add_picture(stream, left, top)
        return True
    except Exception as e:
        logger.warning(f"Could not place image on slide: {e}")
        return False


def add_intro_accent_bar(slide):
    """Add the yellow accent line in a consistent position for all intro slides."""
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.53), Inches(0.46), Inches(0.82), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['bullet_marker']
    accent.line.fill.background()


def _add_agenda_item_paragraph(text_frame, item_text: str):
    """Add one agenda row with yellow bullet, bold chapter prefix, and normal remainder."""
    raw = (item_text or '').strip()
    if not raw:
        return

    match = re.match(r'^\s*((?:Cap[ií]tulo|Chapter|Module)\s+\d+\s*:)\s*(.*)$', raw, flags=re.IGNORECASE)

    para = text_frame.paragraphs[0] if not text_frame.paragraphs[0].text and len(text_frame.paragraphs) == 1 else text_frame.add_paragraph()
    para.level = 0
    para.space_after = Pt(7)

    bullet_run = para.add_run()
    bullet_run.text = "• "
    bullet_run.font.size = Pt(20)
    bullet_run.font.bold = True
    bullet_run.font.name = FONTS['body']
    bullet_run.font.color.rgb = COLORS['bullet_marker']

    if match:
        prefix = match.group(1).strip()
        rest = match.group(2).strip()

        prefix_run = para.add_run()
        prefix_run.text = prefix + (" " if rest else "")
        prefix_run.font.size = Pt(22)
        prefix_run.font.bold = True
        prefix_run.font.name = FONTS['body']
        prefix_run.font.color.rgb = RGBColor(20, 20, 20)

        if rest:
            rest_run = para.add_run()
            rest_run.text = rest
            rest_run.font.size = Pt(22)
            rest_run.font.bold = False
            rest_run.font.name = FONTS['body']
            rest_run.font.color.rgb = RGBColor(20, 20, 20)
    else:
        full_run = para.add_run()
        full_run.text = raw
        full_run.font.size = Pt(22)
        full_run.font.bold = False
        full_run.font.name = FONTS['body']
        full_run.font.color.rgb = RGBColor(20, 20, 20)


def create_intro_cover_slide(prs, layout, slide_html, logo_bytes, ctx):
    """Corporate cover: left title panel, right hero image, countries strip and contact text."""
    slide = prs.slides.add_slide(layout)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Left panel – same colour as background so no visible seam
    left_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(6.3), Inches(7.5))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = COLORS['white']
    left_panel.line.fill.background()

    # Accent bar (fixed position across intro slides)
    add_intro_accent_bar(slide)

    # Title
    title_elem = slide_html.find(class_='intro-cover-title')
    title_text = title_elem.get_text(strip=True) if title_elem else "Curso"
    title_box = slide.shapes.add_textbox(Inches(0.58), Inches(2.05), Inches(5.75), Inches(3.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = RGBColor(0, 60, 120)
    p.alignment = PP_ALIGN.LEFT

    # Right hero image
    hero = slide_html.find('img', class_='intro-cover-main-image')
    hero_bytes = download_image_bytes(hero.get('src') if hero else '', ctx)
    _safe_add_picture(slide, hero_bytes, Inches(6.02), Inches(0.9), width=Inches(6.93), height=Inches(5.75))

    # Countries strip
    countries = slide_html.find('img', class_='intro-cover-countries')
    countries_bytes = download_image_bytes(countries.get('src') if countries else '', ctx)
    _safe_add_picture(slide, countries_bytes, Inches(9.2), Inches(6.53), width=Inches(3.5), height=Inches(0.55))

    # Contact text
    contact = slide_html.find(class_='intro-cover-contact')
    contact_text = contact.get_text(strip=True) if contact else "www.netec.com | servicio@netec.com"
    if contact_text:
        t = slide.shapes.add_textbox(Inches(8.3), Inches(7.0), Inches(4.95), Inches(0.28))
        p = t.text_frame.paragraphs[0]
        p.text = contact_text
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.font.name = FONTS['body']
        p.alignment = PP_ALIGN.RIGHT

    add_logo_bottom_left(slide, logo_bytes)
    return slide


def create_intro_legal_slide(prs, layout, slide_html, logo_bytes, ctx):
    """Corporate legal slide with left icon and right legal paragraphs."""
    slide = prs.slides.add_slide(layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    icon = slide_html.find('img', class_='intro-legal-icon')
    icon_bytes = download_image_bytes(icon.get('src') if icon else '', ctx)
    _safe_add_picture(slide, icon_bytes, Inches(0.9), Inches(1.8), width=Inches(3.2), height=Inches(3.9))

    # Accent bar + title
    add_intro_accent_bar(slide)

    title_elem = slide_html.find(class_='intro-content-title')
    title_text = title_elem.get_text(strip=True) if title_elem else "Propiedad intelectual"
    title_box = slide.shapes.add_textbox(Inches(5.85), Inches(1.0), Inches(6.7), Inches(1.3))
    p = title_box.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.85), Inches(2.35), Inches(6.95), Inches(0.02))
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(195, 195, 195)
    divider.line.fill.background()

    paragraphs = [p.get_text(strip=True) for p in slide_html.select('.intro-paragraphs p') if p.get_text(strip=True)]
    text_box = slide.shapes.add_textbox(Inches(5.85), Inches(2.45), Inches(6.9), Inches(4.7))
    tf = text_box.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = txt
        para.font.size = Pt(18)
        para.font.name = FONTS['body']
        para.font.color.rgb = RGBColor(25, 25, 25)
        para.space_after = Pt(8)

    add_logo_bottom_left(slide, logo_bytes)
    return slide


def create_intro_content_slide(prs, layout, slide_html, logo_bytes, ctx):
    """Corporate content intro slide: left text/list and right asset image."""
    slide = prs.slides.add_slide(layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Accent bar
    add_intro_accent_bar(slide)

    title_elem = slide_html.find(class_='intro-content-title')
    title_text = title_elem.get_text(strip=True) if title_elem else ""
    subtitle_elem = slide_html.find(class_='intro-subtitle')
    subtitle_text = subtitle_elem.get_text(strip=True) if subtitle_elem else ""

    title_box = slide.shapes.add_textbox(Inches(0.83), Inches(1.0), Inches(6.2), Inches(1.2))
    p = title_box.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = RGBColor(0, 0, 0)

    current_top = Inches(2.2)
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.83), current_top, Inches(6.2), Inches(0.7))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(22)
        sp.font.name = FONTS['body']
        sp.font.color.rgb = RGBColor(20, 20, 20)
        current_top = Inches(2.9)

    # Divider
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.83), current_top, Inches(6.6), Inches(0.02))
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(195, 195, 195)
    divider.line.fill.background()
    current_top = current_top + Inches(0.12)

    list_items = [li.get_text(strip=True) for li in slide_html.select('ul.intro-list li') if li.get_text(strip=True)]
    text_box = slide.shapes.add_textbox(Inches(0.9), current_top, Inches(6.35), Inches(4.4))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for txt in list_items:
        para = tf.paragraphs[0] if not tf.paragraphs[0].text and len(tf.paragraphs) == 1 else tf.add_paragraph()
        para.level = 0

        bullet_run = para.add_run()
        bullet_run.text = "• "
        bullet_run.font.size = Pt(22)
        bullet_run.font.bold = True
        bullet_run.font.name = FONTS['body']
        bullet_run.font.color.rgb = COLORS['bullet_marker']

        text_run = para.add_run()
        text_run.text = txt
        text_run.font.size = Pt(21)
        text_run.font.bold = False
        text_run.font.name = FONTS['body']
        text_run.font.color.rgb = RGBColor(20, 20, 20)

        para.space_after = Pt(6)

    asset = slide_html.find('img', class_='intro-right-asset')
    asset_bytes = download_image_bytes(asset.get('src') if asset else '', ctx)
    _safe_add_picture(slide, asset_bytes, Inches(7.6), Inches(1.0), width=Inches(5.2), height=Inches(5.5))

    add_logo_bottom_left(slide, logo_bytes)
    return slide


def create_intro_agenda_slide(prs, layout, slide_html, logo_bytes):
    """Corporate agenda/temario slide with large chapter list and max readability."""
    slide = prs.slides.add_slide(layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    add_intro_accent_bar(slide)

    # Soft horizontal separator
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.82), Inches(13.333), Inches(0.03))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(220, 220, 220)
    sep.line.fill.background()

    title_elem = slide_html.find(class_='intro-content-title')
    title_text = title_elem.get_text(strip=True) if title_elem else "Temario"
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.66), Inches(6.5), Inches(0.95))
    p = title_box.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = RGBColor(0, 0, 0)

    agenda_items = [li.get_text(strip=True) for li in slide_html.select('ul.intro-agenda-list li') if li.get_text(strip=True)]
    list_box = slide.shapes.add_textbox(Inches(1.08), Inches(2.34), Inches(11.0), Inches(3.95))
    tf = list_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for txt in agenda_items:
        _add_agenda_item_paragraph(tf, txt)

    add_logo_bottom_left(slide, logo_bytes)
    return slide


# =============================================================================
# ASSET DOWNLOAD HELPER
# =============================================================================

def _download_asset_image(s3_client, bucket, filename):
    """Download an asset image from S3 logo/Assets/ folder. Returns bytes or None."""
    if not s3_client or not bucket:
        return None
    try:
        key = f"logo/Assets/{filename}"
        response = s3_client.get_object(Bucket=bucket, Key=key)
        return response['Body'].read()
    except Exception as e:
        logger.warning(f"Could not download asset {filename}: {e}")
        return None


def _inject_chapter_summary(prs, layout, logo_bytes, ctx, module_num, resumen_items, module_title):
    """Inject a chapter summary slide programmatically (backward compat for old HTML).
    Layout: title top-left, bullet list left, Resumen_Capitulo.png right, yellow bar at bottom."""
    slide = prs.slides.add_slide(layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_text = "Resumen del capítulo"
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(7.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = RGBColor(34, 34, 34)

    # Summary items
    if resumen_items:
        item_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(7.0), Inches(4.5))
        tf = item_box.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, item in enumerate(resumen_items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            bullet_run = para.add_run()
            bullet_run.text = "• "
            bullet_run.font.size = Pt(17)
            bullet_run.font.bold = True
            bullet_run.font.name = FONTS['body']
            bullet_run.font.color.rgb = RGBColor(34, 34, 34)
            text_run = para.add_run()
            text_run.text = item
            text_run.font.size = Pt(17)
            text_run.font.name = FONTS['body']
            text_run.font.color.rgb = RGBColor(34, 34, 34)
            para.space_after = Pt(6)

    # Resumen_Capitulo.png on right (vertically centered)
    asset_bytes = _download_asset_image(ctx.get('s3_client'), ctx.get('bucket'), 'Resumen_Capitulo.png')
    if asset_bytes:
        try:
            img_stream = io.BytesIO(asset_bytes)
            slide.shapes.add_picture(img_stream, Inches(8.5), Inches(1.8), Inches(3.5), Inches(3.5))
        except Exception as e:
            logger.warning(f"Could not add Resumen_Capitulo.png: {e}")

    # Gray divider line near bottom
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.015)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(199, 199, 199)
    divider.line.fill.background()

    # Yellow accent bar at bottom (partial width, centered-left)
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.15), Inches(4.5), Inches(0.08)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = COLORS['bullet_marker']
    bottom_bar.line.fill.background()

    add_logo_bottom_left(slide, logo_bytes)


def _inject_glossary(prs, layout, logo_bytes, supp):
    """Inject glossary slide(s) programmatically (backward compat for old HTML)."""
    glossary_items = supp.get('_glossary_items', [])
    if not glossary_items:
        return

    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Yellow bar at TOP
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['bullet_marker']
    top_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Glosario"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = COLORS['primary']

    # Two-column layout
    mid = (len(glossary_items) + 1) // 2
    col1_items = glossary_items[:mid]
    col2_items = glossary_items[mid:]

    for col_idx, col_items in enumerate([col1_items, col2_items]):
        left = Inches(0.6) if col_idx == 0 else Inches(6.8)
        box = slide.shapes.add_textbox(left, Inches(1.5), Inches(5.8), Inches(5.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, (term, defn) in enumerate(col_items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            br = para.add_run()
            br.text = "• "
            br.font.size = Pt(15)
            br.font.bold = True
            br.font.name = FONTS['body']
            br.font.color.rgb = COLORS['bullet_marker']
            tr = para.add_run()
            tr.text = term
            tr.font.size = Pt(15)
            tr.font.bold = True
            tr.font.name = FONTS['body']
            tr.font.color.rgb = COLORS['primary']
            if defn:
                dr = para.add_run()
                dr.text = f" {defn}"
                dr.font.size = Pt(14)
                dr.font.name = FONTS['body']
                dr.font.color.rgb = RGBColor(68, 68, 68)
            para.space_after = Pt(3)

    add_logo_bottom_left(slide, logo_bytes)
    logger.info(f"   📖 Injected glossary slide with {len(glossary_items)} terms")


# =============================================================================
# CHAPTER SUMMARY SLIDE
# =============================================================================

def create_chapter_summary_slide(prs, layout, slide_html, logo_bytes, ctx):
    """Chapter summary: white bg, title top-left, bullets left, Resumen_Capitulo.png right,
    yellow accent bar at bottom."""
    slide = prs.slides.add_slide(layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title
    title_elem = slide_html.find(class_='chapter-summary-title')
    title_text = title_elem.get_text(strip=True) if title_elem else "Resumen del capítulo"

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(7.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = RGBColor(34, 34, 34)

    # Summary bullet items
    list_elem = slide_html.find(class_='chapter-summary-list')
    items = []
    if list_elem:
        items = [li.get_text(strip=True) for li in list_elem.find_all('li') if li.get_text(strip=True)]

    if items:
        item_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(7.0), Inches(4.5))
        tf = item_box.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, item in enumerate(items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            bullet_run = para.add_run()
            bullet_run.text = "• "
            bullet_run.font.size = Pt(17)
            bullet_run.font.bold = True
            bullet_run.font.name = FONTS['body']
            bullet_run.font.color.rgb = RGBColor(34, 34, 34)
            text_run = para.add_run()
            text_run.text = item
            text_run.font.size = Pt(17)
            text_run.font.name = FONTS['body']
            text_run.font.color.rgb = RGBColor(34, 34, 34)
            para.space_after = Pt(6)

    # Resumen_Capitulo.png on right side (vertically centered)
    asset_bytes = _download_asset_image(ctx.get('s3_client'), ctx.get('bucket'), 'Resumen_Capitulo.png')
    if asset_bytes:
        try:
            img_stream = io.BytesIO(asset_bytes)
            slide.shapes.add_picture(img_stream, Inches(8.5), Inches(1.8), Inches(3.5), Inches(3.5))
        except Exception as e:
            logger.warning(f"Could not add Resumen_Capitulo.png: {e}")

    # Gray divider line near bottom
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.015)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(199, 199, 199)
    divider.line.fill.background()

    # Yellow accent bar at bottom (partial width)
    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.15), Inches(4.5), Inches(0.08)
    )
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = COLORS['bullet_marker']
    bottom_bar.line.fill.background()

    add_logo_bottom_left(slide, logo_bytes)
    return slide


# =============================================================================
# LAB INTRO SLIDE
# =============================================================================

def create_lab_intro_slide(prs, layout, slide_html, logo_bytes, ctx):
    """Lab intro: white bg, dashed title box, yellow accent, objective, clock bottom-right."""
    slide = prs.slides.add_slide(layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Dashed title box (top area — taller, pushed down)
    from pptx.oxml.ns import qn
    title_box_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.5), Inches(12.1), Inches(2.3)
    )
    title_box_shape.fill.background()  # No fill
    ln = title_box_shape.line
    ln.color.rgb = RGBColor(170, 170, 170)
    ln.width = Pt(1.5)
    ln.dash_style = 3  # dash

    # Title text (centered in box — larger font)
    title_elem = slide_html.find(class_='lab-intro-title') or slide_html.find(class_='slide-title')
    title_text = title_elem.get_text(strip=True) if title_elem else "Lab Activity"
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.7), Inches(1.6))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = RGBColor(17, 17, 17)
    p.alignment = PP_ALIGN.CENTER

    # Yellow accent bar below title box
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.9), Inches(2.4), Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['bullet_marker']
    accent.line.fill.background()

    # Gray divider line
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.1), Inches(12.1), Inches(0.02)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(199, 199, 199)
    div.line.fill.background()

    # Objective section
    obj_elem = slide_html.find(class_='lab-intro-section-body')
    if not obj_elem:
        obj_elem = slide_html.find(class_='lab-intro-objective')
    if not obj_elem:
        for li in slide_html.find_all('li'):
            li_text = li.get_text(strip=True).lower()
            if not li_text.startswith('tiempo') and len(li_text) > 20:
                obj_elem = li
                break

    # Heading label (pushed down to match taller title box)
    heading_elem = slide_html.find(class_='lab-intro-section-heading')
    heading_text = heading_elem.get_text(strip=True) if heading_elem else 'Objetivo:'
    h_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.5), Inches(8.0), Inches(0.5))
    hp = h_box.text_frame.paragraphs[0]
    hp.text = heading_text
    hp.font.size = Pt(18)
    hp.font.bold = True
    hp.font.name = FONTS['title']
    hp.font.color.rgb = RGBColor(34, 34, 34)

    if obj_elem:
        obj_text = obj_elem.get_text(strip=True)
        # Strip checkbox artifacts like '[ ] '
        import re as _re
        obj_text = _re.sub(r'^\[\s*\]\s*', '', obj_text)
        obj_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8.0), Inches(2.0))
        otf = obj_box.text_frame
        otf.word_wrap = True
        op = otf.paragraphs[0]
        op.text = obj_text
        op.font.size = Pt(16)
        op.font.name = FONTS['body']
        op.font.color.rgb = RGBColor(51, 51, 51)
        # Left border effect via a thin shape
        border_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.0), Inches(0.04), Inches(0.6)
        )
        border_line.fill.solid()
        border_line.fill.fore_color.rgb = RGBColor(199, 199, 199)
        border_line.line.fill.background()

    # Clock image + duration (bottom-right)
    asset_bytes = _download_asset_image(ctx.get('s3_client'), ctx.get('bucket'), 'Reloj.png')
    if asset_bytes:
        try:
            img_stream = io.BytesIO(asset_bytes)
            slide.shapes.add_picture(img_stream, Inches(8.5), Inches(5.8), Inches(1.0), Inches(1.0))
        except Exception as e:
            logger.warning(f"Could not add Reloj.png: {e}")

    dur_elem = slide_html.find(class_='lab-intro-duration-label')
    dur_val_elem = slide_html.find(class_='lab-intro-duration-value')
    dur_box_elem = slide_html.find(class_='lab-intro-duration-box')
    if not dur_elem and not dur_box_elem:
        # Backward compat: look for duration text in list items
        for li in slide_html.find_all('li'):
            if 'tiempo' in li.get_text(strip=True).lower():
                dur_elem = li
                break
    if dur_elem or dur_box_elem:
        dur_label_text = dur_elem.get_text(strip=True) if dur_elem else 'Tiempo para esta actividad:'
        dur_value_text = dur_val_elem.get_text(strip=True) if dur_val_elem else ''
        d_box = slide.shapes.add_textbox(Inches(9.6), Inches(5.7), Inches(3.5), Inches(0.5))
        dp = d_box.text_frame.paragraphs[0]
        dp.text = dur_label_text
        dp.font.size = Pt(16)
        dp.font.bold = True
        dp.font.name = FONTS['body']
        dp.font.color.rgb = RGBColor(17, 17, 17)
        if dur_value_text:
            dp2 = d_box.text_frame.add_paragraph()
            dp2.text = dur_value_text
            dp2.font.size = Pt(15)
            dp2.font.name = FONTS['body']
            dp2.font.color.rgb = RGBColor(51, 51, 51)

    add_logo_bottom_left(slide, logo_bytes)
    return slide


# =============================================================================
# GLOSSARY SLIDE
# =============================================================================

def create_glossary_slide(prs, layout, slide_html, logo_bytes):
    """Glossary: white bg, yellow top bar, two-column term grid."""
    slide = prs.slides.add_slide(layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Yellow bar at TOP
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.12)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLORS['bullet_marker']
    top_bar.line.fill.background()

    # Title
    title_elem = slide_html.find(class_='glossary-title')
    title_text = title_elem.get_text(strip=True) if title_elem else "Glosario"

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = COLORS['primary']

    # Glossary items – two-column layout
    items_html = slide_html.find_all(class_='glossary-item')
    items = []
    for it in items_html:
        term_el = it.find(class_='glossary-term')
        def_el = it.find(class_='glossary-def')
        term = term_el.get_text(strip=True) if term_el else ""
        defn = def_el.get_text(strip=True) if def_el else ""
        if term:
            items.append((term, defn))

    # Split into two columns
    mid = (len(items) + 1) // 2
    col1_items = items[:mid]
    col2_items = items[mid:]

    for col_idx, col_items in enumerate([col1_items, col2_items]):
        left = Inches(0.6) if col_idx == 0 else Inches(6.8)
        box = slide.shapes.add_textbox(left, Inches(1.5), Inches(5.8), Inches(5.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, (term, defn) in enumerate(col_items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Bullet
            br = para.add_run()
            br.text = "• "
            br.font.size = Pt(15)
            br.font.bold = True
            br.font.name = FONTS['body']
            br.font.color.rgb = COLORS['bullet_marker']
            # Term (bold)
            tr = para.add_run()
            tr.text = term
            tr.font.size = Pt(15)
            tr.font.bold = True
            tr.font.name = FONTS['body']
            tr.font.color.rgb = COLORS['primary']
            # Definition
            if defn:
                dr = para.add_run()
                dr.text = f" {defn}"
                dr.font.size = Pt(14)
                dr.font.name = FONTS['body']
                dr.font.color.rgb = RGBColor(68, 68, 68)
            para.space_after = Pt(3)

    add_logo_bottom_left(slide, logo_bytes)
    return slide


# =============================================================================
# GRACIAS (CLOSING) SLIDE
# =============================================================================

def create_gracias_slide(prs, layout, slide_html, logo_bytes, ctx):
    """Gracias slide: split layout — white left panel, hero image right."""
    slide = prs.slides.add_slide(layout)

    # Right side: Gracias.png covering right ~55% of slide
    asset_bytes = _download_asset_image(ctx.get('s3_client'), ctx.get('bucket'), 'Gracias.png')
    if asset_bytes:
        try:
            img_stream = io.BytesIO(asset_bytes)
            # Image covers right portion (from ~45% to right edge)
            img_left = Inches(5.8)
            slide.shapes.add_picture(
                img_stream, img_left, Inches(0), Inches(7.533), SLIDE_HEIGHT
            )
        except Exception as e:
            logger.warning(f"Could not add Gracias.png: {e}")

    # White left panel (~50% width, overlaps slightly for curved feel)
    white_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(6.5), SLIDE_HEIGHT
    )
    white_panel.fill.solid()
    white_panel.fill.fore_color.rgb = COLORS['white']
    white_panel.line.fill.background()

    # Curved edge: white oval that softens the right border of the panel
    curve = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(4.5), Inches(-1.0), Inches(4.0), Inches(9.5)
    )
    curve.fill.solid()
    curve.fill.fore_color.rgb = COLORS['white']
    curve.line.fill.background()

    # Yellow accent bar (top-left)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.6), Inches(0.9), Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['bullet_marker']
    accent.line.fill.background()

    # Title
    title_elem = slide_html.find(class_='gracias-title')
    title_text = title_elem.get_text(strip=True) if title_elem else "¡Gracias!"
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5.0), Inches(1.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = FONTS['title']
    p.alignment = PP_ALIGN.LEFT

    # Divider line
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.7), Inches(1.0), Inches(0.025)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(170, 170, 170)
    div.line.fill.background()

    # Subtitle
    sub_elem = slide_html.find(class_='gracias-subtitle')
    sub_text = sub_elem.get_text(strip=True) if sub_elem else "¿Dudas o comentarios?"
    s_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(5.0), Inches(0.8))
    sf = s_box.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.text = sub_text
    sp.font.size = Pt(22)
    sp.font.color.rgb = RGBColor(68, 68, 68)
    sp.font.name = FONTS['body']
    sp.alignment = PP_ALIGN.LEFT

    # Logo at bottom left
    if logo_bytes:
        logo_stream = io.BytesIO(logo_bytes)
        slide.shapes.add_picture(
            logo_stream, Inches(0.8), Inches(6.3), Inches(1.5), Inches(0.6)
        )

    return slide


def create_course_title_slide(prs, layout, slide_html, logo_bytes):
    """Course title: Full gradient, centered title, centered logo at bottom."""
    slide = prs.slides.add_slide(layout)
    
    title_elem = slide_html.find(class_='title')
    title_text = title_elem.get_text(strip=True) if title_elem else "Untitled"
    
    add_gradient_background(slide)
    
    # Title (centered)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = FONTS['title']
    p.alignment = PP_ALIGN.CENTER
    
    add_logo_centered(slide, logo_bytes)
    return slide


def create_module_title_slide(prs, layout, slide_html, logo_bytes, supp=None, module_num=1):
    """Module title: white bg, yellow accent bar top-left, chapter label centered-left,
    chapter name in blue below divider, objectives top-right (heading left-aligned)."""
    slide = prs.slides.add_slide(layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    # Yellow accent bar (top-left corner)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.53), Inches(0.45), Inches(0.82), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['bullet_marker']
    accent.line.fill.background()

    # Parse chapter label and name from the title element
    title_elem = slide_html.find(class_='module-title-chapter')
    chapter_label = title_elem.get_text(strip=True) if title_elem else ""
    name_elem = slide_html.find(class_='module-title-name')
    chapter_name = name_elem.get_text(strip=True) if name_elem else ""

    # Objectives from new HTML classes
    obj_heading_elem = slide_html.find(class_='module-title-obj-heading')
    obj_list_elem = slide_html.find(class_='module-title-obj-list')
    obj_items = []
    if obj_list_elem:
        for li in obj_list_elem.find_all('li'):
            text = li.get_text(strip=True)
            if not text:
                continue
            # Filter glossary-style definitions (term with parens + colon)
            if re.match(r'^[A-Za-z\u00C1-\u00FF]+[^:]{0,40}\)\s*:', text):
                continue
            if text.lower().strip() in _NOISE_HEADINGS:
                continue
            obj_items.append(text)

    # --- Backward-compatible fallback: use supplementary book data ---
    if not chapter_label:
        old_title = slide_html.find(class_='title')
        full = old_title.get_text(strip=True) if old_title else ""

        m = re.match(r'^((?:Cap[ií]tulo|M[oó]dulo|Module|Chapter)\s+\d+)\s*:\s*(.+)$', full, re.IGNORECASE)
        if m:
            chapter_label = m.group(1).strip()
            chapter_name = m.group(2).strip()
        else:
            # Title is just "Capítulo N" without name – enrich from supp
            chapter_label = full or f"Capítulo {module_num}"
            chapter_name = ""
            if supp and module_num in supp:
                mod_info = supp[module_num]
                # Parse name from the full outline title "Módulo N: Name"
                om = re.match(r'^(?:M[oó]dulo|Cap[ií]tulo|Module|Chapter)\s+\d+\s*:\s*(.+)$',
                              mod_info.get('title', ''), re.IGNORECASE)
                if om:
                    chapter_name = om.group(1).strip()
                elif mod_info.get('title') and mod_info['title'] != chapter_label:
                    chapter_name = mod_info['title']

    if not obj_items and supp and module_num in supp:
        obj_items = supp[module_num].get('objectives', [])

    # Chapter label (large, black) – vertically centered on left half
    lbl_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(5.8), Inches(2.5))
    tf = lbl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = chapter_label
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = RGBColor(17, 17, 17)
    p.alignment = PP_ALIGN.LEFT

    # Half-width divider line (left column only)
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.8), Inches(5.8), Inches(0.02)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(199, 199, 199)
    divider.line.fill.background()

    # Chapter name (corporate blue, below divider)
    if chapter_name:
        name_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.5))
        tf = name_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = chapter_name
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.name = FONTS['title']
        p.font.color.rgb = COLORS['primary']  # Corporate blue
        p.alignment = PP_ALIGN.LEFT

    # Objectives on right side
    obj_heading = "Objetivos:"
    if obj_heading_elem:
        obj_heading = obj_heading_elem.get_text(strip=True)

    if obj_items:
        # Heading – LEFT-ALIGNED to align with bullet text
        h_box = slide.shapes.add_textbox(Inches(6.8), Inches(0.35), Inches(5.5), Inches(0.6))
        hp = h_box.text_frame.paragraphs[0]
        hp.text = obj_heading
        hp.font.size = Pt(22)
        hp.font.bold = True
        hp.font.name = FONTS['title']
        hp.font.color.rgb = RGBColor(17, 17, 17)
        hp.alignment = PP_ALIGN.LEFT

        # Objective items (below heading, with right margin)
        obj_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.0), Inches(5.5), Inches(3.7))
        tf = obj_box.text_frame
        tf.word_wrap = True
        tf.clear()
        for i, item in enumerate(obj_items):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            bullet_run = para.add_run()
            bullet_run.text = "• "
            bullet_run.font.size = Pt(16)
            bullet_run.font.bold = True
            bullet_run.font.name = FONTS['body']
            bullet_run.font.color.rgb = COLORS['bullet_marker']

            text_run = para.add_run()
            text_run.text = item
            text_run.font.size = Pt(16)
            text_run.font.bold = False
            text_run.font.name = FONTS['body']
            text_run.font.color.rgb = RGBColor(34, 34, 34)
            para.space_after = Pt(6)

    add_logo_bottom_left(slide, logo_bytes)
    return slide


def create_lesson_title_slide(prs, layout, slide_html, logo_bytes, supp=None, module_num=1, lesson_num=1):
    """Lesson title: white bg, yellow bar at TOP, numbered title, divider, intro text."""
    slide = prs.slides.add_slide(layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']

    title_elem = slide_html.find(class_='title')
    title_text = title_elem.get_text(strip=True) if title_elem else "Untitled"

    # Don't add any numbering prefix — the book title already contains it (e.g. "1.1: Name")
    # Strip any double-prefixed numbering that old code may have added
    _double_prefix = re.match(r'^\d+\.\d+\s+(\d+\.\d+[:\s].+)$', title_text)
    if _double_prefix:
        title_text = _double_prefix.group(1)
    
    # Yellow bar at TOP
    yellow_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, Inches(0.12)
    )
    yellow_bar.fill.solid()
    yellow_bar.fill.fore_color.rgb = COLORS['bullet_marker']
    yellow_bar.line.fill.background()
    
    # Title (centered, dark)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = RGBColor(17, 17, 17)
    p.font.name = FONTS['title']
    p.alignment = PP_ALIGN.CENTER

    # Introduction text (from new HTML classes)
    intro_heading_elem = slide_html.find(class_='lesson-intro-heading')
    intro_heading = intro_heading_elem.get_text(strip=True) if intro_heading_elem else ''
    intro_items = [li.get_text(strip=True) for li in slide_html.select('.lesson-intro-text p') if li.get_text(strip=True)]

    # Backward-compatible fallback: extract intro from book data via title matching
    if not intro_items and supp and module_num in supp:
        supp_lessons = supp[module_num].get('lessons', [])
        matched_lesson = None

        # Try 1: match by "X.Y" number prefix in the title
        _num_m = re.match(r'(\d+\.\d+)', title_text)
        if _num_m:
            target_num = _num_m.group(1)
            for _sl in supp_lessons:
                if target_num + ':' in _sl['title'] or target_num + ' ' in _sl['title'] or _sl['title'].startswith(target_num):
                    matched_lesson = _sl
                    break

        # Try 2: fuzzy match on stripped title text
        if not matched_lesson:
            clean_title = re.sub(r'^\d+\.\d+[:\s]+', '', title_text).strip().lower()
            if clean_title:
                for _sl in supp_lessons:
                    clean_supp = re.sub(r'^\d+\.\d+[:\s]+', '', _sl['title']).strip().lower()
                    if clean_title == clean_supp or clean_title in clean_supp or clean_supp in clean_title:
                        matched_lesson = _sl
                        break

        if matched_lesson:
            intro_text = _extract_intro_from_md(matched_lesson.get('content', ''))
            if intro_text:
                intro_items = [intro_text]

    if intro_items or intro_heading:
        # Divider (left-aligned)
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(3.5),
            Inches(11.2), Inches(0.02)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(199, 199, 199)
        divider.line.fill.background()

        y_pos = 3.7

        # "Introducción" heading bullet (left-aligned)
        if intro_heading:
            heading_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(y_pos), Inches(11.2), Inches(0.5)
            )
            hp = heading_box.text_frame.paragraphs[0]
            bullet_run = hp.add_run()
            bullet_run.text = "\u2022  "
            bullet_run.font.size = Pt(20)
            bullet_run.font.bold = False
            bullet_run.font.name = FONTS['body']
            bullet_run.font.color.rgb = RGBColor(51, 51, 51)
            text_run = hp.add_run()
            text_run.text = intro_heading
            text_run.font.size = Pt(20)
            text_run.font.bold = False
            text_run.font.name = FONTS['body']
            text_run.font.color.rgb = RGBColor(51, 51, 51)
            hp.alignment = PP_ALIGN.LEFT
            y_pos += 0.5

        if intro_items:
            intro_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(y_pos), Inches(11.5), Inches(2.5)
            )
            tf = intro_box.text_frame
            tf.word_wrap = True
            for i, text in enumerate(intro_items):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = text
                para.font.size = Pt(18)
                para.font.name = FONTS['body']
                para.font.color.rgb = RGBColor(68, 68, 68)
                para.alignment = PP_ALIGN.LEFT
                para.space_after = Pt(6)

    add_logo_bottom_left(slide, logo_bytes)
    return slide


def create_module_end_logo_slide(prs, layout, slide_html, logo_bytes):
    """Module end slide: White background, large centered logo."""
    slide = prs.slides.add_slide(layout)
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    add_logo_middle_center(slide, logo_bytes)
    return slide


def create_content_slide(prs, layout, slide_html, logo_bytes, ctx):
    """Create a content slide with header, bullets, and optional image."""
    slide = prs.slides.add_slide(layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    title_elem = slide_html.find(class_='slide-title')
    title_text = title_elem.get_text(strip=True) if title_elem else "Untitled"
    
    subtitle_elem = slide_html.find(class_='slide-subtitle')
    subtitle_text = subtitle_elem.get_text(strip=True) if subtitle_elem else ""
    
    header_height = add_header_bar(slide, title_text, subtitle_text)
    
    content_top_pos = header_height + Inches(0.2) # Start content below header
    
    content_elem = slide_html.find(class_='slide-content')
    if not content_elem:
        add_logo_bottom_left(slide, logo_bytes)
        return slide
    
    # Check for image layout
    has_image = 'image-left' in content_elem.get('class', []) or 'image-right' in content_elem.get('class', [])
    image_position = 'left' if 'image-left' in content_elem.get('class', []) else 'right'
    
    # Find image in image-column or directly in content
    image_column = content_elem.find(class_='image-column')
    if image_column:
        img_elem = image_column.find('img', class_='slide-image')
    else:
        img_elem = content_elem.find('img', class_='slide-image')
    
    if not img_elem:
        img_elem = slide_html.find('img', class_='slide-image')
    
    if img_elem and img_elem.get('src'):
        img_url = img_elem.get('src')
        logger.info(f"Found slide image: {img_url[:80]}...")
        
        # Pass content_top to image function
        ctx['content_top'] = content_top_pos
        add_content_image(slide, img_url, image_position, ctx)
        
        # Adjust bullet area based on image position
        if image_position == 'left':
            bullet_left = Inches(6)
            bullet_width = Inches(6.5)
        else:
            bullet_left = CONTENT_LEFT
            bullet_width = Inches(6)
    else:
        bullet_left = CONTENT_LEFT
        bullet_width = CONTENT_WIDTH
    
    # Content
    heading_elem = content_elem.find(class_='content-heading')
    bullets_elem = content_elem.find('ul', class_='bullets')
    code_elem = content_elem.find(class_='code-block')
    callout_elem = content_elem.find(class_='callout')
    table_elem = content_elem.find('table', class_='slide-table')
    
    current_top = content_top_pos
    
    if heading_elem:
        current_top = add_content_heading(slide, heading_elem, current_top, bullet_left, bullet_width)

    if table_elem:
        current_top = add_table_block(slide, table_elem, current_top, bullet_left, bullet_width)

    if bullets_elem:
        current_top = add_bullets(slide, bullets_elem, current_top, bullet_left, bullet_width)
    
    if code_elem:
        add_code_block(slide, code_elem, current_top)
    
    if callout_elem:
        add_callout(slide, callout_elem)
    
    add_logo_bottom_left(slide, logo_bytes)
    return slide


def add_table_block(slide, table_elem, top, left=None, width=None):
    """Render an HTML <table class='slide-table'> as a native PowerPoint table."""
    if left is None:
        left = CONTENT_LEFT
    if width is None:
        width = CONTENT_WIDTH

    # Parse headers and rows from the HTML table
    headers = []
    thead = table_elem.find('thead')
    if thead:
        header_row = thead.find('tr')
        if header_row:
            headers = [th.get_text(strip=True) for th in header_row.find_all(['th', 'td'])]

    rows = []
    tbody = table_elem.find('tbody')
    row_sources = tbody.find_all('tr') if tbody else table_elem.find_all('tr')
    for tr in row_sources:
        cells = [td.get_text(strip=True) for td in tr.find_all(['td', 'th'])]
        if cells and cells != headers:  # skip header row if no thead
            rows.append(cells)

    if not headers and not rows:
        return top

    num_cols = max(len(headers), max((len(r) for r in rows), default=0))
    num_rows = (1 if headers else 0) + len(rows)

    if num_cols == 0 or num_rows == 0:
        return top

    # Calculate row height and table height
    row_height = Inches(0.45)
    table_height = row_height * num_rows
    max_table_height = Inches(5.0)
    if table_height > max_table_height:
        table_height = max_table_height
        row_height = table_height / num_rows

    col_width = width / num_cols

    tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, table_height)
    table = tbl_shape.table

    # Style the table
    row_idx = 0

    # Header row
    if headers:
        for c, text in enumerate(headers):
            if c < num_cols:
                cell = table.cell(0, c)
                cell.text = text
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['primary']  # #003366
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = COLORS['white']
                p.font.name = FONTS['body']
                p.alignment = PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        row_idx = 1

    # Data rows
    for r, row_data in enumerate(rows):
        ppt_row = row_idx + r
        if ppt_row >= num_rows:
            break
        is_alt = (r % 2 == 1)
        for c, text in enumerate(row_data):
            if c < num_cols:
                cell = table.cell(ppt_row, c)
                cell.text = text
                # Alternate row shading
                cell.fill.solid()
                if is_alt:
                    cell.fill.fore_color.rgb = RGBColor(240, 244, 249)  # very light blue
                else:
                    cell.fill.fore_color.rgb = COLORS['white']
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = COLORS['text_black']
                p.font.name = FONTS['body']
                p.alignment = PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Fill any remaining empty cells (if rows have unequal lengths)
    for r_idx in range(num_rows):
        for c_idx in range(num_cols):
            cell = table.cell(r_idx, c_idx)
            if not cell.text_frame.paragraphs[0].text:
                cell.fill.solid()
                if r_idx == 0 and headers:
                    cell.fill.fore_color.rgb = COLORS['primary']
                elif (r_idx - (1 if headers else 0)) % 2 == 1:
                    cell.fill.fore_color.rgb = RGBColor(240, 244, 249)
                else:
                    cell.fill.fore_color.rgb = COLORS['white']

    logger.info(f"   📊 Added table: {num_rows} rows × {num_cols} cols")
    return top + table_height + Inches(0.3)


def add_gradient_background(slide):
    """Add a full-slide gradient background."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.line.fill.background()
    
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = COLORS['header_dark']
    fill.gradient_stops[1].color.rgb = COLORS['header_light']


def add_header_bar(slide, title_text, subtitle_text=""):
    """
    Add a header region with yellow accent bar on left, dark title, and
    a thin divider line.  Matches the corporate template (no gradient).
    """
    # Determine heights
    is_long_title = len(title_text) > 50
    has_subtitle = bool(subtitle_text)
    base_h = 1.0
    if has_subtitle:
        base_h = 1.5
    if is_long_title:
        base_h = max(base_h, 1.4)

    bar_height = Inches(base_h)

    # Yellow accent bar on left (stops just above the gray divider)
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(0.22),
        Inches(0.08), Inches(base_h - 0.30)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS['bullet_marker']
    accent_bar.line.fill.background()

    # Title text (dark blue)
    title_top = Inches(0.18)
    title_height = Inches(0.9) if not is_long_title else Inches(1.2)

    title_box = slide.shapes.add_textbox(
        Inches(0.65), title_top, Inches(12), title_height
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    p.font.name = FONTS['title']

    # Subtitle (secondary colour)
    if subtitle_text:
        subtitle_top = Inches(0.9) if not is_long_title else Inches(1.1)
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.65), subtitle_top, Inches(12), Inches(0.5)
        )
        tf_sub = subtitle_box.text_frame
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = COLORS['secondary']
        p_sub.font.name = FONTS['body']

    # Divider line under header
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), bar_height,
        Inches(12.533), Inches(0.02)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(208, 208, 208)
    divider.line.fill.background()

    return bar_height



def add_bullets(slide, bullets_elem, top, left=None, width=None):
    """Add bullet list with styled markers."""
    if left is None:
        left = CONTENT_LEFT
    if width is None:
        width = CONTENT_WIDTH
        
    items = bullets_elem.find_all('li', recursive=False)
    
    if not items:
        return top
    
    text_box = slide.shapes.add_textbox(left, top, width, CONTENT_HEIGHT)
    tf = text_box.text_frame
    tf.word_wrap = True
    
    for i, li in enumerate(items):
        text = ""
        for child in li.children:
            if isinstance(child, str):
                text += child.strip()
            elif child.name != 'ul':
                text += child.get_text(strip=True)
        
        if not text:
            continue
            
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Split Bullet Symbol (Yellow) and Text (Black)
        run_symbol = p.add_run()
        run_symbol.text = "▸ "
        run_symbol.font.size = Pt(20)
        run_symbol.font.name = FONTS['body']
        run_symbol.font.color.rgb = COLORS['bullet_marker'] # Yellow Accent
        
        run_text = p.add_run()
        run_text.text = text
        run_text.font.size = Pt(20)
        run_text.font.name = FONTS['body']
        run_text.font.color.rgb = COLORS['text_black'] # Black Text
        
        p.space_before = Pt(8)
        p.space_after = Pt(4)
        
        nested_ul = li.find('ul')
        if nested_ul:
            for nested_li in nested_ul.find_all('li', recursive=False):
                nested_text = nested_li.get_text(strip=True)
                if nested_text:
                    np = tf.add_paragraph()
                    run_symbol_nested = np.add_run()
                    run_symbol_nested.text = "    ○ "
                    run_symbol_nested.font.size = Pt(18)
                    run_symbol_nested.font.name = FONTS['body']
                    run_symbol_nested.font.color.rgb = COLORS['bullet_marker']

                    run_text_nested = np.add_run()
                    run_text_nested.text = nested_text
                    run_text_nested.font.size = Pt(18)
                    run_text_nested.font.name = FONTS['body']
                    run_text_nested.font.color.rgb = COLORS['text_black']
                    np.space_before = Pt(4)
    
    return top + Inches(len(items) * 0.5)


def _parse_colored_spans(html_elem):
    """
    Parse HTML code element preserving <span style="color: #xxx"> tokens.
    Returns list of (text, color_hex_or_None) tuples per line,
    grouped as list-of-lists (one inner list per source line).
    """
    from bs4 import NavigableString, Tag

    # Find the innermost code/pre element that contains the highlighted spans
    code_tag = html_elem.find('code') or html_elem.find('pre') or html_elem
    if code_tag.find('code'):
        code_tag = code_tag.find('code')

    lines = [[]]  # list of lists of (text, color)

    def _walk(node):
        if isinstance(node, NavigableString):
            text = str(node)
            # Split by newlines to maintain line structure
            parts = text.split('\n')
            for i, part in enumerate(parts):
                if i > 0:
                    lines.append([])  # new line
                if part:
                    lines[-1].append((part, None))
        elif isinstance(node, Tag):
            # Extract color from inline style
            color = None
            style = node.get('style', '')
            if style:
                m = re.search(r'color:\s*#([0-9a-fA-F]{3,8})', style)
                if m:
                    color = m.group(1)
                    # Normalize 3-char hex to 6-char
                    if len(color) == 3:
                        color = ''.join(c * 2 for c in color)

            for child in node.children:
                if isinstance(child, NavigableString):
                    text = str(child)
                    parts = text.split('\n')
                    for i, part in enumerate(parts):
                        if i > 0:
                            lines.append([])
                        if part:
                            lines[-1].append((part, color))
                elif isinstance(child, Tag):
                    # Recurse — inner tag may override color
                    _walk(child)

    _walk(code_tag)
    return lines


def add_code_block(slide, code_elem, top):
    """Add a dark-themed code block with syntax coloring from Pygments spans."""
    code_content = code_elem.find('pre') or code_elem.find(class_='code-content')

    # Parse colored spans from the HTML
    has_spans = code_content and code_content.find('span')
    if has_spans:
        colored_lines = _parse_colored_spans(code_content)
    else:
        plain = code_content.get_text() if code_content else ""
        colored_lines = [[(seg, None)] for seg in plain.strip().split('\n')]

    # Limit to 20 lines
    colored_lines = colored_lines[:20]
    num_lines = len(colored_lines)

    # Calculate dynamic height
    line_height = 0.25
    padding = 0.6
    calculated_height = (num_lines * line_height) + padding
    box_height = max(1.5, min(4.5, calculated_height))

    # Dark background rounded rectangle
    code_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        CONTENT_LEFT, top,
        CONTENT_WIDTH, Inches(box_height)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = COLORS['code_bg']
    code_box.line.fill.background()

    # Text box with colored runs
    text_box = slide.shapes.add_textbox(
        CONTENT_LEFT + Inches(0.2), top + Inches(0.2),
        CONTENT_WIDTH - Inches(0.4), Inches(box_height - 0.4)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    default_color = COLORS['code_text']  # #D4D4D4

    for line_idx, segments in enumerate(colored_lines):
        if line_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        if not segments:
            # Empty line — add a single space run to preserve line break
            run = p.add_run()
            run.text = " "
            run.font.size = Pt(12)
            run.font.name = FONTS['code']
            run.font.color.rgb = default_color
            continue

        for text, color_hex in segments:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(12)
            run.font.name = FONTS['code']
            if color_hex:
                try:
                    r = int(color_hex[0:2], 16)
                    g = int(color_hex[2:4], 16)
                    b = int(color_hex[4:6], 16)
                    run.font.color.rgb = RGBColor(r, g, b)
                except (ValueError, IndexError):
                    run.font.color.rgb = default_color
            else:
                run.font.color.rgb = default_color

    logger.info(f"   🎨 Added code block with syntax coloring ({num_lines} lines)")


def add_callout(slide, callout_elem):
    """Add a yellow callout box."""
    text = callout_elem.get_text(strip=True)
    
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        CONTENT_LEFT, Inches(5.5),
        CONTENT_WIDTH, Inches(1)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = COLORS['callout_bg']
    callout.line.fill.background()
    
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.name = FONTS['body']
    p.font.color.rgb = COLORS['text_dark']
    p.alignment = PP_ALIGN.CENTER


def add_content_heading(slide, heading_elem, top, left=None, width=None):
    """Add a bold styled heading (e.g. 'Descripción')."""
    if left is None:
        left = CONTENT_LEFT
    if width is None:
        width = CONTENT_WIDTH
        
    text = heading_elem.get_text(strip=True)
    if not text:
        return top
        
    # Text box for heading
    text_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = FONTS['title']
    p.font.color.rgb = COLORS['primary'] # Use primary blue for headings
    
    return top + Inches(0.6) # Return new top position


def add_logo_middle_center(slide, logo_bytes):
    """Add logo centered both vertically and horizontally."""
    if not logo_bytes:
        return
    
    try:
        logo_stream = io.BytesIO(logo_bytes)
        logo_w = Inches(3) # Larger logo for module end
        logo_h = Inches(1.2)
        
        left = (SLIDE_WIDTH - logo_w) / 2
        top = (SLIDE_HEIGHT - logo_h) / 2
        
        slide.shapes.add_picture(logo_stream, left, top, logo_w, logo_h)
    except Exception as e:
        logger.warning(f"Could not add centered logo: {e}")

