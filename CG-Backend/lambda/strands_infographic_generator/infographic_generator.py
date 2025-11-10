"""
HTML Infographic Generator - Clean, Modern Approach
====================================================
Generates beautiful HTML infographics from course content and exports to editable PPT.

Key Features:
- Clean HTML/CSS layout (no overlaps!)
- Semantic structure that converts well to PPT
- Each concept = one slide with proper text boxes
- Images positioned cleanly
- Export to editable PowerPoint format

Expected event parameters:
    - course_bucket: S3 bucket name
    - project_folder: Project folder path
    - book_version_key: S3 key to specific book version JSON
    - model_provider: 'bedrock' or 'openai'
    - slides_per_lesson: Number of infographic slides per lesson (default: 5)
    - style: 'professional' | 'modern' | 'minimal' (default: 'professional')
"""

import os
import logging
import json
import boto3
import re
from datetime import datetime
from typing import Dict, List, Any, Optional
from botocore.config import Config

# Configure boto3 with extended timeouts
boto_config = Config(
    read_timeout=900,
    connect_timeout=60,
    retries={'max_attempts': 3, 'mode': 'adaptive'}
)

# AWS Clients
s3_client = boto3.client('s3')
bedrock_client = boto3.client('bedrock-runtime', region_name='us-east-1', config=boto_config)
secrets_client = boto3.client('secretsmanager', region_name='us-east-1')

# Model Configuration
DEFAULT_BEDROCK_MODEL = "us.anthropic.claude-sonnet-4-5-20250929-v1:0"
DEFAULT_OPENAI_MODEL = "gpt-5"

# Logging
logger = logging.getLogger("aurora.infographic_generator")
if not logger.handlers:
    h = logging.StreamHandler()
    fmt = logging.Formatter("%(asctime)s %(levelname)s %(name)s: %(message)s")
    h.setFormatter(fmt)
    logger.addHandler(h)
    logger.setLevel(logging.INFO)


def get_secret(secret_name: str) -> dict:
    """Retrieve secret from AWS Secrets Manager."""
    try:
        response = secrets_client.get_secret_value(SecretId=secret_name)
        return json.loads(response['SecretString'])
    except Exception as e:
        logger.error(f"Error retrieving secret {secret_name}: {e}")
        raise


def configure_bedrock_model(model_id: str = DEFAULT_BEDROCK_MODEL):
    """Configure Bedrock model for Strands Agent."""
    try:
        from strands.models import BedrockModel
        
        model = BedrockModel(
            model_id=model_id,
            client=bedrock_client
        )
        logger.info(f"✅ Configured Bedrock model: {model_id}")
        return model
    except Exception as e:
        logger.error(f"Failed to configure Bedrock model: {e}")
        raise


def configure_openai_model(model_id: str = DEFAULT_OPENAI_MODEL) -> Any:
    """Configure OpenAI model for Strands Agent."""
    try:
        from strands.models.openai import OpenAIModel
        
        secret_data = get_secret('aurora/openai-api-key')
        api_key = secret_data.get('api_key')
        
        if not api_key:
            raise ValueError("OpenAI API key not found in secrets")
        
        model = OpenAIModel(
            client_args={"api_key": api_key},
            model_id=model_id,
            streaming=False
        )
        logger.info(f"✅ Configured OpenAI model: {model_id}")
        return model
    except Exception as e:
        logger.error(f"Failed to configure OpenAI model: {e}")
        raise


def load_book_from_s3(bucket: str, book_key: str) -> Dict:
    """Load book JSON from S3."""
    try:
        logger.info(f"📖 Loading book from s3://{bucket}/{book_key}")
        response = s3_client.get_object(Bucket=bucket, Key=book_key)
        book_data = json.loads(response['Body'].read().decode('utf-8'))
        
        # Handle nested module structure
        lessons = []
        if 'modules' in book_data:
            # New structure: modules contain lessons
            for module in book_data.get('modules', []):
                lessons.extend(module.get('lessons', []))
            logger.info(f"✅ Loaded book with {len(lessons)} lessons from {len(book_data.get('modules', []))} modules")
            # Flatten structure for compatibility
            book_data['lessons'] = lessons
        else:
            # Old structure: lessons at top level
            lessons = book_data.get('lessons', [])
            logger.info(f"✅ Loaded book with {len(lessons)} lessons")
        
        # Try to load outline for richer module information
        try:
            project_folder = book_key.split('/')[0]
            outline_key = f"{project_folder}/outline/mini_cisco.yaml"
            logger.info(f"📋 Attempting to load outline: {outline_key}")
            outline_response = s3_client.get_object(Bucket=bucket, Key=outline_key)
            import yaml
            outline_data = yaml.safe_load(outline_response['Body'].read().decode('utf-8'))
            
            # Extract module info from outline
            if 'course' in outline_data and 'modules' in outline_data['course']:
                book_data['outline_modules'] = outline_data['course']['modules']
                logger.info(f"✅ Loaded outline with {len(outline_data['course']['modules'])} modules")
            
            # Extract course metadata
            if 'course' in outline_data:
                course = outline_data['course']
                book_data['course_metadata'] = {
                    'title': course.get('title', ''),
                    'description': course.get('description', ''),
                    'audience': course.get('audience', []),
                    'prerequisites': course.get('prerequisites', []),
                    'learning_outcomes': course.get('learning_outcomes', []),
                    'level': course.get('level', ''),
                    'duration': course.get('total_duration_minutes', 0)
                }
                logger.info(f"✅ Extracted course metadata: {course.get('title', 'N/A')}")
        except Exception as e:
            logger.warning(f"⚠️ Could not load outline: {e}")
        
        return book_data
    except Exception as e:
        logger.error(f"Failed to load book: {e}")
        raise


def extract_images_from_content(content: str) -> List[Dict[str, str]]:
    """Extract image references from lesson content."""
    images = []
    
    # Match markdown image syntax: ![alt](url)
    img_pattern = r'!\[([^\]]*)\]\(([^\)]+)\)'
    matches = re.findall(img_pattern, content)
    
    for alt_text, img_url in matches:
        images.append({
            'alt_text': alt_text,
            'url': img_url,
            'type': 'diagram' if 'diagram' in alt_text.lower() else 'image'
        })
    
    return images


def generate_infographic_structure(
    book_data: Dict,
    model,
    slides_per_lesson: int = 5,
    style: str = 'professional'
) -> Dict:
    """
    Generate infographic structure using AI.
    Each slide = clean HTML section with proper semantic structure.
    """
    from strands import Agent
    
    metadata = book_data.get('metadata', {})
    course_title = metadata.get('title', 'Course Infographic')
    lessons = book_data.get('lessons', [])
    
    logger.info(f"\n🎨 Generating infographic structure for: {course_title}")
    logger.info(f"📊 Lessons: {len(lessons)}, Slides per lesson: {slides_per_lesson}")
    logger.info(f"✨ Style: {style}")
    
    # Create AI agent optimized for infographic generation
    infographic_designer = Agent(
        model=model,
        system_prompt=f"""You are an expert infographic designer specializing in educational content.
Your goal is to create clean, visually appealing slide structures that will be rendered as HTML and exported to PowerPoint.

CRITICAL DESIGN PRINCIPLES:
1. **No Overlaps**: Each slide has clearly separated sections (title, content, visual, summary)
2. **Hierarchy**: Use heading levels (h1, h2, h3) to create visual hierarchy
3. **Whitespace**: Generous spacing between elements - never cram content
4. **Concise Text**: 3-5 bullet points maximum, each under 15 words
5. **One Focus**: Each slide focuses on ONE key concept

SLIDE STRUCTURE (will become HTML):
- title: Clear, descriptive (not generic)
- subtitle: Optional context (1 line)
- content_blocks: Array of content sections
  - type: 'text', 'bullets', 'image', 'quote', 'callout'
  - heading: Section heading
  - items: Array of text items or single text/image
- layout_hint: 'single-column', 'two-column', 'image-focus', 'text-focus'

Style: {style}
- professional: Clean, corporate, blues and grays, structured
- modern: Bold colors, minimal text, high-impact visuals  
- minimal: Maximum whitespace, serif fonts, elegant simplicity

OUTPUT FORMAT (JSON):
{{
    "course_title": "{course_title}",
    "total_slides": 25,
    "style": "{style}",
    "slides": [
        {{
            "slide_number": 1,
            "title": "Compelling Title - Not Generic",
            "subtitle": "Optional one-line context",
            "layout_hint": "single-column",
            "content_blocks": [
                {{
                    "type": "bullets",
                    "heading": "Key Concepts",
                    "items": [
                        "Concise point 1 - under 15 words",
                        "Concise point 2 - under 15 words",
                        "Concise point 3 - under 15 words"
                    ]
                }},
                {{
                    "type": "image",
                    "image_reference": "USE_IMAGE: specific image description",
                    "caption": "Clear caption explaining the visual"
                }}
            ],
            "notes": "Instructor notes"
        }}
    ]
}}

REMEMBER:
- Each slide = ONE key idea
- Maximum 5 bullets per slide
- Use images strategically (not on every slide)
- Generous whitespace is professional
- Short, punchy text beats long paragraphs
""",
        tools=[]
    )
    
    all_slides = []
    slide_counter = 1
    
    # Detect language from course content (simple heuristic)
    sample_text = ' '.join([l.get('title', '') for l in lessons[:3]])
    is_spanish = any(word in sample_text.lower() for word in ['introducción', 'conceptos', 'básicos', 'lección'])
    
    # Title slide with actual course title
    all_slides.append({
        "slide_number": slide_counter,
        "title": course_title,
        "subtitle": f"{len(lessons)} {'Lecciones' if is_spanish else 'Lessons'}",
        "layout_hint": "title",
        "content_blocks": [],
        "notes": "Welcome and course overview"
    })
    slide_counter += 1
    
    # Module overview slide
    if 'modules' in book_data:
        module_titles = []
        for module in book_data.get('modules', []):
            module_title = module.get('module_title', f"Module {module.get('module_number', '')}")
            module_lessons = module.get('lessons', [])
            module_titles.append(f"{module_title} ({len(module_lessons)} {'lecciones' if is_spanish else 'lessons'})")
        
        all_slides.append({
            "slide_number": slide_counter,
            "title": "Módulos del Curso" if is_spanish else "Course Modules",
            "subtitle": "",
            "layout_hint": "single-column",
            "content_blocks": [
                {
                    "type": "bullets",
                    "heading": "",
                    "items": module_titles[:10]  # Limit to 10 modules
                }
            ],
            "notes": "Overview of course structure"
        })
        slide_counter += 1
    
    # Process lessons with timeout guard
    import time as time_module
    start_time = time_module.time()
    MAX_PROCESSING_TIME = 840  # 14 minutes
    
    # Get all lab activity titles to avoid duplicating them
    lab_titles = []
    if 'outline_modules' in book_data:
        for module in book_data.get('outline_modules', []):
            for lab in module.get('lab_activities', []):
                lab_titles.append(lab.get('title', '').lower())
    
    lessons_processed = 0
    for lesson_idx, lesson in enumerate(lessons, 1):
        # Check timeout
        elapsed_time = time_module.time() - start_time
        if elapsed_time > MAX_PROCESSING_TIME:
            logger.warning(f"⚠️ Approaching timeout - processed {lessons_processed}/{len(lessons)} lessons")
            break
        
        lesson_title = lesson.get('title', f'Lesson {lesson_idx}')
        lesson_content = lesson.get('content', '')
        
        # Skip lessons that are lab activities (they'll be added separately from outline)
        is_lab_lesson = (
            'laboratorio' in lesson_title.lower() or
            'lab' in lesson_title.lower() or
            lesson_title.lower() in lab_titles
        )
        
        if is_lab_lesson:
            logger.info(f"\n⏭️  Skipping Lab Lesson {lesson_idx}/{len(lessons)}: {lesson_title} (will be added from outline)")
            continue
        
        logger.info(f"\n📝 Processing Lesson {lesson_idx}/{len(lessons)}: {lesson_title}")
        
        # Extract available images
        available_images = extract_images_from_content(lesson_content)
        logger.info(f"🖼️  Found {len(available_images)} images in lesson content")
        if available_images:
            for img in available_images[:5]:  # Log first 5 images
                logger.info(f"   - Image: {img.get('image_reference', 'N/A')} | Alt: {img.get('alt_text', 'N/A')}")
        
        # Generate slides for this lesson
        # Determine if we should use all content (no limit on slides)
        use_all_content = slides_per_lesson >= 50  # High number indicates "use all content" mode
        
        if use_all_content:
            slide_count_instruction = "as many slides as needed to cover ALL the content comprehensively (no limit)"
            content_coverage_instruction = """- CRITICAL: Include EVERY section, subsection, and important detail from the content
- Do NOT skip or summarize any content - every concept must have its own slide  
- PRESERVE the content text AS-IS - do not rewrite or rephrase bullet points
- Use the EXACT images referenced in the markdown (by their image IDs/descriptions)
- Maintain the ORIGINAL structure and organization of the content
- Better to have more slides with clear, focused content than fewer slides with cramped information"""
        else:
            slide_count_instruction = f"{slides_per_lesson} slides"
            content_coverage_instruction = "- Each slide = ONE key concept\n- Focus on the most important topics"
        
        lesson_prompt = f"""
Create {slide_count_instruction} for this lesson.

LESSON {lesson_idx}: {lesson_title}

CONTENT:
{lesson_content}

Available images: {len(available_images)}
Image descriptions: {[img['alt_text'] for img in available_images]}

REQUIREMENTS:
{content_coverage_instruction}
- 3-5 bullets maximum per slide
- Use images strategically (reference by description)
- Clear hierarchy with headings
- Generous whitespace (don't cram content)
- Layout hint for each slide
- IMPORTANT: The LAST slide must be a lesson summary with key takeaways
- Language: Use the same language as the lesson content (Spanish/English)
"""
        
        # Call AI with retry logic
        max_retries = 3
        retry_delay = 10
        ai_response = None
        
        for attempt in range(max_retries):
            try:
                logger.info(f"🔄 Attempt {attempt + 1}/{max_retries}...")
                ai_response = infographic_designer(lesson_prompt)
                logger.info(f"✅ Generated slides for lesson {lesson_idx}")
                break
            except Exception as e:
                if "timed out" in str(e).lower() and attempt < max_retries - 1:
                    logger.warning(f"⚠️ Timeout, retrying in {retry_delay}s...")
                    time_module.sleep(retry_delay)
                    retry_delay *= 2
                else:
                    raise
        
        if ai_response is None:
            raise Exception(f"Failed to generate slides for lesson {lesson_idx}")
        
        # Parse AI response
        if hasattr(ai_response, 'output'):
            ai_response = ai_response.output
        elif hasattr(ai_response, 'text'):
            ai_response = ai_response.text
        
        ai_response = str(ai_response).strip()
        if ai_response.startswith('```json'):
            ai_response = ai_response[7:]
        elif ai_response.startswith('```'):
            ai_response = ai_response[3:]
        if ai_response.endswith('```'):
            ai_response = ai_response[:-3]
        
        try:
            lesson_slides = json.loads(ai_response.strip())
        except json.JSONDecodeError as e:
            logger.error(f"JSON parse error for lesson {lesson_idx}: {e}")
            # Save for debugging
            with open(f"/tmp/failed_json_lesson_{lesson_idx}.txt", "w") as f:
                f.write(ai_response)
            raise
        
        # Add slides from this lesson
        current_module = lesson.get('module_number', 1)
        for slide_data in lesson_slides.get('slides', []):
            slide_data['slide_number'] = slide_counter
            slide_data['lesson_number'] = lesson_idx
            slide_data['lesson_title'] = lesson_title
            slide_data['module_number'] = current_module
            all_slides.append(slide_data)
            slide_counter += 1
        
        # NOTE: Module summary, lab slides, and "Gracias" slide will be added later
        # when processing slides in create_pptx_from_structure()
        
        lessons_processed += 1
        logger.info(f"✅ Completed lesson {lesson_idx}/{len(lessons)} - Total slides: {len(all_slides)}")
    
    # NOTE: No course summary slide is added here
    # Each module ends with: Last Lesson Summary → Labs → Gracias
    # This provides natural closure for each module
    
    completion_status = "complete" if lessons_processed == len(lessons) else "partial"
    
    infographic_structure = {
        "course_title": course_title,
        "total_slides": len(all_slides),
        "total_lessons": len(lessons),
        "lessons_processed": lessons_processed,
        "completion_status": completion_status,
        "style": style,
        "generated_at": datetime.now().isoformat(),
        "slides": all_slides,
        "outline_modules": book_data.get('outline_modules', []),  # Pass outline info
        "course_metadata": book_data.get('course_metadata', {})  # Pass course metadata
    }
    
    logger.info(f"\n✅ Generated {len(all_slides)} infographic slides (no course summary)")
    logger.info(f"📋 Structure: Lessons with summaries + Labs/Gracias per module")
    return infographic_structure
    return infographic_structure


def generate_html_from_structure(structure: Dict) -> str:
    """
    Convert infographic structure to clean HTML.
    Each slide = semantic HTML section.
    """
    style = structure.get('style', 'professional')
    
    # Color schemes
    color_schemes = {
        'professional': {
            'primary': '#003366',
            'secondary': '#4682B4',
            'accent': '#FFC000',
            'bg': '#F8F9FA',
            'text': '#333333'
        },
        'modern': {
            'primary': '#212121',
            'secondary': '#607D8B',
            'accent': '#009688',
            'bg': '#FAFAFA',
            'text': '#212121'
        },
        'minimal': {
            'primary': '#000000',
            'secondary': '#757575',
            'accent': '#E0E0E0',
            'bg': '#FFFFFF',
            'text': '#212121'
        }
    }
    
    colors = color_schemes.get(style, color_schemes['professional'])
    
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{structure.get('course_title', 'Course Infographic')}</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            color: {colors['text']};
            background: {colors['bg']};
            line-height: 1.6;
        }}
        
        /* Toolbar for editing controls */
        .toolbar {{
            position: fixed;
            top: 0;
            left: 0;
            right: 0;
            background: #2c3e50;
            padding: 15px 20px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.2);
            z-index: 1000;
            display: flex;
            justify-content: space-between;
            align-items: center;
        }}
        
        .toolbar-title {{
            color: white;
            font-size: 18px;
            font-weight: 600;
        }}
        
        .toolbar-buttons {{
            display: flex;
            gap: 10px;
        }}
        
        .btn {{
            padding: 10px 20px;
            border: none;
            border-radius: 5px;
            font-size: 14px;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s;
        }}
        
        .btn-primary {{
            background: {colors['accent']};
            color: #333;
        }}
        
        .btn-primary:hover {{
            background: {colors['primary']};
            color: white;
        }}
        
        .btn-secondary {{
            background: #34495e;
            color: white;
        }}
        
        .btn-secondary:hover {{
            background: #2c3e50;
        }}
        
        .btn-success {{
            background: #27ae60;
            color: white;
        }}
        
        .btn-success:hover {{
            background: #229954;
        }}
        
        /* Content area with top margin for toolbar */
        .content {{
            margin-top: 80px;
        }}
        
        .slide {{
            width: 1280px;
            height: 720px;
            margin: 20px auto;
            background: white;
            padding: 60px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            page-break-after: always;
            display: flex;
            flex-direction: column;
            position: relative;
        }}
        
        /* Editable elements styling */
        [contenteditable="true"] {{
            outline: none;
            transition: background 0.3s;
        }}
        
        [contenteditable="true"]:hover {{
            background: rgba(255, 255, 0, 0.1);
            cursor: text;
        }}
        
        [contenteditable="true"]:focus {{
            background: rgba(255, 255, 0, 0.2);
            padding: 5px;
            margin: -5px;
            border-radius: 3px;
        }}
        
        .slide-title {{
            font-size: 48px;
            font-weight: bold;
            color: {colors['primary']};
            margin-bottom: 10px;
        }}
        
        .slide-subtitle {{
            font-size: 24px;
            color: {colors['secondary']};
            margin-bottom: 40px;
        }}
        
        .content-block {{
            margin-bottom: 30px;
        }}
        
        .block-heading {{
            font-size: 32px;
            font-weight: 600;
            color: {colors['primary']};
            margin-bottom: 15px;
        }}
        
        .bullets {{
            list-style: none;
            padding: 0;
        }}
        
        .bullets li {{
            font-size: 24px;
            padding: 12px 0 12px 40px;
            position: relative;
        }}
        
        .bullets li:before {{
            content: "▸";
            color: {colors['accent']};
            font-size: 28px;
            position: absolute;
            left: 0;
        }}
        
        .image-container {{
            text-align: center;
            margin: 20px 0;
        }}
        
        .image-container img {{
            max-width: 100%;
            max-height: 400px;
            object-fit: contain;
        }}
        
        .image-caption {{
            font-size: 18px;
            color: {colors['secondary']};
            margin-top: 10px;
            font-style: italic;
        }}
        
        .callout {{
            background: {colors['accent']}20;
            border-left: 4px solid {colors['accent']};
            padding: 20px;
            font-size: 22px;
        }}
        
        .two-column {{
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 40px;
        }}
        
        @media print {{
            .toolbar {{
                display: none;
            }}
            
            .content {{
                margin-top: 0;
            }}
            
            .slide {{
                page-break-after: always;
                margin: 0;
                box-shadow: none;
                width: 100%;
                height: auto;
            }}
            
            [contenteditable="true"]:hover,
            [contenteditable="true"]:focus {{
                background: transparent !important;
                padding: 0 !important;
                margin: 0 !important;
            }}
        }}
    </style>
</head>
<body>
    <!-- Editing Toolbar -->
    <div class="toolbar">
        <div class="toolbar-title">
            📊 {structure.get('course_title', 'Course Infographic')} - Editable Version
        </div>
        <div class="toolbar-buttons">
            <button class="btn btn-secondary" onclick="toggleEditMode()">
                <span id="edit-btn-text">🔒 Lock Editing</span>
            </button>
            <button class="btn btn-success" onclick="saveChanges()">
                💾 Save Changes
            </button>
            <button class="btn btn-primary" onclick="downloadPDF()">
                📄 Download PDF
            </button>
        </div>
    </div>
    
    <div class="content">
"""
    
    # Generate slides with editable content
    for slide in structure.get('slides', []):
        layout = slide.get('layout_hint', 'single-column')
        
        html += f'<div class="slide" data-slide="{slide.get("slide_number")}">\n'
        html += f'  <h1 class="slide-title" contenteditable="true">{slide.get("title", "")}</h1>\n'
        
        if slide.get('subtitle'):
            html += f'  <p class="slide-subtitle" contenteditable="true">{slide["subtitle"]}</p>\n'
        
        # Content blocks
        content_blocks = slide.get('content_blocks', [])
        
        if layout == 'two-column' and len(content_blocks) >= 2:
            html += '  <div class="two-column">\n'
            for block in content_blocks[:2]:
                html += '    <div>\n'
                html += generate_content_block_html(block)
                html += '    </div>\n'
            html += '  </div>\n'
            # Remaining blocks
            for block in content_blocks[2:]:
                html += generate_content_block_html(block)
        else:
            for block in content_blocks:
                html += generate_content_block_html(block)
        
        html += '</div>\n\n'
    
    html += """
    </div> <!-- End content -->
    
    <script>
        let editingEnabled = true;
        
        // Toggle edit mode
        function toggleEditMode() {
            editingEnabled = !editingEnabled;
            const editables = document.querySelectorAll('[contenteditable]');
            const btnText = document.getElementById('edit-btn-text');
            
            editables.forEach(el => {
                el.contentEditable = editingEnabled;
            });
            
            if (editingEnabled) {
                btnText.textContent = '🔒 Lock Editing';
            } else {
                btnText.textContent = '🔓 Enable Editing';
            }
        }
        
        // Save changes (downloads edited HTML)
        function saveChanges() {
            const html = document.documentElement.outerHTML;
            const blob = new Blob([html], { type: 'text/html' });
            const url = URL.createObjectURL(blob);
            const a = document.createElement('a');
            a.href = url;
            a.download = 'infographic_edited.html';
            document.body.appendChild(a);
            a.click();
            document.body.removeChild(a);
            URL.revokeObjectURL(url);
            
            alert('✅ Changes saved! The edited HTML file has been downloaded.');
        }
        
        // Download as PDF
        function downloadPDF() {
            alert('📄 PDF Download Instructions:\\n\\n1. Press Ctrl+P (Windows/Linux) or Cmd+P (Mac)\\n2. Select "Save as PDF" as the destination\\n3. Click Save\\n\\nThe PDF will preserve all your edits and formatting!');
            window.print();
        }
        
        // Auto-save to localStorage every 30 seconds
        setInterval(() => {
            if (editingEnabled) {
                const content = document.querySelector('.content').innerHTML;
                localStorage.setItem('infographic_autosave', content);
                console.log('✅ Auto-saved at', new Date().toLocaleTimeString());
            }
        }, 30000);
        
        // Restore from localStorage on load
        window.addEventListener('load', () => {
            const saved = localStorage.getItem('infographic_autosave');
            if (saved && confirm('📝 Found auto-saved changes. Restore them?')) {
                document.querySelector('.content').innerHTML = saved;
            }
        });
        
        // Warn before leaving if there are unsaved changes
        let initialContent = document.querySelector('.content').innerHTML;
        window.addEventListener('beforeunload', (e) => {
            const currentContent = document.querySelector('.content').innerHTML;
            if (currentContent !== initialContent && editingEnabled) {
                e.preventDefault();
                e.returnValue = 'You have unsaved changes. Are you sure you want to leave?';
            }
        });
    </script>
</body>
</html>
"""
    
    return html


def generate_content_block_html(block: Dict) -> str:
    """Generate HTML for a content block with editable content."""
    block_type = block.get('type', 'text')
    html = '<div class="content-block">\n'
    
    if block.get('heading'):
        html += f'  <h2 class="block-heading" contenteditable="true">{block["heading"]}</h2>\n'
    
    if block_type == 'bullets':
        html += '  <ul class="bullets">\n'
        for item in block.get('items', []):
            html += f'    <li contenteditable="true">{item}</li>\n'
        html += '  </ul>\n'
    
    elif block_type == 'image':
        html += '  <div class="image-container">\n'
        # Note: In full implementation, resolve image URLs from S3
        image_ref = block.get('image_reference', '')
        html += f'    <div class="image-placeholder" style="background:#f0f0f0;height:300px;display:flex;align-items:center;justify-content:center;">\n'
        html += f'      <span style="color:#999;">{image_ref}</span>\n'
        html += '    </div>\n'
        if block.get('caption'):
            html += f'    <p class="image-caption" contenteditable="true">{block["caption"]}</p>\n'
        html += '  </div>\n'
    
    elif block_type == 'callout':
        text = block.get('text', '')
        html += f'  <div class="callout" contenteditable="true">{text}</div>\n'
    
    elif block_type == 'text':
        text = block.get('text', '')
        html += f'  <p style="font-size:24px;line-height:1.6;" contenteditable="true">{text}</p>\n'
    
    html += '</div>\n'
    return html


def reorganize_slides_for_correct_order(structure: Dict) -> Dict:
    """
    Reorganize slides to ensure correct order:
    - Remove standalone "Gracias" slides (they'll be added after labs)
    - Remove lab lesson slides (they're handled separately from outline)
    - Remove course summary slides (not needed - each module has its own ending)
    - Group by module and ensure: lessons → summary → labs → gracias
    """
    slides = structure.get('slides', [])
    outline_modules = structure.get('outline_modules', [])
    
    # Get all lab activity titles
    lab_titles_lower = []
    if outline_modules:
        for module in outline_modules:
            for lab in module.get('lab_activities', []):
                lab_titles_lower.append(lab.get('title', '').lower())
    
    reorganized = []
    removed_count = 0
    
    for slide in slides:
        layout_hint = slide.get('layout_hint', '')
        title = slide.get('title', '').lower()
        
        # Skip standalone "Gracias" slides - they'll be added after labs
        if layout_hint == 'thanks' or 'gracias' in title or 'thank you' in title:
            logger.info(f"⏭️  Removing standalone 'Gracias' slide: {slide.get('title')}")
            removed_count += 1
            continue
        
        # Skip course summary slides - not needed
        if 'resumen del curso' in title or 'course summary' in title:
            logger.info(f"⏭️  Removing course summary slide: {slide.get('title')}")
            removed_count += 1
            continue
        
        # Skip lab lesson slides - they'll be added from outline
        if 'laboratorio' in title or 'lab' in title or title in lab_titles_lower:
            logger.info(f"⏭️  Removing lab lesson slide: {slide.get('title')}")
            removed_count += 1
            continue
        
        reorganized.append(slide)
    
    logger.info(f"📋 Reorganized slides: {len(slides)} → {len(reorganized)} (removed {removed_count} slides)")
    structure['slides'] = reorganized
    return structure


def convert_html_to_pptx(html_content: str, structure: Dict, course_bucket: str = None, project_folder: str = None, template_key: str = None) -> bytes:
    """
    Convert HTML infographic to editable PowerPoint.
    Each HTML slide becomes a PPT slide with text boxes and images from S3.
    Optionally uses a custom template from S3.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from PIL import Image
    import io
    
    logger.info("🎨 Converting HTML to editable PowerPoint with images...")
    
    # FIRST: Reorganize slides to remove duplicates and ensure correct order
    structure = reorganize_slides_for_correct_order(structure)
    
    # Try to load template from S3
    use_template = False
    if template_key:
        try:
            logger.info(f"📥 Loading template from S3: {template_key}")
            template_response = s3_client.get_object(Bucket='crewai-course-artifacts', Key=template_key)
            template_stream = io.BytesIO(template_response['Body'].read())
            prs = Presentation(template_stream)
            use_template = True
            logger.info(f"✅ Template loaded! {len(prs.slide_layouts)} layouts available")
        except Exception as e:
            logger.warning(f"⚠️ Failed to load template: {e}. Using default design.")
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)
    
    # Get layouts (template or default)
    if use_template:
        # Template layout mapping - Full set
        title_slide_layout = prs.slide_layouts[0]  # Portada del curso
        property_layout = prs.slide_layouts[1]  # Propiedad intelectual
        description_layout = prs.slide_layouts[2]  # Descripción del curso
        objectives_layout = prs.slide_layouts[3]  # Objetivos del curso
        prerequisites_layout = prs.slide_layouts[4]  # Prerrequisitos del curso
        audience_layout = prs.slide_layouts[5]  # Audiencia del curso
        outline_layout = prs.slide_layouts[6]  # Temario del curso
        presentation_layout = prs.slide_layouts[7]  # Presentación del grupo
        module_layout = prs.slide_layouts[8]  # Portada del capítulo
        lesson_layout = prs.slide_layouts[9]  # Nombre del tema 2
        content_layout = prs.slide_layouts[10]  # Contenido - General
        summary_layout = prs.slide_layouts[11]  # Resumen del capítulo
        lab_layout = prs.slide_layouts[12]  # Descripción de la práctica
        thanks_layout = prs.slide_layouts[13]  # Final de la presentación
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    else:
        title_slide_layout = prs.slide_layouts[0]
        blank_layout = prs.slide_layouts[6]
        module_layout = title_slide_layout
        content_layout = blank_layout
        summary_layout = blank_layout
        property_layout = blank_layout
        description_layout = blank_layout
        objectives_layout = blank_layout
        prerequisites_layout = blank_layout
        audience_layout = blank_layout
        outline_layout = blank_layout
        presentation_layout = blank_layout
        lesson_layout = blank_layout
        lab_layout = blank_layout
        thanks_layout = blank_layout
    
    # Color scheme based on style
    style = structure.get('style', 'professional')
    color_map = {
        'professional': {'primary': RGBColor(0, 51, 102), 'accent': RGBColor(255, 192, 0)},
        'modern': {'primary': RGBColor(33, 33, 33), 'accent': RGBColor(0, 150, 136)},
        'minimal': {'primary': RGBColor(0, 0, 0), 'accent': RGBColor(117, 117, 117)}
    }
    colors = color_map.get(style, color_map['professional'])
    
    # Image cache to avoid duplicate downloads
    image_cache = {}
    
    def download_image_from_s3(image_ref: str) -> Optional[bytes]:
        """Download image from S3 based on reference (URL or image ID)."""
        try:
            if image_ref in image_cache:
                return image_cache[image_ref]
            
            # Check if image_ref is a full S3 URL
            if image_ref.startswith('http'):
                # Extract bucket and key from URL
                # Format: https://bucket.s3.amazonaws.com/key or https://s3.amazonaws.com/bucket/key
                import re
                url_match = re.search(r'https://([^/]+)\.s3\.amazonaws\.com/(.+)', image_ref)
                if url_match:
                    bucket = url_match.group(1)
                    key = url_match.group(2)
                    logger.info(f"📥 Downloading image from URL: s3://{bucket}/{key}")
                    response = s3_client.get_object(Bucket=bucket, Key=key)
                    image_bytes = response['Body'].read()
                    image_cache[image_ref] = image_bytes
                    return image_bytes
                else:
                    logger.warning(f"⚠️ Could not parse S3 URL: {image_ref}")
                    return None
            
            # Legacy: Extract image filename from reference (e.g., "01-01-0001.png")
            import re
            match = re.search(r'(\d{2}-\d{2}-\d{4})', image_ref)
            if not match:
                logger.warning(f"⚠️ Could not extract image ID from: {image_ref}")
                return None
            
            image_id = match.group(1)
            image_key = f"{project_folder}/images/{image_id}.png"
            
            logger.info(f"📥 Downloading image: s3://{course_bucket}/{image_key}")
            response = s3_client.get_object(Bucket=course_bucket, Key=image_key)
            image_bytes = response['Body'].read()
            image_cache[image_ref] = image_bytes
            return image_bytes
        except Exception as e:
            logger.warning(f"⚠️ Failed to download image {image_ref}: {e}")
            return None
    
    def get_image_dimensions(image_bytes: bytes) -> tuple:
        """Get image dimensions in inches for layout."""
        try:
            img = Image.open(io.BytesIO(image_bytes))
            width_px, height_px = img.size
            # Assume 96 DPI for screen images
            width_in = width_px / 96.0
            height_in = height_px / 96.0
            return (width_in, height_in)
        except:
            return (6.0, 4.0)  # Default
    
    def add_image_to_slide(slide, image_bytes: bytes, left: float, top: float, width: float, height: float):
        """Add image to slide at specified position."""
        try:
            pic_stream = io.BytesIO(image_bytes)
            pic = slide.shapes.add_picture(pic_stream, Inches(left), Inches(top), 
                                          width=Inches(width), height=Inches(height))
            logger.info(f"🖼️  Added image at ({left:.2f}, {top:.2f}) size {width:.2f}x{height:.2f}")
            return pic
        except Exception as e:
            logger.error(f"❌ Failed to add image: {e}")
            return None
    
    # Create slides
    slide_count = 0
    last_lesson_number = None  # Track lesson changes
    last_module_number = None  # Track module changes
    
    # Get course metadata from structure
    course_metadata = structure.get('course_metadata', {})
    outline_modules = structure.get('outline_modules', [])
    
    for slide_idx, slide_data in enumerate(structure.get('slides', [])):
        layout_hint = slide_data.get('layout_hint', 'single-column')
        slide_title = slide_data.get('title', 'Untitled')[:60]
        
        # Strip AI-generated continuation markers from titles
        # Our split logic will add proper continuation markers when needed
        import re
        original_title = slide_data.get('title', '')
        cleaned_title = re.sub(r'\s*\(cont\.?\s*\d+\)\s*$', '', original_title, flags=re.IGNORECASE)
        if cleaned_title != original_title:
            logger.info(f"🧹 Cleaned AI continuation marker from title: '{original_title}' → '{cleaned_title}'")
            slide_data['title'] = cleaned_title
            slide_title = cleaned_title[:60]
        
        logger.info(f"\n{'='*80}")
        logger.info(f"🔄 Processing slide {slide_idx + 1}/{len(structure.get('slides', []))}")
        logger.info(f"   Title: {slide_title}")
        logger.info(f"   Layout: {layout_hint}")
        logger.info(f"   Current slide_count: {slide_count}")
        logger.info(f"{'='*80}")
        
        if layout_hint == 'title':
            # TITLE SLIDE - Use course name from outline
            slide = prs.slides.add_slide(title_slide_layout)
            course_title = course_metadata.get('title', slide_data.get('title', 'Generated Course Book'))
            
            # Set title
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                slide.shapes.title.text = course_title
            else:
                # Find first text placeholder
                for shape in slide.placeholders:
                    if shape.has_text_frame:
                        shape.text = course_title
                        break
            
            # Set subtitle if available
            if len(slide.placeholders) > 1 and slide_data.get('subtitle'):
                for i, shape in enumerate(slide.placeholders):
                    if i > 0 and shape.has_text_frame:
                        shape.text = slide_data.get('subtitle', '')
                        break
            slide_count += 1
            logger.info(f"✅ Title slide #{slide_count} created")
            
            # ADD INTRO SLIDES after title
            if use_template and course_metadata:
                logger.info(f"🔧 Adding intro slides - course_metadata available: {bool(course_metadata)}")
                
                # 1. Propiedad Intelectual (no changes needed)
                slide = prs.slides.add_slide(property_layout)
                slide_count += 1
                logger.info(f"✅ Added Propiedad Intelectual slide")
                
                # 2. Course Description
                if course_metadata.get('description'):
                    slide = prs.slides.add_slide(description_layout)
                    logger.info(f"🔧 Description slide added, looking for placeholders...")
                    for shape in slide.placeholders:
                        if shape.has_text_frame:
                            # Placeholder 10 is the body text
                            if shape.placeholder_format.idx == 10:
                                shape.text = course_metadata['description']
                                logger.info(f"✅ Replaced course description in placeholder {shape.placeholder_format.idx}")
                                break
                    slide_count += 1
                
                # 3. Objectives / Learning Outcomes
                if course_metadata.get('learning_outcomes'):
                    slide = prs.slides.add_slide(objectives_layout)
                    for shape in slide.placeholders:
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            # Keep title, add objectives as bullets
                            if 'objetivo' in shape.text.lower():
                                continue  # Keep the title
                            else:
                                text_frame.clear()
                                for outcome in course_metadata['learning_outcomes']:
                                    p = text_frame.add_paragraph() if text_frame.paragraphs else text_frame.paragraphs[0]
                                    p.text = outcome
                                    p.level = 0
                                break
                    slide_count += 1
                
                # 4. Prerequisites
                if course_metadata.get('prerequisites'):
                    slide = prs.slides.add_slide(prerequisites_layout)
                    for shape in slide.placeholders:
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            if 'prerequisito' in shape.text.lower():
                                continue  # Keep the title
                            else:
                                text_frame.clear()
                                for prereq in course_metadata['prerequisites']:
                                    p = text_frame.add_paragraph() if text_frame.paragraphs else text_frame.paragraphs[0]
                                    p.text = prereq
                                    p.level = 0
                                break
                    slide_count += 1
                
                # 5. Audience
                if course_metadata.get('audience'):
                    slide = prs.slides.add_slide(audience_layout)
                    for shape in slide.placeholders:
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            if 'audiencia' in shape.text.lower():
                                continue  # Keep the title
                            else:
                                text_frame.clear()
                                for aud in course_metadata['audience']:
                                    p = text_frame.add_paragraph() if text_frame.paragraphs else text_frame.paragraphs[0]
                                    p.text = aud
                                    p.level = 0
                                break
                    slide_count += 1
                
                # 6. Course Outline / Temario
                if outline_modules:
                    slide = prs.slides.add_slide(outline_layout)
                    for shape in slide.placeholders:
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            if 'temario' in shape.text.lower():
                                continue  # Keep the title
                            else:
                                text_frame.clear()
                                for idx, mod in enumerate(outline_modules, 1):
                                    p = text_frame.add_paragraph() if text_frame.paragraphs else text_frame.paragraphs[0]
                                    p.text = f"Módulo {idx}: {mod.get('title', '')}"
                                    p.level = 0
                                break
                    slide_count += 1
                
                # 7. Presentación del grupo (no changes needed)
                slide = prs.slides.add_slide(presentation_layout)
                slide_count += 1
            
        elif layout_hint == 'module' or 'módulo' in slide_data.get('title', '').lower():
            # MODULE TITLE SLIDE - "Portada del capítulo" (Layout 8)
            slide = prs.slides.add_slide(module_layout)
            logger.info(f"🔧 Creating module slide for module {slide_data.get('module_number', '?')}")
            
            # Extract module number and get name from outline or slide data
            module_num = slide_data.get('module_number', 1)
            
            # Try to get module name from outline in structure
            module_title = None
            if outline_modules:
                if module_num <= len(outline_modules):
                    module_title = outline_modules[module_num - 1].get('title', '')
            
            if not module_title:
                module_title = slide_data.get('title', f'Módulo {module_num}')
            
            logger.info(f"🔧 Module {module_num}: {module_title}")
            
            # Replace ALL text placeholders with proper content
            for shape in slide.placeholders:
                if shape.has_text_frame:
                    try:
                        idx = shape.placeholder_format.idx
                        logger.info(f"  📝 Module placeholder idx {idx}")
                        # Placeholder 0: "Capítulo/Módulo X" - module number
                        if idx == 0:
                            shape.text = f"Módulo {module_num}"
                            logger.info(f"✅ Replaced module number in placeholder 0: Módulo {module_num}")
                        # Placeholder 10: "Nombre del capítulo" - module title (editable text box in template)
                        elif idx == 10:
                            shape.text = module_title
                            logger.info(f"✅ Replaced module title in placeholder 10: {module_title}")
                    except Exception as e:
                        logger.warning(f"Could not replace module text in placeholder {idx}: {e}")
            
            slide_count += 1
            logger.info(f"✅ Module title slide #{slide_count} created")
            
            # Note: Lesson title slides will be added before each lesson's content slides
            # They are NOT added here to avoid duplication
            
        elif layout_hint == 'summary' or 'resumen' in slide_data.get('title', '').lower():
            # SUMMARY SLIDE - "Resumen del capítulo" (Layout 11)
            # This is a LESSON summary (end of each lesson)
            
            # Check if this is the last lesson of a module to add labs/gracias AFTER this
            current_module = slide_data.get('module_number')
            current_lesson = slide_data.get('lesson_number')
            is_last_lesson_of_module = False
            
            # Look ahead to determine if this is the last lesson of the module
            if slide_idx + 1 < len(structure.get('slides', [])):
                next_slide = structure.get('slides', [])[slide_idx + 1]
                next_module = next_slide.get('module_number')
                next_lesson = next_slide.get('lesson_number')
                
                # This is the last lesson if next slide is in different module
                is_last_lesson_of_module = (next_module != current_module)
                
                logger.info(f"📊 Lesson summary M{current_module}L{current_lesson}, next M{next_module}L{next_lesson}, last_lesson={is_last_lesson_of_module}")
            else:
                # This is the last slide overall
                is_last_lesson_of_module = True
                logger.info(f"📊 Lesson summary is last slide overall - last_lesson=True")
            
            # Only create lesson summary slide, NOT module summary
            # Module summary will be replaced by labs + gracias
            logger.info(f"➕ Creating lesson summary slide for M{current_module}L{current_lesson}")
            
            try:
                slide = prs.slides.add_slide(summary_layout)
                
                # Use simple "Resumen" title
                summary_title = "Resumen"
                
                # Find and replace the title placeholder, then populate content
                content_added = False
                for shape in slide.placeholders:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        
                        # Check if this is the title placeholder
                        if 'resumen' in shape.text.lower():
                            # Replace template title
                            shape.text = summary_title
                            continue
                        
                        # This is the content placeholder
                        text_frame.clear()
                        text_frame.word_wrap = True
                        
                        # Set margins to avoid footer overlap
                        from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
                        text_frame.auto_size = MSO_AUTO_SIZE.NONE
                        text_frame.vertical_anchor = MSO_ANCHOR.TOP
                        
                        # Get placeholder dimensions
                        placeholder_height = shape.height
                        FOOTER_MARGIN = Inches(1.5)  # Larger margin for summary slides with bigger footer
                        max_summary_height = placeholder_height - FOOTER_MARGIN
                        
                        logger.info(f"📏 Summary placeholder height: {placeholder_height/914400:.2f}in, max: {max_summary_height/914400:.2f}in")
                        
                        # Font sizes and spacing IN EMUS (for proper comparison)
                        # Convert points to EMUs: 1 point = 12700 EMUs
                        heading_size = 22
                        bullet_size = 18
                        heading_height = Pt(heading_size).emu + Pt(18).emu  # heading + space after + space before
                        bullet_height = Pt(bullet_size).emu + Pt(26).emu    # bullet + space after + line height
                        
                        # Estimate and split content if needed
                        all_blocks = slide_data.get('content_blocks', [])
                        current_blocks = []
                        summary_slides_content = []
                        current_height = 0
                        
                        logger.info(f"📋 Summary has {len(all_blocks)} content blocks to process")
                        
                        for block_idx, block in enumerate(all_blocks):
                            if block.get('type') == 'bullets':
                                # Estimate height for this block IN EMUS
                                block_height = 0
                                if block.get('heading'):
                                    block_height += heading_height
                                num_items = len(block.get('items', []))
                                block_height += num_items * bullet_height
                                
                                logger.info(f"  Block {block_idx+1}: {num_items} items, height={block_height/914400:.2f}in, current_total={current_height/914400:.2f}in, would_be={(current_height + block_height)/914400:.2f}in")
                                
                                # Check if adding this block would exceed space
                                if current_height + block_height > max_summary_height and current_blocks:
                                    # Save current blocks and start new slide
                                    logger.info(f"📄 Summary content exceeds limit ({(current_height + block_height)/914400:.2f}in > {max_summary_height/914400:.2f}in), splitting")
                                    summary_slides_content.append(current_blocks)
                                    current_blocks = [block]
                                    current_height = block_height
                                else:
                                    current_blocks.append(block)
                                    current_height += block_height
                        
                        # Add remaining blocks
                        if current_blocks:
                            summary_slides_content.append(current_blocks)
                        
                        logger.info(f"📊 Summary split into {len(summary_slides_content)} slide(s)")
                        
                        # Add content for FIRST summary slide
                        if summary_slides_content:
                            first_summary_blocks = summary_slides_content[0]
                            first_para = True
                            
                            for block in first_summary_blocks:
                                if block.get('type') == 'bullets':
                                    # Add heading if present - NO BULLET for heading
                                    if block.get('heading'):
                                        if first_para:
                                            p = text_frame.paragraphs[0]
                                            first_para = False
                                        else:
                                            p = text_frame.add_paragraph()
                                        p.text = block['heading']
                                        p.font.size = Pt(22)
                                        p.font.bold = True
                                        p.space_after = Pt(12)
                                        p.space_before = Pt(6)
                                    
                                    # Add bullet items - DON'T add "•" since template adds it
                                    for item in block.get('items', []):
                                        p = text_frame.add_paragraph()
                                        # Strip existing bullets from item text
                                        item_text = item.lstrip('• ').lstrip('- ')
                                        p.text = item_text
                                        p.level = 0
                                        p.font.size = Pt(18)
                                        p.space_after = Pt(8)
                        
                        # Create continuation summary slides if needed
                        for cont_idx, cont_blocks in enumerate(summary_slides_content[1:], start=2):
                            logger.info(f"📄 Creating summary continuation slide {cont_idx}")
                            cont_slide = prs.slides.add_slide(summary_layout)
                            
                            # Set title
                            for shape in cont_slide.placeholders:
                                if shape.has_text_frame and 'resumen' in shape.text.lower():
                                    shape.text = f"{summary_title} (cont. {cont_idx})"
                                    break
                            
                            # Add content to continuation slide
                            for shape in cont_slide.placeholders:
                                if shape.has_text_frame and 'resumen' not in shape.text.lower():
                                    cont_text_frame = shape.text_frame
                                    cont_text_frame.clear()
                                    cont_text_frame.word_wrap = True
                                    cont_text_frame.auto_size = MSO_AUTO_SIZE.NONE
                                    cont_text_frame.vertical_anchor = MSO_ANCHOR.TOP
                                    
                                    first_para_cont = True
                                    for block in cont_blocks:
                                        if block.get('type') == 'bullets':
                                            if block.get('heading'):
                                                if first_para_cont:
                                                    p = cont_text_frame.paragraphs[0]
                                                    first_para_cont = False
                                                else:
                                                    p = cont_text_frame.add_paragraph()
                                                p.text = block['heading']
                                                p.font.size = Pt(22)
                                                p.font.bold = True
                                                p.space_after = Pt(12)
                                                p.space_before = Pt(6)
                                            
                                            for item in block.get('items', []):
                                                p = cont_text_frame.add_paragraph()
                                                item_text = item.lstrip('• ').lstrip('- ')
                                                p.text = item_text
                                                p.level = 0
                                                p.font.size = Pt(18)
                                                p.space_after = Pt(8)
                                    break
                            
                            slide_count += 1
                            logger.info(f"✅ Summary continuation slide #{slide_count} created")
                        
                        content_added = True
                        break
                
                if not content_added:
                    logger.warning(f"⚠️ No text placeholder found in summary layout")
                
                slide_count += 1
                logger.info(f"✅ Lesson summary slide created (slide #{slide_count})")
                
            except Exception as e:
                logger.error(f"❌ Failed to create lesson summary slide: {e}")
                # Create a fallback blank slide to maintain slide count
                logger.info(f"🔧 Creating fallback slide to maintain count")
                slide = prs.slides.add_slide(content_layout)
                # Add error message as title
                for shape in slide.placeholders:
                    if shape.has_text_frame:
                        shape.text = f"Error: Failed to create summary slide"
                        break
                slide_count += 1
            
            # ADD LAB ACTIVITY SLIDES and GRACIAS only if this is the LAST LESSON of the module
            if is_last_lesson_of_module and use_template and outline_modules and current_module and current_module <= len(outline_modules):
                logger.info(f"🎯 LAST LESSON OF MODULE {current_module} - Adding labs and Gracias")
                module_data = outline_modules[current_module - 1]
                
                lab_activities = module_data.get('lab_activities', [])
                logger.info(f"📝 Module {current_module} has {len(lab_activities)} lab activities")
                
                # Add lab slides
                for lab_idx, lab in enumerate(lab_activities, 1):
                    try:
                        lab_slide = prs.slides.add_slide(lab_layout)
                        lab_title = lab.get('title', '')
                        lab_duration = lab.get('duration_minutes', 0)
                        
                        logger.info(f"🧪 Adding lab {lab_idx}/{len(lab_activities)}: {lab_title}")
                        
                        # Populate lab slide placeholders
                        for shape in lab_slide.placeholders:
                            if shape.has_text_frame:
                                idx = shape.placeholder_format.idx
                                # Placeholder 0: "Nombre de la actividad"
                                if idx == 0:
                                    shape.text = lab_title
                                    logger.info(f"  ✅ Lab title set in placeholder 0: {lab_title}")
                                # Placeholder 12: "XX minutos."
                                elif idx == 12:
                                    shape.text = f"{lab_duration} minutos."
                                    logger.info(f"  ✅ Lab duration set in placeholder 12: {lab_duration} min")
                        slide_count += 1
                        logger.info(f"✅ Lab slide #{slide_count} created successfully")
                        
                    except Exception as e:
                        logger.error(f"❌ Failed to create lab slide: {e}")
                        # Create fallback slide
                        try:
                            fallback_slide = prs.slides.add_slide(content_layout)
                            for shape in fallback_slide.placeholders:
                                if shape.has_text_frame:
                                    shape.text = f"Lab: {lab_title}"
                                    break
                            slide_count += 1
                            logger.info(f"🔧 Created fallback lab slide #{slide_count}")
                        except:
                            logger.error(f"❌ Could not create fallback slide")
                
                # ADD "GRACIAS" SLIDE after labs (end of module)
                try:
                    logger.info(f"🎉 Adding 'Gracias' slide after module {current_module}")
                    gracias_slide = prs.slides.add_slide(thanks_layout)
                    slide_count += 1
                    logger.info(f"✅ 'Gracias' slide #{slide_count} added successfully")
                except Exception as e:
                    logger.error(f"❌ Failed to create Gracias slide: {e}")
                    # Create fallback slide
                    try:
                        fallback_slide = prs.slides.add_slide(title_slide_layout)
                        for shape in fallback_slide.placeholders:
                            if shape.has_text_frame:
                                shape.text = "¡Gracias!"
                                break
                        slide_count += 1
                        logger.info(f"🔧 Created fallback Gracias slide #{slide_count}")
                    except:
                        logger.error(f"❌ Could not create fallback Gracias slide")

            
        elif layout_hint == 'thanks' or 'gracias' in slide_data.get('title', '').lower():
            # THANKS SLIDE - This should have been removed in preprocessing
            # But if it somehow got through, skip it COMPLETELY
            logger.warning(f"⚠️ Found 'Gracias' slide that should have been removed: {slide_data.get('title')}")
            logger.warning(f"⚠️ Skipping this slide entirely - NOT incrementing slide_count")
            continue  # Skip to next iteration without creating any slide
            
        else:
            # CONTENT SLIDE - Check if we need to add a lesson title slide first
            current_lesson = slide_data.get('lesson_number')
            current_module = slide_data.get('module_number')
            
            # Add lesson title slide when lesson changes
            if use_template and current_lesson and current_lesson != last_lesson_number:
                try:
                    lesson_slide = prs.slides.add_slide(lesson_layout)
                    
                    # Get module and lesson info from outline
                    if outline_modules and current_module and current_module <= len(outline_modules):
                        module_data = outline_modules[current_module - 1]
                        module_title = module_data.get('title', f'Módulo {current_module}')
                        
                        # Find the lesson by index
                        lessons = module_data.get('lessons', [])
                        lesson_idx = current_lesson - 1 if current_module == 1 else current_lesson - sum(len(outline_modules[i].get('lessons', [])) for i in range(current_module - 1)) - 1
                        
                        if 0 <= lesson_idx < len(lessons):
                            lesson_title = lessons[lesson_idx].get('title', f'Lección {current_lesson}')
                        else:
                            lesson_title = f'Lección {current_lesson}'
                    else:
                        module_title = f'Módulo {current_module}' if current_module else 'Módulo'
                        lesson_title = f'Lección {current_lesson}'
                    
                    # Populate lesson slide placeholders (indices switched in new template)
                    for shape in lesson_slide.placeholders:
                        if shape.has_text_frame:
                            idx = shape.placeholder_format.idx
                            # Placeholder 0: NOW "Nombre del capítulo o módulo" (module - was lesson)
                            if idx == 0:
                                shape.text = module_title
                                logger.info(f"✅ Replaced module in lesson slide placeholder 0: {module_title}")
                            # Placeholder 10: NOW "Nombre de la Lección" (lesson - was module)
                            elif idx == 10:
                                shape.text = lesson_title
                                logger.info(f"✅ Replaced lesson title in placeholder 10: {lesson_title}")
                    
                    slide_count += 1
                    last_lesson_number = current_lesson
                    logger.info(f"✅ Lesson title slide #{slide_count} created for L{current_lesson}")
                    
                except Exception as e:
                    logger.error(f"❌ Failed to create lesson title slide: {e}")
                    # Create fallback slide
                    try:
                        fallback_slide = prs.slides.add_slide(title_slide_layout)
                        for shape in fallback_slide.placeholders:
                            if shape.has_text_frame:
                                shape.text = f"Lección {current_lesson}"
                                break
                        slide_count += 1
                        last_lesson_number = current_lesson
                        logger.info(f"🔧 Created fallback lesson title slide #{slide_count}")
                    except:
                        logger.error(f"❌ Could not create fallback lesson title slide")
            
            # REGULAR CONTENT SLIDE - "Contenido - General" (Layout 10)
            try:
                if use_template:
                    slide = prs.slides.add_slide(content_layout)
                    use_template_placeholders = True
                else:
                    slide = prs.slides.add_slide(blank_layout)
                    use_template_placeholders = False
                    
                    # Background for non-template
                    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
                    bg.fill.solid()
                    bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
                
                logger.info(f"➕ Creating content slide: {slide_data.get('title', 'Untitled')[:50]}")
                
            except Exception as e:
                logger.error(f"❌ Failed to create content slide: {e}")
                # Create fallback blank slide
                try:
                    slide = prs.slides.add_slide(blank_layout if use_template else prs.slide_layouts[0])
                    use_template_placeholders = False
                    logger.info(f"🔧 Created fallback blank slide")
                except Exception as e2:
                    logger.error(f"❌ Could not create fallback slide: {e2}")
                    # Skip this slide entirely
                    continue
            
            # Set title using template placeholder or text box
            try:
                if use_template_placeholders:
                    # Replace "Desarrollo del contenido" with actual title
                    topic_title = slide_data.get('title', '')
                    for shape in slide.placeholders:
                        if shape.has_text_frame:
                            try:
                                current_text = shape.text.strip()
                                # Replace template text with actual content
                                if 'desarrollo' in current_text.lower() or not current_text:
                                    shape.text = topic_title
                                    break
                            except:
                                pass
                else:
                    # Create text box for title
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3),
                        Inches(12.3), Inches(0.8)
                    )
                    title_frame = title_box.text_frame
                    title_para = title_frame.paragraphs[0]
                    title_para.text = slide_data.get('title', '')
                    title_para.font.size = Pt(40)
                    title_para.font.bold = True
                    title_para.font.color.rgb = colors['primary']
            except Exception as e:
                logger.warning(f"⚠️ Failed to set slide title: {e}")
            
            # Determine if slide has images and text
            image_blocks = [b for b in slide_data.get('content_blocks', []) if b.get('type') == 'image']
            text_blocks = [b for b in slide_data.get('content_blocks', []) if b.get('type') == 'bullets']
            has_images = len(image_blocks) > 0
            has_text = len(text_blocks) > 0
            
            # IMAGE-ONLY LAYOUT (centered, large)
            if has_images and not has_text and course_bucket and project_folder:
                logger.info(f"📐 Using centered image layout for slide {slide_count + 1}")
                
                # Download first image
                image_ref = image_blocks[0].get('image_reference', '')
                image_bytes = download_image_from_s3(image_ref)
                
                if image_bytes:
                    # Calculate image dimensions for centered display
                    img_width_in, img_height_in = get_image_dimensions(image_bytes)
                    
                    # Scale image to fit most of the slide (larger for image-only)
                    max_img_width = 10.0
                    max_img_height = 5.5
                    
                    scale = min(max_img_width / img_width_in, max_img_height / img_height_in)
                    img_width = img_width_in * scale
                    img_height = img_height_in * scale
                    
                    # Center the image horizontally
                    img_left = (13.333 - img_width) / 2
                    img_top = 1.5
                    
                    add_image_to_slide(slide, image_bytes, img_left, img_top, img_width, img_height)
                    logger.info(f"🎨 Centered image: {img_width:.2f}x{img_height:.2f} at x={img_left:.2f}")
                
                slide_count += 1
                continue  # Skip text processing for image-only slides
            
            # Layout strategy: side-by-side if image+text, full-width if text-only
            logger.info(f"📐 Processing content blocks for slide")
            if has_images and has_text and course_bucket and project_folder:
                # IMAGE + TEXT LAYOUT (side-by-side)
                logger.info(f"📐 Using side-by-side layout for slide {slide_count + 1}")
                
                # Download first image
                image_ref = image_blocks[0].get('image_reference', '')
                image_bytes = download_image_from_s3(image_ref)
                
                if image_bytes:
                    # Calculate image dimensions (right side)
                    img_width_in, img_height_in = get_image_dimensions(image_bytes)
                    
                    # Scale image to fit right half
                    max_img_width = 5.5
                    max_img_height = 5.5
                    
                    scale = min(max_img_width / img_width_in, max_img_height / img_height_in)
                    img_width = img_width_in * scale
                    img_height = img_height_in * scale
                    
                    # Position image on right side
                    img_left = 7.0
                    img_top = 1.5
                    
                    add_image_to_slide(slide, image_bytes, img_left, img_top, img_width, img_height)
                    
                    # Text content on left side (NO OVERLAP!)
                    text_left = 0.8
                    text_top = 1.5
                    text_width = 5.5
                    text_height = 5.5
                else:
                    # No image downloaded, use full width for text
                    text_left = 0.8
                    text_top = 1.5
                    text_width = 11.5
                    text_height = 5.5
            else:
                # TEXT-ONLY LAYOUT (full width)
                text_left = 0.8
                text_top = 1.5
                text_width = 11.5
                text_height = 5.5
            
            # Add content blocks (bullets, text, etc.)
            # Strategy: Max 2 content blocks per slide, create additional slides if needed
            content_blocks = [b for b in slide_data.get('content_blocks', []) if b.get('type') == 'bullets']
            
            # Skip creating slide if no content blocks (prevent blank slides)
            if not content_blocks:
                logger.warning(f"⚠️ No content blocks for slide '{slide_data.get('title', 'Unknown')}', skipping")
                continue
            
            # Split content blocks into groups of max 2 blocks per slide
            blocks_per_slide = 2
            block_groups = []
            for i in range(0, len(content_blocks), blocks_per_slide):
                block_groups.append(content_blocks[i:i + blocks_per_slide])
            
            # Process each group (create multiple slides if needed)
            for group_idx, blocks_group in enumerate(block_groups):
                # For subsequent groups, create new slide
                if group_idx > 0:
                    logger.info(f"📄 Creating continuation slide {group_idx + 1} for topic")
                    if use_template:
                        slide = prs.slides.add_slide(content_layout)
                        use_template_placeholders = True
                    else:
                        slide = prs.slides.add_slide(blank_layout)
                        use_template_placeholders = False
                        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
                        bg.fill.solid()
                        bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
                    
                    # Set title for continuation slide
                    if use_template_placeholders:
                        topic_title = slide_data.get('title', '') + f" (cont. {group_idx + 1})"
                        for shape in slide.placeholders:
                            if shape.has_text_frame:
                                try:
                                    current_text = shape.text.strip()
                                    if 'desarrollo' in current_text.lower() or not current_text:
                                        shape.text = topic_title
                                        break
                                except:
                                    pass
                    
                    slide_count += 1
            
            if use_template_placeholders:
                # Use template content placeholder (idx 11 for "Contenido - General" layout)
                content_placeholder = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 11 and shape.has_text_frame:
                        content_placeholder = shape
                        break
                
                if content_placeholder:
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()
                    text_frame.word_wrap = True
                    
                    # FIX 4: Set vertical alignment to TOP (not centered)
                    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
                    text_frame.auto_size = MSO_AUTO_SIZE.NONE
                    text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Align text to top
                    
                    # Get placeholder dimensions to calculate available space
                    placeholder_height = content_placeholder.height
                    placeholder_top = content_placeholder.top
                    
                    # CRITICAL: Set bottom margin to avoid footer overlap
                    # Standard PowerPoint slide is 7.5" tall
                    # Footer typically starts at 6.8" from top
                    # Leave 1.2" margin from bottom to avoid logo overlap (increased for safety)
                    FOOTER_MARGIN = Inches(1.2)  # Increased safety margin for footer/logo
                    max_content_height = placeholder_height - FOOTER_MARGIN
                    
                    logger.info(f"📏 Placeholder height: {placeholder_height/914400:.2f}in, Max content height: {max_content_height/914400:.2f}in")
                    
                    # Use consistent, readable font sizes
                    heading_size = 22  # Slightly smaller to fit more
                    bullet_size = 18
                    spacer_size = 10  # FIX 1: Reduced from 16 to 10 for tighter spacing
                    heading_space_after = 8  # FIX 1: Reduced from 12 to 8
                    bullet_space_after = 4  # FIX 1: Reduced from 8 to 4
                    
                    # Helper function to estimate content height
                    def estimate_block_height(block, is_first=False):
                        """Estimate the height needed for a content block with safety margin"""
                        height = 0
                        if not is_first and block.get('heading'):
                            height += spacer_size  # Spacer before block
                        if block.get('heading'):
                            # Account for potential line wrapping in headings
                            height += heading_size + heading_space_after + 8  # Extra space for wrapping
                        num_items = len(block.get('items', []))
                        # Assume average bullet might wrap to 1.5 lines (conservative estimate)
                        height += num_items * (bullet_size + bullet_space_after + 10)  # Increased for wrapping
                        return height
                    
                    # Split content blocks intelligently based on available space
                    current_slide_blocks = []
                    slides_content = []  # List of lists of blocks per slide
                    current_height = 0
                    
                    for block_idx, block in enumerate(blocks_group):
                        block_height = estimate_block_height(block, is_first=(len(current_slide_blocks) == 0))
                        
                        # Check if adding this block would exceed available space
                        if current_height + block_height > max_content_height and current_slide_blocks:
                            # Start a new slide
                            logger.info(f"📄 Content exceeds slide limit ({current_height + block_height:.0f}pt > {max_content_height/914400*72:.0f}pt), splitting to new slide")
                            slides_content.append(current_slide_blocks)
                            current_slide_blocks = [block]
                            current_height = block_height
                        else:
                            # Add to current slide
                            current_slide_blocks.append(block)
                            current_height += block_height
                    
                    # Add remaining blocks
                    if current_slide_blocks:
                        slides_content.append(current_slide_blocks)
                    
                    logger.info(f"📊 Split content into {len(slides_content)} slide(s)")
                    
                    # Process FIRST slide's content (current slide)
                    if slides_content:
                        first_slide_blocks = slides_content[0]
                        
                        # Add content blocks for this slide group
                        first_para = True
                        for block_idx, block in enumerate(first_slide_blocks):
                            if block.get('type') == 'bullets':
                                # Add spacing between blocks (except first)
                                if block_idx > 0:
                                    spacer = text_frame.add_paragraph()
                                    spacer.text = ""
                                    spacer.space_after = Pt(spacer_size)
                                
                                # Add heading
                                if block.get('heading'):
                                    if not first_para:
                                        p = text_frame.add_paragraph()
                                    else:
                                        p = text_frame.paragraphs[0]
                                        first_para = False
                                    
                                    # Strip bullet prefix if present in heading text
                                    heading_text = block['heading']
                                    if heading_text.startswith('• '):
                                        heading_text = heading_text[2:]
                                    
                                    p.text = heading_text
                                    p.font.size = Pt(heading_size)
                                    p.font.bold = True
                                    
                                    # FIX 3: Remove bullet from heading by adding buNone element
                                    try:
                                        # Access the paragraph XML element
                                        pPr = p._element.get_or_add_pPr()
                                        # Add buNone element to explicitly remove bullet
                                        from lxml import etree
                                        buNone = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                                        if buNone is None:
                                            # Create buNone element to remove bullet
                                            buNone = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                                            pPr.insert(0, buNone)
                                    except Exception as e:
                                        # Fallback: at least remove indent
                                        try:
                                            from pptx.util import Inches
                                            p.paragraph_format.left_indent = Inches(0)
                                            p.paragraph_format.first_line_indent = Inches(0)
                                        except:
                                            pass
                                    
                                    p.space_after = Pt(heading_space_after)
                                    p.space_before = Pt(4) if block_idx > 0 else Pt(0)
                                    # FIX 1: Reduce line spacing within heading
                                    p.line_spacing = 1.0  # Single spacing
                                
                                # Add bullets
                                for item in block.get('items', []):
                                    p = text_frame.add_paragraph()
                                    p.text = item
                                    p.level = 0  # Bullet items keep level 0
                                    p.font.size = Pt(bullet_size)
                                    p.space_after = Pt(bullet_space_after)
                                    # FIX 1: Reduce line spacing within bullets
                                    p.line_spacing = 1.0  # Single spacing
                    
                    # Create CONTINUATION SLIDES for remaining content
                    for continuation_idx, continuation_blocks in enumerate(slides_content[1:], start=2):
                        logger.info(f"📄 Creating continuation slide {continuation_idx} for long content")
                        
                        # Create new slide with same layout
                        continuation_slide = prs.slides.add_slide(content_layout)
                        
                        # Set title to indicate continuation
                        topic_title = slide_data.get('title', '') + f" (cont. {continuation_idx})"
                        for shape in continuation_slide.placeholders:
                            if shape.has_text_frame:
                                try:
                                    current_text = shape.text.strip()
                                    if 'desarrollo' in current_text.lower() or not current_text:
                                        shape.text = topic_title
                                        break
                                except:
                                    pass
                        
                        # Find content placeholder in continuation slide
                        continuation_content_placeholder = None
                        for shape in continuation_slide.placeholders:
                            if shape.placeholder_format.idx == 11 and shape.has_text_frame:
                                continuation_content_placeholder = shape
                                break
                        
                        if continuation_content_placeholder:
                            cont_text_frame = continuation_content_placeholder.text_frame
                            cont_text_frame.clear()
                            cont_text_frame.word_wrap = True
                            cont_text_frame.auto_size = MSO_AUTO_SIZE.NONE
                            cont_text_frame.vertical_anchor = MSO_ANCHOR.TOP
                            
                            # Add continuation content
                            first_para_cont = True
                            for block_idx, block in enumerate(continuation_blocks):
                                if block.get('type') == 'bullets':
                                    # Add spacing between blocks (except first)
                                    if block_idx > 0:
                                        spacer = cont_text_frame.add_paragraph()
                                        spacer.text = ""
                                        spacer.space_after = Pt(spacer_size)
                                    
                                    # Add heading
                                    if block.get('heading'):
                                        if not first_para_cont:
                                            p = cont_text_frame.add_paragraph()
                                        else:
                                            p = cont_text_frame.paragraphs[0]
                                            first_para_cont = False
                                        
                                        heading_text = block['heading']
                                        if heading_text.startswith('• '):
                                            heading_text = heading_text[2:]
                                        
                                        p.text = heading_text
                                        p.font.size = Pt(heading_size)
                                        p.font.bold = True
                                        
                                        # Remove bullet from heading
                                        try:
                                            pPr = p._element.get_or_add_pPr()
                                            from lxml import etree
                                            buNone = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                                            if buNone is None:
                                                buNone = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}buNone')
                                                pPr.insert(0, buNone)
                                        except:
                                            try:
                                                from pptx.util import Inches
                                                p.paragraph_format.left_indent = Inches(0)
                                                p.paragraph_format.first_line_indent = Inches(0)
                                            except:
                                                pass
                                        
                                        p.space_after = Pt(heading_space_after)
                                        p.space_before = Pt(4) if block_idx > 0 else Pt(0)
                                        p.line_spacing = 1.0
                                    
                                    # Add bullets
                                    for item in block.get('items', []):
                                        p = cont_text_frame.add_paragraph()
                                        p.text = item
                                        p.level = 0
                                        p.font.size = Pt(bullet_size)
                                        p.space_after = Pt(bullet_space_after)
                                        p.line_spacing = 1.0
                        
                        slide_count += 1
                        logger.info(f"✅ Continuation slide #{slide_count} created successfully")
                else:
                    # Fallback to manual text boxes
                    use_template_placeholders = False
            
            if not use_template_placeholders:
                # Original manual text box approach
                y_offset = text_top
                
                # CRITICAL: Define maximum y position to avoid footer overlap
                # Standard slide is 7.5" tall, footer/logo typically at ~6.8"
                # Increased safety margin to prevent any overlap
                max_y_position = 6.0  # Leave 1.5" margin from bottom for footer/logo
                
                for block_idx, block in enumerate(slide_data.get('content_blocks', [])):
                    if block.get('type') == 'bullets':
                        # Add extra spacing before heading (except first block)
                        if block_idx > 0 and block.get('heading'):
                            y_offset += 0.3  # Moderate space before new section
                        
                        # Calculate estimated height needed for this block
                        num_items = len(block.get('items', []))
                        has_heading = bool(block.get('heading'))
                        estimated_height = 0.35 if has_heading else 0.15  # Reduced from 0.4/0.2
                        estimated_height += (num_items * 0.28)  # Reduced from 0.35 inches per bullet
                        
                        # Ensure we don't go off the slide or overlap footer
                        available_height = max_y_position - y_offset
                        if available_height < 0.5:
                            logger.warning(f"⚠️ Not enough space for block {block_idx} (y={y_offset:.2f}, max={max_y_position}), skipping to avoid footer overlap")
                            break
                        
                        block_height = min(estimated_height, available_height)
                        
                        # Add text box for bullets (EDITABLE!)
                        text_box = slide.shapes.add_textbox(
                            Inches(text_left), Inches(y_offset),
                            Inches(text_width), Inches(block_height)
                        )
                        text_frame = text_box.text_frame
                        text_frame.word_wrap = True
                        text_frame.auto_size = None  # Disable auto-resize
                        
                        # Heading with more spacing
                        if block.get('heading'):
                            p = text_frame.paragraphs[0]
                            p.text = block['heading']
                            p.font.size = Pt(24)  # Reduced from 26
                            p.font.bold = True
                            p.font.color.rgb = colors['primary']
                            p.space_before = Pt(8)  # Reduced from 12
                            p.space_after = Pt(10)  # Reduced from 12
                        
                        # Bullets
                        for item in block.get('items', []):
                            p = text_frame.add_paragraph()
                            p.text = item
                            p.level = 0
                            p.font.size = Pt(18)
                            p.space_after = Pt(6)  # Reduced from 8
                        
                        # ALWAYS advance y_offset to prevent overlaps - REDUCED gap
                        y_offset += block_height + 0.15  # Reduced from 0.25
            
            # Increment slide count ONLY for content slides that were actually created
            slide_count += 1
            logger.info(f"✅ Content slide #{slide_count} created successfully")
    
    # Final validation and logging
    actual_slide_count = len(prs.slides)
    logger.info(f"\n{'='*80}")
    logger.info(f"📊 FINAL PPT STATISTICS")
    logger.info(f"{'='*80}")
    logger.info(f"Total slides created: {actual_slide_count}")
    logger.info(f"Slide counter value: {slide_count}")
    logger.info(f"Input slides processed: {len(structure.get('slides', []))}")
    
    if actual_slide_count != slide_count:
        logger.warning(f"⚠️ MISMATCH: Actual slides ({actual_slide_count}) ≠ Counter ({slide_count})")
    else:
        logger.info(f"✅ Slide count verification passed")
    
    # Log slide breakdown
    logger.info(f"\n📋 Slide Breakdown:")
    for idx, slide in enumerate(prs.slides, 1):
        slide_title = "Untitled"
        try:
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                slide_title = slide.shapes.title.text[:50]
            else:
                # Try to get title from first text placeholder
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame.text:
                        slide_title = shape.text_frame.text[:50]
                        break
        except:
            pass
        logger.info(f"  Slide {idx}: {slide_title}")
    
    logger.info(f"{'='*80}\n")
    
    # Save to bytes
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    logger.info(f"✅ Created editable PowerPoint with {len(structure.get('slides', []))} slides")
    return pptx_buffer.getvalue()


def lambda_handler(event, context):
    """Main Lambda handler for infographic generation."""
    try:
        logger.info("=" * 80)
        logger.info("🎨 INFOGRAPHIC GENERATOR - HTML TO EDITABLE PPT")
        logger.info("=" * 80)
        
        # Parse event
        if 'body' in event:
            body = json.loads(event['body']) if isinstance(event['body'], str) else event['body']
        else:
            body = event
        
        # Extract parameters
        course_bucket = body.get('course_bucket')
        project_folder = body.get('project_folder')
        book_version_key = body.get('book_version_key')
        book_type = body.get('book_type', 'theory')  # 'theory' or 'lab'
        model_provider = body.get('model_provider', 'bedrock').lower()
        slides_per_lesson = int(body.get('slides_per_lesson', 5))
        style = body.get('style', 'professional')
        
        # Log what we received
        logger.info(f"📥 Received book_version_key: {book_version_key}")
        logger.info(f"📥 Book type: {book_type}")
        logger.info(f"📥 Project folder: {project_folder}")
        
        # Batch processing
        lesson_start = int(body.get('lesson_start', 1))
        lesson_end = body.get('lesson_end')
        max_lessons_per_batch = int(body.get('max_lessons_per_batch', 10))
        
        if not course_bucket or not project_folder:
            raise ValueError("course_bucket and project_folder required")
        
        # Auto-discover book version if not provided
        if not book_version_key:
            logger.info("🔎 No book_version_key provided, starting auto-discovery...")
            # Try multiple possible book folder structures
            possible_folders = [
                f"{project_folder}/{book_type}-book/",  # New structure
                f"{project_folder}/book/",              # Legacy structure
                f"{project_folder}/"                     # Root folder
            ]
            
            book_files = []
            for book_folder in possible_folders:
                logger.info(f"🔍 Searching in {book_folder}")
                
                try:
                    response = s3_client.list_objects_v2(
                        Bucket=course_bucket,
                        Prefix=book_folder,
                        Delimiter='/'
                    )
                    
                    # Find book files (book_version_*.json or Generated_Course_Book_data.json)
                    folder_files = [
                        obj['Key'] for obj in response.get('Contents', [])
                        if obj['Key'].endswith('.json') and (
                            'book_version' in obj['Key'] or 
                            'Generated_Course_Book_data' in obj['Key'] or
                            'Book_data' in obj['Key']
                        )
                    ]
                    
                    if folder_files:
                        book_files = folder_files
                        logger.info(f"✅ Found {len(book_files)} book file(s) in {book_folder}")
                        break
                        
                except Exception as e:
                    logger.warning(f"Could not search {book_folder}: {e}")
                    continue
            
            if not book_files:
                raise ValueError(f"No book found. Searched in: {', '.join(possible_folders)}")
            
            # Sort by timestamp (newest first) and pick the first one
            book_files.sort(reverse=True)
            book_version_key = book_files[0]
            logger.info(f"✅ Auto-discovered: {book_version_key}")
        else:
            logger.info(f"✅ Using provided book_version_key: {book_version_key}")
        
        logger.info(f"📊 Slides per lesson: {slides_per_lesson}")
        logger.info(f"🎨 Style: {style}")
        logger.info(f"📦 Batch: lessons {lesson_start} to {lesson_end or 'end'}")
        
        # Configure AI model
        if model_provider == 'openai':
            model = configure_openai_model()
        else:
            model = configure_bedrock_model()
        
        # Load book data
        logger.info(f"📚 Loading book from S3: {course_bucket}/{book_version_key}")
        book_data = load_book_from_s3(course_bucket, book_version_key)
        
        # Apply batching
        total_lessons = len(book_data.get('lessons', []))
        if lesson_end:
            lesson_end = min(int(lesson_end), total_lessons)
        else:
            lesson_end = min(lesson_start + max_lessons_per_batch - 1, total_lessons)
        
        original_lessons = book_data.get('lessons', [])
        book_data['lessons'] = original_lessons[lesson_start-1:lesson_end]
        
        logger.info(f"📖 Processing lessons {lesson_start}-{lesson_end} of {total_lessons}")
        
        # Generate infographic structure
        structure = generate_infographic_structure(
            book_data, model, slides_per_lesson, style
        )
        
        # Save structure to S3
        structure_key = f"{project_folder}/infographics/infographic_structure.json"
        s3_client.put_object(
            Bucket=course_bucket,
            Key=structure_key,
            Body=json.dumps(structure, indent=2),
            ContentType='application/json'
        )
        logger.info(f"💾 Saved structure: s3://{course_bucket}/{structure_key}")
        
        # Generate HTML
        html_content = generate_html_from_structure(structure)
        html_key = f"{project_folder}/infographics/infographic.html"
        s3_client.put_object(
            Bucket=course_bucket,
            Key=html_key,
            Body=html_content,
            ContentType='text/html'
        )
        logger.info(f"💾 Saved HTML: s3://{course_bucket}/{html_key}")
        
        # Convert to editable PPT with images and template
        template_key = "PPT_Templates/Esandar_Aumentado_3.pptx"  # Your custom template
        pptx_bytes = convert_html_to_pptx(html_content, structure, course_bucket, project_folder, template_key)
        
        course_title = structure.get('course_title', 'presentation')
        sanitized_filename = course_title.replace(' ', '_').replace('/', '-')
        pptx_key = f"{project_folder}/infographics/{sanitized_filename}.pptx"
        
        s3_client.put_object(
            Bucket=course_bucket,
            Key=pptx_key,
            Body=pptx_bytes,
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        logger.info(f"💾 Saved editable PPT: s3://{course_bucket}/{pptx_key}")
        
        # Return success
        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Headers': 'Content-Type',
                'Access-Control-Allow-Methods': 'POST,OPTIONS',
                'Content-Type': 'application/json'
            },
            'body': json.dumps({
                'message': 'Infographic generated successfully',
                'course_title': structure['course_title'],
                'total_slides': structure['total_slides'],
                'completion_status': structure.get('completion_status', 'complete'),
                'structure_s3_key': structure_key,
                'html_s3_key': html_key,
                'pptx_s3_key': pptx_key
            })
        }
        
    except Exception as e:
        logger.error(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        
        return {
            'statusCode': 500,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Content-Type': 'application/json'
            },
            'body': json.dumps({
                'error': str(e),
                'traceback': traceback.format_exc()
            })
        }
