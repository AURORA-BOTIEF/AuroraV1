#!/usr/bin/env python3
"""
Local PPT Layout Testing Tool
Run this script locally to test different slide layouts without deploying to AWS.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import sys

def create_test_presentation():
    """Create a test PPT with various layout configurations"""
    
    print("🎨 Creating test PowerPoint presentation...")
    
    prs = Presentation()
    # Set 16:9 widescreen format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    colors = {
        'primary': RGBColor(0, 51, 102),      # Dark Blue
        'secondary': RGBColor(68, 114, 196),  # Medium Blue
        'accent': RGBColor(255, 192, 0)       # Gold
    }
    
    # Layout configurations to test
    configs = [
        {
            'name': 'Conservative (More Text)',
            'title_height': 0.8,
            'content_left': 0.5,
            'content_width': 7.0,
            'image_left': 8.0,
            'image_width': 4.8,
            'image_height': 4.5
        },
        {
            'name': 'Balanced 50-50',
            'title_height': 0.8,
            'content_left': 0.5,
            'content_width': 6.0,
            'image_left': 7.0,
            'image_width': 5.8,
            'image_height': 4.8
        },
        {
            'name': 'Image Focus',
            'title_height': 0.8,
            'content_left': 0.5,
            'content_width': 5.0,
            'image_left': 6.0,
            'image_width': 6.8,
            'image_height': 5.0
        }
    ]
    
    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title_shape = slide.shapes.title
    subtitle = slide.placeholders[1]
    title_shape.text = "Test de Layout PPT"
    subtitle.text = "Prueba de diferentes configuraciones de diapositivas"
    
    # Center title
    title_shape.left = Inches((13.333 - 12) / 2)
    title_shape.width = Inches(12)
    title_shape.top = Inches(2.5)
    
    # Format title
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = colors['primary']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Format subtitle
    subtitle.left = Inches((13.333 - 12) / 2)
    subtitle.width = Inches(12)
    subtitle.top = Inches(4.5)
    if subtitle.has_text_frame:
        subtitle_para = subtitle.text_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = colors['secondary']
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Create test slides for each configuration
    for config_idx, config in enumerate(configs, 1):
        print(f"  Creating test slide {config_idx}: {config['name']}")
        
        # Use blank layout for full control
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Add title
        title = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.3),
            Inches(12.5),
            Inches(config['title_height'])
        )
        title_frame = title.text_frame
        title_frame.text = f"Config {config_idx}: {config['name']}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = colors['primary']
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add title background
        title.fill.solid()
        title.fill.fore_color.rgb = RGBColor(248, 249, 250)
        
        # Add content text box
        content_top = Inches(0.3 + config['title_height'] + 0.2)
        content = slide.shapes.add_textbox(
            Inches(config['content_left']),
            content_top,
            Inches(config['content_width']),
            Inches(5.5)
        )
        
        text_frame = content.text_frame
        text_frame.word_wrap = True
        
        # Add sample bullets
        bullets = [
            "Este es el primer punto de la diapositiva",
            "Segunda viñeta con texto de ejemplo",
            "Tercera viñeta para probar el espaciado",
            "Cuarta viñeta con contenido más largo para ver cómo se comporta el texto",
            "Quinta viñeta final"
        ]
        
        for i, bullet_text in enumerate(bullets):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = f"▸ {bullet_text}"
            p.font.size = Pt(14)
            p.font.name = 'Calibri'
            p.space_before = Pt(6)
            p.space_after = Pt(4)
            p.font.color.rgb = RGBColor(33, 37, 41)
        
        # Draw border around content area (for visualization)
        content.line.color.rgb = RGBColor(200, 200, 200)
        content.line.width = Pt(1)
        
        # Add image placeholder
        image_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(config['image_left']),
            content_top,
            Inches(config['image_width']),
            Inches(config['image_height'])
        )
        
        # Style image placeholder
        fill = image_shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        line = image_shape.line
        line.color.rgb = colors['secondary']
        line.width = Pt(2)
        
        # Add text to placeholder
        img_text_frame = image_shape.text_frame
        img_text_frame.text = "IMAGEN\nAQUÍ"
        img_para = img_text_frame.paragraphs[0]
        img_para.alignment = PP_ALIGN.CENTER
        img_para.font.size = Pt(24)
        img_para.font.bold = True
        img_para.font.color.rgb = colors['secondary']
        
        # Add dimension annotations
        dims_text = (
            f"Dimensiones:\n"
            f"• Texto: {config['content_width']:.1f}\" × 5.5\" @ {config['content_left']:.1f}\"\n"
            f"• Imagen: {config['image_width']:.1f}\" × {config['image_height']:.1f}\" @ {config['image_left']:.1f}\"\n"
            f"• Gap: {config['image_left'] - config['content_left'] - config['content_width']:.1f}\""
        )
        
        dims = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(6.5),
            Inches(12.5),
            Inches(0.8)
        )
        dims_frame = dims.text_frame
        dims_frame.text = dims_text
        dims_para = dims_frame.paragraphs[0]
        dims_para.font.size = Pt(10)
        dims_para.font.color.rgb = RGBColor(100, 100, 100)
    
    # Save the presentation
    output_file = 'test_layouts.pptx'
    prs.save(output_file)
    print(f"\n✅ Test presentation created: {output_file}")
    print(f"\n📋 Instructions:")
    print(f"1. Open {output_file} in PowerPoint")
    print(f"2. Review each configuration slide")
    print(f"3. Choose the one that looks best")
    print(f"4. Tell me which config number you prefer (1, 2, or 3)")
    print(f"5. I'll apply those exact dimensions to the production code")
    print(f"\n💡 Tip: You can also manually adjust dimensions and tell me custom values!")
    
    return output_file

if __name__ == "__main__":
    try:
        create_test_presentation()
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
