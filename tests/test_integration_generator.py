import importlib.util
import io
import os
import sys


def load_module():
    path = os.path.join(os.path.dirname(__file__), "..", "CG-Backend", "lambda", "strands_ppt_generator", "strands_ppt_generator.py")
    path = os.path.abspath(path)
    spec = importlib.util.spec_from_file_location("strands_ppt_generator", path)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def make_png_bytes():
    try:
        from PIL import Image
    except Exception:
        return None
    im = Image.new('RGB', (100, 60), color=(200, 200, 210))
    out = io.BytesIO()
    im.save(out, format='PNG')
    return out.getvalue()


def test_dedicated_image_slide_fallback(monkeypatch):
    mod = load_module()

    png = make_png_bytes()
    assert png is not None, "Pillow required for this integration test"

    # Monkeypatch requests.get to return the image bytes for HTTP download
    class FakeResp:
        def __init__(self, content):
            self.status_code = 200
            self.content = content

    def fake_get(url, timeout=10):
        return FakeResp(png)

    monkeypatch.setattr('requests.get', fake_get)

    # Force reflow helper to indicate it cannot avoid overlap (None)
    monkeypatch.setattr(mod, 'adjust_image_rect_to_avoid_overlap_module', lambda *args, **kwargs: None)

    # Build a minimal presentation structure with one content slide that references an HTTP image
    presentation_structure = {
        'presentation_title': 'IntegrationTest',
        'style': 'professional',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'content',
                'title': 'Test Slide',
                'bullets': ['One', 'Two', 'Three'],
                'image_url': 'http://example.com/image.png',
                'image_reference': 'USE_IMAGE: example'
            }
        ]
    }

    # Run generator
    pptx_bytes = mod.generate_pptx_file(presentation_structure, book_data={})
    assert pptx_bytes[:2] == b'PK'

    # Inspect resulting PPTX to ensure two slides were produced (content + dedicated image slide)
    from pptx import Presentation
    prs = Presentation(io.BytesIO(pptx_bytes))
    # Expect at least 2 slides (content + dedicated image slide)
    assert len(prs.slides) >= 2

    # Ensure at least one slide contains a picture shape
    pic_count = 0
    for s in prs.slides:
        for sp in s.shapes:
            # python-pptx shapes with image have attribute .image when present
            try:
                if hasattr(sp, 'image') and sp.image is not None:
                    pic_count += 1
            except Exception:
                continue

    assert pic_count >= 1
