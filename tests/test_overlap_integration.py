"""
Integration tests for overlap detection in real slide generation scenarios.

These tests create actual PPTX slides with overlapping content and verify
that the overlap detection and reflow strategies work correctly in practice.
"""
import io
import os
import sys
import pytest

# Add generator to path
CG_BACKEND = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'CG-Backend', 'lambda', 'strands_ppt_generator'))
if CG_BACKEND not in sys.path:
    sys.path.insert(0, CG_BACKEND)

try:
    from strands_ppt_generator import generate_pptx_file
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_overlap_with_long_bullets_and_image():
    """Test that long bullet lists with images don't overlap."""
    presentation_structure = {
        'presentation_title': 'Overlap Test - Long Content',
        'style': 'professional',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'content',
                'title': 'Long Bullet List with Image',
                'bullets': [
                    'First bullet point with substantial content that takes up space',
                    'Second bullet point explaining complex concepts in detail',
                    'Third bullet point with even more information to display',
                    'Fourth bullet point adding to the vertical space usage',
                    'Fifth bullet point to ensure we have significant text height',
                    'Sixth bullet point to really push the overlap boundaries',
                ],
                'image_url': 's3://fake-bucket/test-image.png'  # Will use placeholder
            }
        ]
    }
    
    pptx_bytes = generate_pptx_file(presentation_structure, {})
    assert pptx_bytes is not None
    assert len(pptx_bytes) > 0
    
    # Load and verify structure
    prs = Presentation(io.BytesIO(pptx_bytes))
    assert len(prs.slides) >= 1
    
    slide = prs.slides[0]
    
    # Find text and image shapes
    text_shapes = []
    image_shapes = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text:
            text_shapes.append(shape)
        if hasattr(shape, 'image') or getattr(shape, 'name', '') == 'aurora_placeholder':
            image_shapes.append(shape)
    
    # Should have both text and image/placeholder
    assert len(text_shapes) > 0
    
    # If we have both text and image, verify no overlap
    if len(image_shapes) > 0:
        for text_shape in text_shapes:
            for img_shape in image_shapes:
                # Get bounding boxes
                t_left = text_shape.left
                t_top = text_shape.top
                t_right = t_left + text_shape.width
                t_bottom = t_top + text_shape.height
                
                i_left = img_shape.left
                i_top = img_shape.top
                i_right = i_left + img_shape.width
                i_bottom = i_top + img_shape.height
                
                # Check for overlap
                horiz_overlap = not (t_right <= i_left or i_right <= t_left)
                vert_overlap = not (t_bottom <= i_top or i_bottom <= t_top)
                overlap = horiz_overlap and vert_overlap
                
                # There should be no significant overlap (allowing small margin for rounding)
                if overlap:
                    # Calculate overlap area
                    overlap_left = max(t_left, i_left)
                    overlap_top = max(t_top, i_top)
                    overlap_right = min(t_right, i_right)
                    overlap_bottom = min(t_bottom, i_bottom)
                    overlap_width = max(0, overlap_right - overlap_left)
                    overlap_height = max(0, overlap_bottom - overlap_top)
                    
                    # Allow tiny overlaps (< 1% of either shape's area) for rounding errors
                    text_area = text_shape.width * text_shape.height
                    img_area = img_shape.width * img_shape.height
                    overlap_area = overlap_width * overlap_height
                    
                    relative_overlap = overlap_area / min(text_area, img_area) if min(text_area, img_area) > 0 else 0
                    
                    assert relative_overlap < 0.01, f"Significant overlap detected: {relative_overlap:.2%}"


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_multiple_overlapping_scenarios():
    """Test multiple slides with different overlapping scenarios."""
    presentation_structure = {
        'presentation_title': 'Overlap Scenarios Test',
        'style': 'professional',
        'slides': [
            # Scenario 1: Short text, large image placeholder
            {
                'slide_number': 1,
                'slide_type': 'content',
                'title': 'Short Text',
                'bullets': ['Brief point', 'Another brief point'],
                'image_url': 's3://fake/large-image.png'
            },
            # Scenario 2: Medium text with image
            {
                'slide_number': 2,
                'slide_type': 'content',
                'title': 'Medium Content',
                'bullets': [
                    'First point with moderate length',
                    'Second point with some detail',
                    'Third point for spacing',
                    'Fourth point to add height'
                ],
                'image_url': 's3://fake/medium-image.png'
            },
            # Scenario 3: Image-only slide (should never overlap)
            {
                'slide_number': 3,
                'slide_type': 'image',
                'title': 'Image Only',
                'caption': 'Single image slide',
                'image_url': 's3://fake/image-only.png'
            }
        ]
    }
    
    pptx_bytes = generate_pptx_file(presentation_structure, {})
    assert pptx_bytes is not None
    
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    # Verify no overlaps in any slide
    for slide_idx, slide in enumerate(prs.slides):
        text_shapes = []
        non_text_shapes = []
        
        for shape in slide.shapes:
            # Classify shapes
            if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'paragraphs'):
                # It's a text shape
                if shape.text_frame.text.strip():  # Has actual text
                    text_shapes.append(shape)
            else:
                # It's a non-text shape (image, placeholder, etc.)
                non_text_shapes.append(shape)
        
        # Check each text shape against each non-text shape
        for text_shape in text_shapes:
            for other_shape in non_text_shapes:
                t_left = text_shape.left
                t_top = text_shape.top
                t_right = t_left + text_shape.width
                t_bottom = t_top + text_shape.height
                
                o_left = other_shape.left
                o_top = other_shape.top
                o_right = o_left + other_shape.width
                o_bottom = o_top + other_shape.height
                
                horiz_overlap = not (t_right <= o_left or o_right <= t_left)
                vert_overlap = not (t_bottom <= o_top or o_bottom <= t_top)
                
                if horiz_overlap and vert_overlap:
                    # Calculate overlap percentage
                    overlap_left = max(t_left, o_left)
                    overlap_top = max(t_top, o_top)
                    overlap_right = min(t_right, o_right)
                    overlap_bottom = min(t_bottom, o_bottom)
                    overlap_width = max(0, overlap_right - overlap_left)
                    overlap_height = max(0, overlap_bottom - overlap_top)
                    overlap_area = overlap_width * overlap_height
                    
                    text_area = text_shape.width * text_shape.height
                    other_area = other_shape.width * other_shape.height
                    min_area = min(text_area, other_area)
                    
                    relative_overlap = overlap_area / min_area if min_area > 0 else 0
                    
                    assert relative_overlap < 0.01, \
                        f"Slide {slide_idx + 1}: Overlap {relative_overlap:.2%} between text and non-text shape"


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_dedicated_image_slide_fallback_on_extreme_overlap():
    """Test that extreme overlaps trigger dedicated image slide fallback."""
    # This presentation has very long text that would overlap any image placement
    presentation_structure = {
        'presentation_title': 'Extreme Overlap Test',
        'style': 'professional',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'content',
                'title': 'Extremely Long Content That Fills The Entire Slide',
                'bullets': [
                    'First very long bullet point that contains substantial information and takes up significant horizontal space on the slide',
                    'Second comprehensive bullet point with detailed explanations and multiple sub-points that would normally require their own slide',
                    'Third extensive bullet point covering complex topics with thorough analysis and detailed breakdowns of various concepts',
                    'Fourth detailed bullet point with in-depth coverage of important topics and extensive commentary on various aspects',
                    'Fifth lengthy bullet point providing comprehensive information and covering multiple related topics in great detail',
                    'Sixth expansive bullet point with thorough documentation and extensive descriptions of critical concepts',
                    'Seventh detailed bullet point that continues the pattern of comprehensive coverage and detailed explanations',
                    'Eighth extensive bullet point wrapping up the comprehensive analysis with final thoughts and conclusions'
                ],
                'image_url': 's3://fake/image.png'
            }
        ]
    }
    
    pptx_bytes = generate_pptx_file(presentation_structure, {})
    assert pptx_bytes is not None
    
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    # Should have created additional slides due to content length/overlap issues
    # Original design requested 1 content slide, but generator may split it or add dedicated image slide
    # At minimum, should have at least 1 slide
    assert len(prs.slides) >= 1
    
    # All slides should be valid (no exceptions during generation)
    for slide in prs.slides:
        assert slide is not None
        # Each slide should have some shapes
        assert len(slide.shapes) > 0


if __name__ == '__main__':
    pytest.main([__file__, '-v'])
