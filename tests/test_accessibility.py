"""
Tests for accessibility features in PPT generation.

Verifies that images have proper alt text and metadata for screen readers.
"""
import io
import os
import sys
import pytest

# Add generator to path
CG_BACKEND = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'CG-Backend', 'lambda', 'strands_ppt_generator'))
if CG_BACKEND not in sys.path:
    sys.path.insert(0, CG_BACKEND)

try:
    from strands_ppt_generator import generate_pptx_file
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def get_image_alt_text(picture_shape):
    """Extract alt text from a picture shape.
    
    Args:
        picture_shape: A python-pptx Picture shape
        
    Returns:
        str: The alt text, or None if not found
    """
    try:
        # Access the underlying XML element
        pic_element = picture_shape._element
        
        # Find nvPicPr or nvSpPr (non-visual properties)
        nv_props = None
        for child in pic_element:
            if 'nvPicPr' in child.tag or 'nvSpPr' in child.tag:
                nv_props = child
                break
        
        if nv_props is not None:
            # Find cNvPr within nvPicPr or nvSpPr
            for child in nv_props:
                if 'cNvPr' in child.tag:
                    # Get the descr attribute (alt text)
                    return child.get('descr')
        
        return None
    except Exception:
        return None


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_content_slide_image_has_alt_text():
    """Verify that images on content slides have alt text."""
    presentation_structure = {
        'presentation_title': 'Alt Text Test',
        'style': 'professional',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'content',
                'title': 'Network Architecture',
                'caption': 'Diagram showing network topology',
                'bullets': ['Key component 1', 'Key component 2'],
                'image_url': 's3://fake/network-diagram.png'
            }
        ]
    }
    
    pptx_bytes = generate_pptx_file(presentation_structure, {})
    assert pptx_bytes is not None
    
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[0]
    
    # Find picture shapes
    picture_found = False
    alt_text_found = False
    
    for shape in slide.shapes:
        # Check if it's a picture (has image attribute or is placeholder)
        if hasattr(shape, 'image') or getattr(shape, 'name', '') == 'aurora_placeholder':
            picture_found = True
            
            # Extract alt text
            alt_text = get_image_alt_text(shape)
            
            if alt_text:
                alt_text_found = True
                print(f"✅ Found alt text: '{alt_text}'")
                
                # Verify alt text is meaningful (not empty)
                assert len(alt_text) > 0, "Alt text should not be empty"
                
                # Should contain relevant keywords from caption or title
                text_lower = alt_text.lower()
                assert 'diagram' in text_lower or 'network' in text_lower or 'placeholder' in text_lower, \
                    f"Alt text should be descriptive: got '{alt_text}'"
    
    assert picture_found, "Should have found at least one picture shape"
    assert alt_text_found, "At least one picture should have alt text"


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_image_slide_has_alt_text():
    """Verify that image-only slides have alt text."""
    presentation_structure = {
        'presentation_title': 'Image Slide Alt Text Test',
        'style': 'professional',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'image',
                'title': 'System Architecture',
                'caption': 'Complete system overview diagram',
                'image_url': 's3://fake/architecture.png'
            }
        ]
    }
    
    pptx_bytes = generate_pptx_file(presentation_structure, {})
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[0]
    
    picture_found = False
    alt_text_found = False
    
    for shape in slide.shapes:
        if hasattr(shape, 'image') or getattr(shape, 'name', '') == 'aurora_placeholder':
            picture_found = True
            alt_text = get_image_alt_text(shape)
            
            if alt_text:
                alt_text_found = True
                print(f"✅ Image slide alt text: '{alt_text}'")
                assert len(alt_text) > 0
                
                # Should reference the content
                text_lower = alt_text.lower()
                assert any(keyword in text_lower for keyword in ['system', 'architecture', 'overview', 'diagram', 'placeholder']), \
                    f"Alt text should describe the image: got '{alt_text}'"
    
    assert picture_found, "Image slide should have a picture"
    assert alt_text_found, "Image should have alt text"


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_placeholder_has_alt_text():
    """Verify that placeholder images have descriptive alt text."""
    presentation_structure = {
        'presentation_title': 'Placeholder Alt Text Test',
        'style': 'professional',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'content',
                'title': 'Missing Image Example',
                'caption': 'Database schema',
                'bullets': ['Point 1', 'Point 2'],
                'image_url': 's3://nonexistent-bucket/missing-image.png'
            }
        ]
    }

    pptx_bytes = generate_pptx_file(presentation_structure, {})
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[0]

    placeholder_found = False
    alt_text_found = False

    for shape in slide.shapes:
        # Look for aurora_placeholder or any placeholder in the name
        shape_name = getattr(shape, 'name', '')
        if shape_name == 'aurora_placeholder' or 'placeholder' in shape_name.lower():
            placeholder_found = True
            alt_text = get_image_alt_text(shape)

            if alt_text:
                alt_text_found = True
                print(f"✅ Placeholder alt text: '{alt_text}'")
                assert len(alt_text) > 0

                # Placeholder alt text should indicate it's a placeholder
                text_lower = alt_text.lower()
                assert 'placeholder' in text_lower, \
                    f"Placeholder alt text should mention 'placeholder': got '{alt_text}'"

    assert placeholder_found, "Should have found aurora_placeholder or placeholder shape"
    assert alt_text_found, "Placeholder should have alt text"
@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_multiple_images_all_have_alt_text():
    """Verify that multiple images in a presentation all have alt text."""
    presentation_structure = {
        'presentation_title': 'Multiple Images Test',
        'style': 'professional',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'content',
                'title': 'First Slide',
                'caption': 'Network diagram',
                'bullets': ['Point A'],
                'image_url': 's3://fake/img1.png'
            },
            {
                'slide_number': 2,
                'slide_type': 'image',
                'title': 'Second Slide',
                'caption': 'Architecture overview',
                'image_url': 's3://fake/img2.png'
            },
            {
                'slide_number': 3,
                'slide_type': 'content',
                'title': 'Third Slide',
                'caption': 'Data flow',
                'bullets': ['Point B', 'Point C'],
                'image_url': 's3://fake/img3.png'
            }
        ]
    }
    
    pptx_bytes = generate_pptx_file(presentation_structure, {})
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    total_pictures = 0
    pictures_with_alt_text = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'image') or getattr(shape, 'name', '') == 'aurora_placeholder':
                total_pictures += 1
                alt_text = get_image_alt_text(shape)
                
                if alt_text and len(alt_text) > 0:
                    pictures_with_alt_text += 1
                    print(f"✅ Slide {slide_idx + 1}: Alt text '{alt_text[:50]}...'")
                else:
                    print(f"⚠️ Slide {slide_idx + 1}: No alt text found")
    
    assert total_pictures > 0, "Should have found pictures in the presentation"
    
    # At least 80% of pictures should have alt text (allowing for edge cases)
    coverage = pictures_with_alt_text / total_pictures if total_pictures > 0 else 0
    assert coverage >= 0.8, \
        f"At least 80% of pictures should have alt text: {pictures_with_alt_text}/{total_pictures} = {coverage:.1%}"


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_alt_text_length_reasonable():
    """Verify that alt text is concise (not too long or too short)."""
    presentation_structure = {
        'presentation_title': 'Alt Text Length Test',
        'style': 'professional',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'content',
                'title': 'Test Slide',
                'caption': 'A diagram showing the complete end-to-end workflow from initial data collection through processing, analysis, visualization, and final reporting with detailed annotations',
                'bullets': ['Point 1'],
                'image_url': 's3://fake/test.png'
            }
        ]
    }
    
    pptx_bytes = generate_pptx_file(presentation_structure, {})
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[0]
    
    for shape in slide.shapes:
        if hasattr(shape, 'image') or getattr(shape, 'name', '') == 'aurora_placeholder':
            alt_text = get_image_alt_text(shape)
            
            if alt_text:
                length = len(alt_text)
                print(f"Alt text length: {length} characters")
                
                # Alt text should be reasonable length
                assert length >= 5, "Alt text should be at least 5 characters"
                assert length <= 300, f"Alt text should be max 300 characters (got {length})"


if __name__ == '__main__':
    pytest.main([__file__, '-v'])
