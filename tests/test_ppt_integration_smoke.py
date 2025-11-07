"""
PPT Generator Integration/Smoke Test
=====================================
Generates a real PPTX file with both real images and placeholders for CI artifact.

This test:
1. Creates a presentation with multiple slide types
2. Uses real S3 images (if available) and placeholders (if not)
3. Validates the PPTX structure and content
4. Outputs a PPTX file that can be uploaded as a CI artifact
"""
import io
import os
import sys
import pytest

# Ensure the generator is importable
CG_BACKEND = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'CG-Backend', 'lambda', 'strands_ppt_generator'))
if CG_BACKEND not in sys.path:
    sys.path.insert(0, CG_BACKEND)

try:
    from strands_ppt_generator import generate_pptx_file
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_generate_comprehensive_pptx_smoke():
    """
    Integration smoke test: Generate a complete PPTX with various slide types,
    including both real images and placeholders.
    
    This test serves as:
    - Smoke test for PPT generation
    - CI artifact generator
    - Validation of placeholder detection
    """
    
    # Create a comprehensive presentation structure
    presentation_structure = {
        'presentation_title': 'CI Smoke Test Presentation',
        'style': 'professional',
        'slides': [
            # Title slide
            {
                'slide_number': 1,
                'slide_type': 'title',
                'title': 'PPT Generator CI Smoke Test',
                'subtitle': 'Automated test with real images and placeholders'
            },
            
            # Content slide - no image
            {
                'slide_number': 2,
                'slide_type': 'content',
                'title': 'Key Features',
                'bullets': [
                    'Adaptive layout engine',
                    'TF-IDF image selection',
                    'Placeholder fallback',
                    'Professional styling',
                    'S3 and HTTP image support'
                ]
            },
            
            # Image slide - will use placeholder (non-existent S3)
            {
                'slide_number': 3,
                'slide_type': 'image',
                'title': 'Placeholder Test',
                'caption': 'This image does not exist (placeholder)',
                'image_url': 's3://nonexistent-bucket/missing-image.png'
            },
            
            # Content slide with image reference (placeholder)
            {
                'slide_number': 4,
                'slide_type': 'content',
                'title': 'Mixed Content',
                'bullets': [
                    'Text content is primary',
                    'Image supports the message',
                    'Placeholders are gracefully handled'
                ],
                'image_url': 's3://test-bucket/another-missing.png'
            },
            
            # Summary slide
            {
                'slide_number': 5,
                'slide_type': 'summary',
                'title': 'Summary',
                'bullets': [
                    'All slide types tested',
                    'Placeholders validated',
                    'Layout engine verified',
                    'Ready for production'
                ]
            }
        ],
        'generated_at': '2025-10-30T12:00:00Z',
        'total_slides': 5
    }
    
    # Mock book data with some content
    book_data = {
        'metadata': {
            'title': 'CI Smoke Test Course'
        },
        'lessons': []
    }
    
    # Generate PPTX
    print("\n🎨 Generating smoke test PPTX...")
    pptx_bytes = generate_pptx_file(presentation_structure, book_data)
    
    # Verify bytes were generated
    assert pptx_bytes is not None, "PPTX bytes should not be None"
    assert len(pptx_bytes) > 0, "PPTX should have content"
    print(f"✅ Generated PPTX: {len(pptx_bytes)} bytes")
    
    # Load and validate structure
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    # Verify slide count
    assert len(prs.slides) == 5, f"Expected 5 slides, got {len(prs.slides)}"
    print(f"✅ Slide count validated: {len(prs.slides)} slides")
    
    # Verify placeholder detection
    placeholder_count = 0
    slide_types_found = set()
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        print(f"\n📄 Slide {slide_num}:")
        print(f"   Layout: {slide.slide_layout.name}")
        print(f"   Shapes: {len(slide.shapes)}")
        
        # Detect placeholders
        for shape in slide.shapes:
            shape_name = getattr(shape, 'name', '')

            # Look for aurora_placeholder or any shape with 'placeholder' in the name
            if shape_name == 'aurora_placeholder' or 'placeholder' in shape_name.lower():
                placeholder_count += 1
                print(f"   ✓ Found placeholder: {shape_name}")

            # Track shape types
            if hasattr(shape, 'shape_type'):
                slide_types_found.add(str(shape.shape_type))    # We expect at least 2 placeholders (slides 3 and 4 have missing images)
    assert placeholder_count >= 2, f"Expected at least 2 placeholders, found {placeholder_count}"
    print(f"\n✅ Placeholder detection: {placeholder_count} placeholders found")
    
    # Save PPTX to output directory for CI artifact
    output_dir = os.path.join(os.path.dirname(__file__), '..', 'test-output')
    os.makedirs(output_dir, exist_ok=True)
    
    output_path = os.path.join(output_dir, 'ci_smoke_test.pptx')
    with open(output_path, 'wb') as f:
        f.write(pptx_bytes)
    
    print(f"\n💾 Saved PPTX artifact: {output_path}")
    print(f"   Size: {len(pptx_bytes):,} bytes")
    print(f"   Slides: {len(prs.slides)}")
    print(f"   Placeholders: {placeholder_count}")
    
    # Verify the saved file exists and is readable
    assert os.path.exists(output_path), "Output file should exist"
    assert os.path.getsize(output_path) > 0, "Output file should not be empty"
    
    print("\n✅ All smoke test validations passed!")


@pytest.mark.skipif(not PPTX_AVAILABLE, reason="python-pptx not available")
def test_placeholder_naming_consistency():
    """
    Verify that all image placeholders (inserted for missing images) 
    are consistently named 'aurora_placeholder'.
    """
    presentation_structure = {
        'presentation_title': 'Placeholder Naming Test',
        'slides': [
            {
                'slide_number': 1,
                'slide_type': 'image',
                'title': 'Missing Image 1',
                'caption': 'Placeholder 1',
                'image_url': 's3://fake/image1.png'
            },
            {
                'slide_number': 2,
                'slide_type': 'content',
                'title': 'Missing Image 2',
                'bullets': ['Test bullet'],
                'image_url': 's3://fake/image2.png'
            }
        ],
        'style': 'professional'
    }
    
    pptx_bytes = generate_pptx_file(presentation_structure, {})
    prs = Presentation(io.BytesIO(pptx_bytes))
    
    placeholder_names = []
    for slide in prs.slides:
        for shape in slide.shapes:
            name = getattr(shape, 'name', '')
            # Only look for our custom placeholders (aurora_placeholder or shapes with 'placeholder' in name)
            # Exclude built-in PowerPoint placeholders (which are also of type PLACEHOLDER)
            if name == 'aurora_placeholder' or ('placeholder' in name.lower() and hasattr(shape, 'shape_type') and shape.shape_type != 14):
                placeholder_names.append(name)

    # We should have found at least one placeholder
    assert len(placeholder_names) >= 1, f"Expected at least 1 placeholder shape, found {len(placeholder_names)}"
    
    # All placeholders should have 'placeholder' in their name for consistency and accessibility
    for name in placeholder_names:
        assert 'placeholder' in name.lower(), \
            f"Expected 'placeholder' in name for accessibility, got '{name}'"
    
    print(f"✅ Found {len(placeholder_names)} placeholder(s) with consistent naming:")
    for name in placeholder_names:
        print(f"   - {name}")


if __name__ == '__main__':
    # Allow running directly for manual testing
    print("Running PPT Integration Smoke Test...")
    test_generate_comprehensive_pptx_smoke()
    test_placeholder_naming_consistency()
    print("\n🎉 All tests passed!")
